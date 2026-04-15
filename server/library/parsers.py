"""Library 解析器: 将各格式文件解析为 ParseResult

P0 阶段: PDF (PyMuPDF) + DOCX + Markdown + 纯文本 + Excel (基础)
独立实现, 不依赖 server/handlers/parsers.py (CLI 模式无 Mixin 可用)
"""

from __future__ import annotations

import csv
import io
import json
import logging
import re
import time
from pathlib import Path
from typing import Optional

from .models import DocumentMetadata, ParseResult, Section

logger = logging.getLogger(__name__)

# ── 可选依赖 ──

try:
    import fitz  # PyMuPDF
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

try:
    from pptx import Presentation as PptxPresentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


# ── 解析器注册表 ──

_PARSERS: dict[str, callable] = {}


def _register(ext: str):
    """注册解析器装饰器"""
    def wrapper(fn):
        _PARSERS[ext] = fn
        return fn
    return wrapper


def parse_file(file_path: str) -> ParseResult:
    """解析单个文件, 返回 ParseResult"""
    start = time.perf_counter()
    p = Path(file_path)
    ext = p.suffix.lower()

    parser = _PARSERS.get(ext)
    if parser is None:
        # 尝试当纯文本
        parser = _parse_text

    result = parser(p)
    result.parse_time_ms = int((time.perf_counter() - start) * 1000)

    # 填充通用 metadata
    stat = p.stat()
    result.metadata = DocumentMetadata(
        source_path=str(p.absolute()),
        source_name=p.name,
        file_size=stat.st_size,
        modified_at=str(stat.st_mtime),
        word_count=len(result.content),
        content_hash="",  # Scanner 已算, 这里不重复
    )

    if not result.title:
        result.title = p.stem

    if not result.summary:
        result.summary = result.content[:500]

    return result


# ──────────────────────────────────────────────
# PDF — PyMuPDF (快速, P0)
# ──────────────────────────────────────────────

@_register(".pdf")
def _parse_pdf(path: Path) -> ParseResult:
    if not HAS_PYMUPDF:
        raise RuntimeError("pymupdf 未安装, 请运行: pip install pymupdf")

    doc = fitz.open(str(path))
    pages: list[str] = []
    sections: list[Section] = []

    for i, page in enumerate(doc):
        text = page.get_text("text")
        if text.strip():
            pages.append(text)

    doc.close()

    full_text = "\n\n".join(pages)

    # 尝试提取标题 (第一行非空文本)
    title = ""
    for line in full_text.split("\n"):
        line = line.strip()
        if len(line) > 2 and len(line) < 100:
            title = line
            break

    # 尝试按标题模式拆分章节
    sections = _extract_sections(full_text)

    return ParseResult(
        title=title,
        content=full_text,
        summary=full_text[:500],
        sections=sections,
        parser_name="pymupdf",
        metadata=DocumentMetadata("", "", 0, page_count=len(pages)),
    )


# ──────────────────────────────────────────────
# Word (.docx)
# ──────────────────────────────────────────────

@_register(".docx")
def _parse_docx(path: Path) -> ParseResult:
    if not HAS_DOCX:
        raise RuntimeError("python-docx 未安装, 请运行: pip install python-docx")

    doc = DocxDocument(str(path))
    parts: list[str] = []
    sections: list[Section] = []
    current_section: Optional[dict] = None

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # 检测标题样式
        style_name = (para.style.name or "").lower()
        heading_level = 0
        if style_name.startswith("heading"):
            try:
                heading_level = int(style_name.replace("heading", "").strip())
            except ValueError:
                heading_level = 1

        if heading_level > 0:
            # 保存前一个章节
            if current_section:
                content = "\n".join(current_section["lines"])
                sections.append(Section(
                    title=current_section["title"],
                    content=content,
                    level=current_section["level"],
                    char_count=len(content),
                ))
            current_section = {"title": text, "level": heading_level, "lines": []}
            parts.append(f"{'#' * heading_level} {text}")
        else:
            parts.append(text)
            if current_section is not None:
                current_section["lines"].append(text)

    # 最后一个章节
    if current_section:
        content = "\n".join(current_section["lines"])
        sections.append(Section(
            title=current_section["title"],
            content=content,
            level=current_section["level"],
            char_count=len(content),
        ))

    # 解析表格
    for table in doc.tables:
        table_rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            table_rows.append("| " + " | ".join(cells) + " |")
        if table_rows:
            # 插入 Markdown 表头分隔
            if len(table_rows) > 1:
                header = table_rows[0]
                sep = "| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |"
                table_rows.insert(1, sep)
            parts.append("\n".join(table_rows))

    full_text = "\n\n".join(parts)
    title = sections[0].title if sections else ""

    return ParseResult(
        title=title,
        content=full_text,
        summary=full_text[:500],
        sections=sections,
        parser_name="python-docx",
    )


# ──────────────────────────────────────────────
# Excel (.xlsx) — 基础解析 (P0)
# ──────────────────────────────────────────────

@_register(".xlsx")
@_register(".xls")
def _parse_xlsx(path: Path) -> ParseResult:
    if not HAS_XLSX:
        raise RuntimeError("openpyxl 未安装, 请运行: pip install openpyxl")

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    key_data: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"## Sheet: {sheet_name}")

        rows: list[list[str]] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append(cells)

        if not rows:
            continue

        # 转为 Markdown 表格
        if rows:
            header = "| " + " | ".join(rows[0]) + " |"
            sep = "| " + " | ".join(["---"] * len(rows[0])) + " |"
            body = "\n".join("| " + " | ".join(r) + " |" for r in rows[1:])
            parts.append(f"{header}\n{sep}\n{body}")

    wb.close()
    full_text = "\n\n".join(parts)

    return ParseResult(
        title=path.stem,
        content=full_text,
        summary=full_text[:500],
        key_data=key_data,
        parser_name="openpyxl",
        doc_type="data",
    )


# ──────────────────────────────────────────────
# PPT (.pptx) — P1, 但基础实现放这里
# ──────────────────────────────────────────────

@_register(".pptx")
def _parse_pptx(path: Path) -> ParseResult:
    if not HAS_PPTX:
        raise RuntimeError("python-pptx 未安装, 请运行: pip install python-pptx")

    prs = PptxPresentation(str(path))
    slides: list[str] = []

    for i, slide in enumerate(prs.slides):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            slides.append(f"## 幻灯片 {i + 1}\n" + "\n".join(texts))

    full_text = "\n\n".join(slides)

    return ParseResult(
        title=path.stem,
        content=full_text,
        summary=full_text[:500],
        parser_name="python-pptx",
        doc_type="slides",
    )


# ──────────────────────────────────────────────
# Markdown
# ──────────────────────────────────────────────

@_register(".md")
@_register(".markdown")
def _parse_markdown(path: Path) -> ParseResult:
    text = path.read_text(encoding="utf-8", errors="replace")
    sections = _extract_sections(text)

    # 提取标题
    title = ""
    for line in text.split("\n"):
        if line.startswith("# "):
            title = line[2:].strip()
            break

    return ParseResult(
        title=title,
        content=text,
        summary=text[:500],
        sections=sections,
        parser_name="markdown",
    )


# ──────────────────────────────────────────────
# 纯文本 / CSV / TSV
# ──────────────────────────────────────────────

@_register(".txt")
@_register(".text")
def _parse_text(path: Path) -> ParseResult:
    text = path.read_text(encoding="utf-8", errors="replace")
    return ParseResult(
        title=path.stem,
        content=text,
        summary=text[:500],
        parser_name="plaintext",
    )


@_register(".csv")
def _parse_csv(path: Path) -> ParseResult:
    text = path.read_text(encoding="utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows = list(reader)

    if not rows:
        return ParseResult(title=path.stem, content="", summary="", parser_name="csv")

    # 转 Markdown 表格
    header = "| " + " | ".join(rows[0]) + " |"
    sep = "| " + " | ".join(["---"] * len(rows[0])) + " |"
    body_lines = ["| " + " | ".join(r) + " |" for r in rows[1:]]
    md = f"{header}\n{sep}\n" + "\n".join(body_lines)

    return ParseResult(
        title=path.stem,
        content=md,
        summary=md[:500],
        parser_name="csv",
        doc_type="data",
    )


@_register(".tsv")
def _parse_tsv(path: Path) -> ParseResult:
    text = path.read_text(encoding="utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text), delimiter="\t")
    rows = list(reader)

    if not rows:
        return ParseResult(title=path.stem, content="", summary="", parser_name="tsv")

    header = "| " + " | ".join(rows[0]) + " |"
    sep = "| " + " | ".join(["---"] * len(rows[0])) + " |"
    body_lines = ["| " + " | ".join(r) + " |" for r in rows[1:]]
    md = f"{header}\n{sep}\n" + "\n".join(body_lines)

    return ParseResult(
        title=path.stem,
        content=md,
        summary=md[:500],
        parser_name="tsv",
        doc_type="data",
    )


# ──────────────────────────────────────────────
# 工具函数
# ──────────────────────────────────────────────

# 中文/英文标题模式
_HEADING_RE = re.compile(
    r"^(#{1,6})\s+(.+)$"          # Markdown 标题
    r"|^(第[一二三四五六七八九十\d]+[章节条])\s*(.+)$"   # 中文章节
    r"|^([一二三四五六七八九十]+、)\s*(.+)$"              # 中文序号
    r"|^(\d+\.)\s+(.+)$",          # 数字序号
    re.MULTILINE,
)


def _extract_sections(text: str) -> list[Section]:
    """从文本中按标题模式提取章节结构"""
    sections: list[Section] = []
    lines = text.split("\n")
    current_title = ""
    current_level = 1
    current_lines: list[str] = []

    for line in lines:
        heading_match = _HEADING_RE.match(line.strip())
        if heading_match:
            # 保存前一章节
            if current_lines or current_title:
                content = "\n".join(current_lines)
                sections.append(Section(
                    title=current_title,
                    content=content,
                    level=current_level,
                    char_count=len(content),
                    has_tables="|" in content and "---" in content,
                    has_data=bool(re.search(r"\d+[%％万亿元]", content)),
                ))

            # 解析标题层级
            groups = heading_match.groups()
            if groups[0]:  # Markdown # heading
                current_level = len(groups[0])
                current_title = groups[1]
            elif groups[2]:  # 第X章
                current_level = 1
                current_title = f"{groups[2]} {groups[3]}"
            elif groups[4]:  # 一、
                current_level = 2
                current_title = f"{groups[4]} {groups[5]}"
            elif groups[6]:  # 1.
                current_level = 3
                current_title = f"{groups[6]} {groups[7]}"
            current_lines = []
        else:
            current_lines.append(line)

    # 最后一段
    if current_lines or current_title:
        content = "\n".join(current_lines)
        sections.append(Section(
            title=current_title,
            content=content,
            level=current_level,
            char_count=len(content),
            has_tables="|" in content and "---" in content,
            has_data=bool(re.search(r"\d+[%％万亿元]", content)),
        ))

    return sections
