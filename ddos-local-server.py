#!/usr/bin/env python3
"""
DD-OS Native Server v3.0
独立运行的本地 AI 操作系统后端

功能:
    - 文件操作 (读/写/列目录)
    - 命令执行 (Shell)
    - 任务管理 (后台执行)
    - 记忆持久化

用法:
    python ddos-local-server.py [--port 3001] [--path ~/clawd]

API:
    GET  /status              - 服务状态
    GET  /files               - 列出所有文件
    GET  /file/<name>         - 获取文件内容
    GET  /skills              - 获取技能列表
    GET  /memories            - 获取记忆数据
    GET  /all                 - 获取所有数据
    POST /api/tools/execute   - 执行工具 (新)
    POST /task/execute        - 执行任务 (兼容旧接口)
    GET  /task/status/<id>    - 查询任务状态
"""
from __future__ import annotations

import os
import sys
import re
import json
import argparse
import threading
import time
import uuid
import subprocess
import shlex
import shutil
from pathlib import Path
from http.server import ThreadingHTTPServer, BaseHTTPRequestHandler
from urllib.parse import unquote, urlparse, parse_qs
from datetime import datetime, timedelta
from concurrent.futures import ThreadPoolExecutor, Future

# 🌐 全局绕过系统代理: Windows 注册表可能配有本地代理 (如 Clash/V2Ray 的 127.0.0.1:13658)
# Python 的 urllib/requests 默认读取系统代理，如果代理未运行则所有外发 HTTP 请求超时
# 在服务启动时全局安装无代理的 opener，确保 urllib.request.urlopen() 直连
import urllib.request as _urllib_req
_no_proxy_handler = _urllib_req.ProxyHandler({})
_opener = _urllib_req.build_opener(_no_proxy_handler)
_urllib_req.install_opener(_opener)

# PyYAML (skill-executor/parser.py 已依赖)
try:
    import yaml
    HAS_YAML = True
except ImportError:
    HAS_YAML = False

# MCP 客户端支持
try:
    from skills.mcp_manager import MCPClientManager
    HAS_MCP = True
except ImportError:
    HAS_MCP = False
    MCPClientManager = None

# 文件解析 (可选依赖，缺失时降级)
try:
    import pdfplumber
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation as PptxPresentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import pytesseract
    from PIL import Image
    HAS_OCR = True
except ImportError:
    HAS_OCR = False

import base64
import io
import sqlite3

VERSION = "4.1.0"

# 🏠 应用根目录 (兼容 PyInstaller frozen 模式)
# PyInstaller --onefile 模式下 __file__ 指向临时解压目录，需用 sys.executable 定位实际路径
if getattr(sys, 'frozen', False):
    APP_DIR = Path(sys.executable).parent.resolve()
else:
    APP_DIR = Path(__file__).parent.resolve()


# ============================================
# V2: SQLite 数据库初始化
# ============================================

def init_sqlite_db(db_path: Path) -> sqlite3.Connection:
    """初始化 SQLite 数据库，创建 V2 所需的表"""
    conn = sqlite3.connect(str(db_path), check_same_thread=False)
    conn.row_factory = sqlite3.Row
    conn.execute("PRAGMA journal_mode=WAL")
    conn.execute("PRAGMA foreign_keys=ON")

    conn.executescript("""
        -- 会话表
        CREATE TABLE IF NOT EXISTS sessions (
            id TEXT PRIMARY KEY,
            title TEXT NOT NULL DEFAULT '',
            type TEXT NOT NULL DEFAULT 'general',
            nexus_id TEXT,
            created_at INTEGER NOT NULL,
            updated_at INTEGER NOT NULL,
            last_message_preview TEXT DEFAULT ''
        );

        -- 消息表
        CREATE TABLE IF NOT EXISTS messages (
            id TEXT PRIMARY KEY,
            session_id TEXT NOT NULL,
            role TEXT NOT NULL,
            content TEXT NOT NULL,
            timestamp INTEGER NOT NULL,
            metadata TEXT,
            FOREIGN KEY (session_id) REFERENCES sessions(id) ON DELETE CASCADE
        );
        CREATE INDEX IF NOT EXISTS idx_messages_session ON messages(session_id, timestamp);

        -- 检查点表 (断点续作)
        CREATE TABLE IF NOT EXISTS checkpoints (
            session_id TEXT PRIMARY KEY,
            data TEXT NOT NULL,
            created_at INTEGER NOT NULL,
            FOREIGN KEY (session_id) REFERENCES sessions(id) ON DELETE CASCADE
        );

        -- 记忆表 (FTS5 全文搜索)
        CREATE TABLE IF NOT EXISTS memory (
            id TEXT PRIMARY KEY,
            source TEXT NOT NULL DEFAULT 'ephemeral',
            content TEXT NOT NULL,
            nexus_id TEXT,
            tags TEXT DEFAULT '[]',
            metadata TEXT DEFAULT '{}',
            created_at INTEGER NOT NULL
        );
        CREATE INDEX IF NOT EXISTS idx_memory_source ON memory(source);
        CREATE INDEX IF NOT EXISTS idx_memory_nexus ON memory(nexus_id);

        -- FTS5 虚拟表 (全文搜索)
        CREATE VIRTUAL TABLE IF NOT EXISTS memory_fts USING fts5(
            content,
            tags,
            content='memory',
            content_rowid='rowid'
        );

        -- 自动同步 FTS 索引的触发器
        CREATE TRIGGER IF NOT EXISTS memory_ai AFTER INSERT ON memory BEGIN
            INSERT INTO memory_fts(rowid, content, tags)
            VALUES (new.rowid, new.content, new.tags);
        END;
        CREATE TRIGGER IF NOT EXISTS memory_ad AFTER DELETE ON memory BEGIN
            INSERT INTO memory_fts(memory_fts, rowid, content, tags)
            VALUES ('delete', old.rowid, old.content, old.tags);
        END;
        CREATE TRIGGER IF NOT EXISTS memory_au AFTER UPDATE ON memory BEGIN
            INSERT INTO memory_fts(memory_fts, rowid, content, tags)
            VALUES ('delete', old.rowid, old.content, old.tags);
            INSERT INTO memory_fts(rowid, content, tags)
            VALUES (new.rowid, new.content, new.tags);
        END;

        -- 评分表
        CREATE TABLE IF NOT EXISTS nexus_scoring (
            nexus_id TEXT PRIMARY KEY,
            scoring_data TEXT NOT NULL,
            updated_at INTEGER NOT NULL
        );
    """)

    # V3: 安全地添加 confidence 列 (如果不存在)
    try:
        conn.execute("SELECT confidence FROM memory LIMIT 1")
    except sqlite3.OperationalError:
        conn.execute("ALTER TABLE memory ADD COLUMN confidence REAL DEFAULT 0.5")
        print("[SQLite] Added 'confidence' column to memory table")

    conn.commit()
    print(f"[SQLite] Database initialized at {db_path}")
    return conn


# 全局数据库连接 (线程安全 WAL 模式)
_db_conn: sqlite3.Connection | None = None
_db_lock = threading.Lock()

# 🛡️ 安全配置
DANGEROUS_COMMANDS = {'rm -rf /', 'format', 'mkfs', 'dd if=/dev/zero'}
MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB 最大文件大小
MAX_OUTPUT_SIZE = 512 * 1024      # 512KB 最大输出
PLUGIN_TIMEOUT = 60               # 插件执行超时(秒)


def safe_utf8_truncate(text: str, max_bytes: int) -> str:
    """UTF-8 安全截断，不破坏多字节字符"""
    encoded = text.encode('utf-8')
    if len(encoded) <= max_bytes:
        return text
    safe_idx = max_bytes
    while safe_idx > 0 and (encoded[safe_idx] & 0xC0) == 0x80:
        safe_idx -= 1
    return encoded[:safe_idx].decode('utf-8') + f"\n[已截断至约 {max_bytes // 1024}KB]"

# 🌐 静态文件 MIME 类型映射
MIME_TYPES = {
    '.html': 'text/html',
    '.css': 'text/css',
    '.js': 'application/javascript',
    '.json': 'application/json',
    '.png': 'image/png',
    '.jpg': 'image/jpeg',
    '.jpeg': 'image/jpeg',
    '.gif': 'image/gif',
    '.svg': 'image/svg+xml',
    '.ico': 'image/x-icon',
    '.woff': 'font/woff',
    '.woff2': 'font/woff2',
    '.ttf': 'font/ttf',
    '.eot': 'application/vnd.ms-fontobject',
}


# ============================================
# 🧮 文本相似度计算 (用于 Nexus 去重)
# ============================================

def calculate_text_similarity(text1: str, text2: str) -> float:
    """计算两个文本的 N-gram Jaccard 相似度"""
    if not text1 or not text2:
        return 0.0
    
    def get_ngrams(text: str) -> set:
        text = text.lower()
        # 清理符号，保留中英文和数字
        text = re.sub(r'[^\w\s\u4e00-\u9fff]', '', text)
        chars = list(text.replace(' ', ''))
        if len(chars) < 2:
            return set(chars)
        # 提取单字和相邻双字词 (Bi-gram)
        bigrams = [''.join(chars[i:i+2]) for i in range(len(chars)-1)]
        return set(chars + bigrams)
    
    set1 = get_ngrams(text1)
    set2 = get_ngrams(text2)
    if not set1 or not set2:
        return 0.0
    
    intersection = set1.intersection(set2)
    union = set1.union(set2)
    return len(intersection) / len(union)


# ============================================
# 🔌 SKILL.md Frontmatter 解析
# ============================================

def parse_skill_frontmatter(skill_md_path: Path) -> dict:
    """从 SKILL.md 提取 YAML frontmatter 元数据 (支持嵌套 metadata.openclaw)"""
    try:
        content = skill_md_path.read_text(encoding='utf-8')
    except Exception:
        return {}

    match = re.match(r'^---\s*\n(.*?)\n---\s*\n', content, re.DOTALL)
    if not match:
        # 无 frontmatter，提取第一段非标题文本作为 description
        desc = ''
        for line in content.split('\n'):
            line = line.strip()
            if line and not line.startswith('#'):
                desc = line[:200]
                break
        return {'description': desc}

    if HAS_YAML:
        try:
            raw = yaml.safe_load(match.group(1)) or {}
        except Exception:
            return {}
    else:
        # 无 PyYAML 时用简单正则提取 key: value
        raw = {}
        for line in match.group(1).split('\n'):
            m = re.match(r'^(\w+)\s*:\s*(.+)$', line.strip())
            if m:
                key, val = m.group(1), m.group(2).strip()
                # 简单处理数组 [a, b, c]
                if val.startswith('[') and val.endswith(']'):
                    val = [v.strip().strip('"\'') for v in val[1:-1].split(',') if v.strip()]
                # 尝试 JSON 解析 (处理 metadata 等嵌套字段)
                elif val.startswith('{'):
                    try:
                        val = json.loads(val)
                    except Exception:
                        pass
                raw[key] = val
        return _flatten_openclaw_metadata(raw)

    return _flatten_openclaw_metadata(raw)


def _flatten_openclaw_metadata(raw: dict) -> dict:
    """从 frontmatter 中提取 metadata.openclaw 字段并扁平化到顶层"""
    metadata = raw.get('metadata', {})
    if isinstance(metadata, dict):
        openclaw = metadata.get('openclaw', {})
        if isinstance(openclaw, dict):
            # 提升 openclaw 字段到顶层 (不覆盖已有字段)
            for key in ('emoji', 'primaryEnv', 'requires', 'install', 'anyBins'):
                if key in openclaw and key not in raw:
                    raw[key] = openclaw[key]
            # requires 中的 anyBins 提升
            oc_requires = openclaw.get('requires', {})
            if isinstance(oc_requires, dict):
                if 'requires' not in raw:
                    raw['requires'] = oc_requires
                else:
                    # 合并
                    existing = raw['requires'] if isinstance(raw['requires'], dict) else {}
                    for k, v in oc_requires.items():
                        if k not in existing:
                            existing[k] = v
                    raw['requires'] = existing
            oc_install = openclaw.get('install', [])
            if oc_install and 'install' not in raw:
                raw['install'] = oc_install
        # 保留原始 metadata.openclaw 用于序列化
        raw['_openclaw'] = openclaw
    return raw


def skill_name_to_tool_name(name: str) -> str:
    """将 skill 名称标准化为工具名 (kebab-case -> snake_case)"""
    return name.replace('-', '_').replace(' ', '_').lower()


# ============================================
# 🌌 NEXUS.md 解析
# ============================================

def parse_nexus_frontmatter(nexus_md_path: Path) -> dict:
    """从 NEXUS.md 提取 YAML frontmatter 元数据"""
    try:
        content = nexus_md_path.read_text(encoding='utf-8')
    except Exception:
        return {}

    match = re.match(r'^---\s*\n(.*?)\n---\s*\n', content, re.DOTALL)
    if not match:
        return {}

    if HAS_YAML:
        try:
            return yaml.safe_load(match.group(1)) or {}
        except Exception:
            return {}
    else:
        result = {}
        for line in match.group(1).split('\n'):
            m = re.match(r'^(\w+)\s*:\s*(.+)$', line.strip())
            if m:
                key, val = m.group(1), m.group(2).strip()
                if val.startswith('[') and val.endswith(']'):
                    val = [v.strip().strip('"\'') for v in val[1:-1].split(',') if v.strip()]
                result[key] = val
        return result


def extract_nexus_body(nexus_md_path: Path) -> str:
    """从 NEXUS.md 提取 Markdown 正文 (跳过 frontmatter)"""
    try:
        content = nexus_md_path.read_text(encoding='utf-8')
    except Exception:
        return ''

    # 去掉 frontmatter
    match = re.match(r'^---\s*\n.*?\n---\s*\n', content, re.DOTALL)
    if match:
        return content[match.end():].strip()
    return content.strip()


def update_nexus_frontmatter(nexus_md_path: Path, updates: dict):
    """更新 NEXUS.md 的 frontmatter 字段 (保留 body 不变)"""
    body = extract_nexus_body(nexus_md_path)
    frontmatter = parse_nexus_frontmatter(nexus_md_path)
    frontmatter.update(updates)

    # 重建 YAML frontmatter
    lines = ['---']
    for key, val in frontmatter.items():
        if isinstance(val, list):
            lines.append(f'{key}:')
            for item in val:
                lines.append(f'  - {item}')
        elif isinstance(val, dict):
            lines.append(f'{key}:')
            for k, v in val.items():
                lines.append(f'  {k}: {v}')
        else:
            lines.append(f'{key}: {val}')
    lines.append('---')
    lines.append('')
    lines.append(body)

    nexus_md_path.write_text('\n'.join(lines), encoding='utf-8')


def count_experience_entries(exp_dir: Path) -> int:
    """统计经验目录中的条目数，用于 XP 计算"""
    xp = 0
    successes = exp_dir / 'successes.md'
    failures = exp_dir / 'failures.md'
    if successes.exists():
        try:
            lines = successes.read_text(encoding='utf-8').split('\n')
            xp += sum(1 for l in lines if l.startswith('### ')) * 10
        except Exception:
            pass
    if failures.exists():
        try:
            lines = failures.read_text(encoding='utf-8').split('\n')
            xp += sum(1 for l in lines if l.startswith('### ')) * 5
        except Exception:
            pass
    return xp


# ============================================
# 🔌 动态工具注册表
# ============================================

class ToolRegistry:
    """动态工具发现与注册 - 支持内置工具 + 插件工具 + 指令型工具 + MCP工具"""

    def __init__(self, clawd_path: Path, project_path: Path = None):
        self.clawd_path = clawd_path
        # 项目目录 (脚本/exe 所在目录)，用于加载内置技能
        self.project_path = project_path or APP_DIR
        self.builtin_tools: dict = {}      # name -> callable
        self.plugin_tools: dict = {}       # name -> ToolSpec dict (有 execute.py)
        self.instruction_tools: dict = {}  # name -> InstructionSpec (纯 SKILL.md)
        self.mcp_tools: dict = {}          # name -> MCPToolSpec dict (MCP 服务器)
        self.mcp_manager: 'MCPClientManager | None' = None

    def register_builtin(self, name: str, handler):
        """注册内置工具"""
        self.builtin_tools[name] = handler

    def _get_skills_dirs(self) -> list[Path]:
        """获取所有技能目录 (用户目录 + 项目目录)"""
        dirs = []
        # 用户数据目录的技能 (优先级高，可覆盖内置)
        user_skills = self.clawd_path / 'skills'
        if user_skills.exists():
            dirs.append(user_skills)
        # 项目目录的内置技能
        project_skills = self.project_path / 'skills'
        if project_skills.exists() and project_skills != user_skills:
            dirs.append(project_skills)
        return dirs

    def scan_plugins(self):
        """递归扫描 skills/ 目录，统一从 SKILL.md frontmatter 注册可执行插件 + 指令型技能"""
        skills_dirs = self._get_skills_dirs()
        if not skills_dirs:
            return

        plugin_count = 0
        instruction_count = 0

        seen_dirs: set = set()
        seen_tools: set = set()  # 防止重复注册同名工具

        for skills_dir in skills_dirs:
            # ── 统一扫描 SKILL.md ──
            for skill_md in skills_dir.rglob('SKILL.md'):
                skill_dir = skill_md.parent
                dir_key = str(skill_dir.resolve())

                if dir_key in seen_dirs:
                    continue
                seen_dirs.add(dir_key)

                try:
                    frontmatter = parse_skill_frontmatter(skill_md)
                    original_name = frontmatter.get('name', skill_dir.name)
                    executable = frontmatter.get('executable', '')
                    runtime = frontmatter.get('runtime', 'python')

                    if executable:
                        # ── 可执行技能 (有 executable 字段) ──
                        exe_path = skill_dir / executable
                        if not exe_path.exists():
                            print(f"[ToolRegistry] Warning: {exe_path} not found for skill '{original_name}', skipping")
                            continue

                        tools_list = frontmatter.get('tools', [])
                        if tools_list:
                            # 多工具技能: frontmatter 中有 tools 数组
                            for tool_spec in tools_list:
                                tool_name = tool_spec.get('toolName', '')
                                if not tool_name:
                                    continue
                                if tool_name in seen_tools:
                                    continue
                                if tool_name in self.builtin_tools:
                                    print(f"[ToolRegistry] Warning: plugin '{tool_name}' conflicts with builtin, skipping")
                                    continue

                                self.plugin_tools[tool_name] = {
                                    'name': tool_name,
                                    'exe_path': str(exe_path),
                                    'runtime': tool_spec.get('runtime', runtime),
                                    'inputs': tool_spec.get('inputs', {}),
                                    'outputs': tool_spec.get('outputs', {}),
                                    'description': tool_spec.get('description', ''),
                                    'dangerLevel': tool_spec.get('dangerLevel', frontmatter.get('dangerLevel', 'safe')),
                                    'version': frontmatter.get('version', '1.0.0'),
                                    'skill_dir': str(skill_dir),
                                    'keywords': tool_spec.get('keywords', frontmatter.get('keywords', [])),
                                }
                                seen_tools.add(tool_name)
                                plugin_count += 1
                                print(f"[ToolRegistry] Registered plugin: {tool_name} ({exe_path.name})")
                        else:
                            # 单工具技能: toolName = skill_name_to_tool_name(name)
                            tool_name = skill_name_to_tool_name(original_name)
                            if tool_name in seen_tools:
                                continue
                            if tool_name in self.builtin_tools:
                                print(f"[ToolRegistry] Warning: plugin '{tool_name}' conflicts with builtin, skipping")
                                continue

                            self.plugin_tools[tool_name] = {
                                'name': tool_name,
                                'exe_path': str(exe_path),
                                'runtime': runtime,
                                'inputs': frontmatter.get('inputs', {}),
                                'outputs': {},
                                'description': frontmatter.get('description', ''),
                                'dangerLevel': frontmatter.get('dangerLevel', 'safe'),
                                'version': frontmatter.get('version', '1.0.0'),
                                'skill_dir': str(skill_dir),
                                'keywords': frontmatter.get('tags', frontmatter.get('keywords', [])),
                            }
                            seen_tools.add(tool_name)
                            plugin_count += 1
                            print(f"[ToolRegistry] Registered plugin: {tool_name} ({exe_path.name})")
                    else:
                        # ── 指令型技能 (无 executable) ──
                        tool_name = skill_name_to_tool_name(original_name)

                        if tool_name in self.builtin_tools or tool_name in self.plugin_tools or tool_name in seen_tools:
                            print(f"[ToolRegistry] Warning: instruction skill '{tool_name}' conflicts, skipping")
                            continue

                        self.instruction_tools[tool_name] = {
                            'name': tool_name,
                            'original_name': original_name,
                            'skill_path': str(skill_md),
                            'skill_dir': str(skill_dir),
                            'description': frontmatter.get('description', ''),
                            'inputs': frontmatter.get('inputs', {}),
                            'keywords': frontmatter.get('tags', frontmatter.get('keywords', [])),
                            'dangerLevel': 'safe',
                            'version': frontmatter.get('version', '1.0.0'),
                        }
                        seen_tools.add(tool_name)
                        instruction_count += 1
                        print(f"[ToolRegistry] Registered instruction skill: {tool_name} (from {skills_dir.name})")

                except Exception as e:
                    print(f"[ToolRegistry] Error loading {skill_md}: {e}")

            # ── Deprecated fallback: manifest.json (兼容无 SKILL.md 的第三方技能) ──
            for manifest_path in skills_dir.rglob('manifest.json'):
                skill_dir = manifest_path.parent
                dir_key = str(skill_dir.resolve())
                if dir_key in seen_dirs:
                    continue
                seen_dirs.add(dir_key)

                print(f"[ToolRegistry] ⚠️ DEPRECATED: {manifest_path} has no SKILL.md, please migrate to SKILL.md format")

                try:
                    spec = json.loads(manifest_path.read_text(encoding='utf-8'))
                    tools_list = spec.get('tools', [])
                    if not tools_list:
                        tools_list = [spec]

                    for tool_spec in tools_list:
                        tool_name = tool_spec.get('toolName', '')
                        executable = tool_spec.get('executable', spec.get('executable', 'execute.py'))

                        if not tool_name or tool_name in seen_tools:
                            continue

                        exe_path = skill_dir / executable
                        if not exe_path.exists():
                            continue

                        if tool_name in self.builtin_tools:
                            continue

                        self.plugin_tools[tool_name] = {
                            'name': tool_name,
                            'exe_path': str(exe_path),
                            'runtime': tool_spec.get('runtime', spec.get('runtime', 'python')),
                            'inputs': tool_spec.get('inputs', {}),
                            'outputs': tool_spec.get('outputs', {}),
                            'description': tool_spec.get('description', ''),
                            'dangerLevel': tool_spec.get('dangerLevel', spec.get('dangerLevel', 'safe')),
                            'version': tool_spec.get('version', spec.get('version', '1.0.0')),
                            'skill_dir': str(skill_dir),
                            'keywords': tool_spec.get('keywords', spec.get('keywords', [])),
                        }
                        seen_tools.add(tool_name)
                        plugin_count += 1

                except Exception as e:
                    print(f"[ToolRegistry] Error loading deprecated manifest {manifest_path}: {e}")

        total = plugin_count + instruction_count
        if total > 0:
            print(f"[ToolRegistry] {total} tool(s) registered ({plugin_count} plugins, {instruction_count} instruction skills)")

    def scan_mcp_servers(self):
        """扫描并连接 MCP 服务器"""
        if not HAS_MCP:
            print("[ToolRegistry] MCP support not available (missing mcp_manager)")
            return

        self.mcp_manager = MCPClientManager(clawd_path=self.clawd_path)
        count = self.mcp_manager.initialize_all()

        if count > 0:
            # 注册 MCP 工具
            mcp_tool_count = 0
            for tool_info in self.mcp_manager.get_all_tools():
                tool_name = tool_info['name']
                # 冲突检查
                if tool_name in self.builtin_tools or tool_name in self.plugin_tools or tool_name in self.instruction_tools:
                    print(f"[ToolRegistry] Warning: MCP tool '{tool_name}' conflicts with existing tool, skipping")
                    continue

                self.mcp_tools[tool_name] = {
                    'name': tool_name,
                    'server': tool_info.get('server', ''),
                    'description': tool_info.get('description', ''),
                    'inputs': tool_info.get('inputs', {}),
                    'dangerLevel': 'safe',
                    'version': '1.0.0',
                }
                mcp_tool_count += 1

            print(f"[ToolRegistry] {mcp_tool_count} MCP tool(s) registered from {count} server(s)")

    def is_registered(self, name: str) -> bool:
        return name in self.builtin_tools or name in self.plugin_tools or name in self.instruction_tools or name in self.mcp_tools

    def get_plugin(self, name: str) -> dict | None:
        return self.plugin_tools.get(name)

    def get_instruction(self, name: str) -> dict | None:
        return self.instruction_tools.get(name)

    def get_mcp_tool(self, name: str) -> dict | None:
        return self.mcp_tools.get(name)

    def list_all(self) -> list:
        """返回所有已注册工具（内置+插件+指令型+MCP）"""
        # 内置工具元数据 (为有特殊参数的工具提供描述)
        BUILTIN_META = {
            'nexusBindSkill': {
                'description': '为当前 Nexus 绑定新技能依赖',
                'inputs': {
                    'nexusId': {'type': 'string', 'description': 'Nexus ID', 'required': True},
                    'skillId': {'type': 'string', 'description': '要绑定的技能 ID', 'required': True},
                },
            },
            'nexusUnbindSkill': {
                'description': '从当前 Nexus 移除技能依赖',
                'inputs': {
                    'nexusId': {'type': 'string', 'description': 'Nexus ID', 'required': True},
                    'skillId': {'type': 'string', 'description': '要移除的技能 ID', 'required': True},
                },
            },
            'parseFile': {
                'description': '解析文档文件（PDF/DOCX/PPTX）或对图像进行OCR文字识别，返回提取的文本内容',
                'inputs': {
                    'filePath': {'type': 'string', 'description': '文件路径（支持 .pdf .docx .pptx .png .jpg 等格式）', 'required': True},
                },
            },
            'generateSkill': {
                'description': '动态生成 Python SKILL 并保存。当现有工具无法完成任务时，用此工具创建新能力',
                'inputs': {
                    'name': {'type': 'string', 'description': '技能名称 (kebab-case，如 ppt-maker)', 'required': True},
                    'description': {'type': 'string', 'description': '技能功能描述', 'required': True},
                    'pythonCode': {'type': 'string', 'description': 'Python 实现代码（必须包含 main() 函数）', 'required': True},
                    'nexusId': {'type': 'string', 'description': '关联的 Nexus ID（可选，指定后保存到 Nexus 目录）', 'required': False},
                    'triggers': {'type': 'array', 'description': '触发关键词列表（可选）', 'required': False},
                },
            },
            'browser_navigate': {
                'description': '使用浏览器导航到指定 URL，返回页面标题和文本内容。支持 JavaScript 渲染的动态页面',
                'inputs': {
                    'url': {'type': 'string', 'description': '要访问的网页 URL', 'required': True},
                    'waitUntil': {'type': 'string', 'description': '等待条件: domcontentloaded(默认) / networkidle / load', 'required': False},
                },
            },
            'browser_click': {
                'description': '点击浏览器当前页面上的元素（需先用 browser_navigate 打开页面）',
                'inputs': {
                    'selector': {'type': 'string', 'description': 'CSS 选择器或文本选择器，如 "button.submit" 或 "text=登录"', 'required': True},
                },
            },
            'browser_fill': {
                'description': '在浏览器当前页面的输入框中填写内容',
                'inputs': {
                    'selector': {'type': 'string', 'description': 'CSS 选择器，如 "input[name=search]" 或 "#username"', 'required': True},
                    'value': {'type': 'string', 'description': '要填写的文本内容', 'required': True},
                },
            },
            'browser_extract': {
                'description': '提取浏览器当前页面的文本内容（支持指定选择器提取局部内容）',
                'inputs': {
                    'selector': {'type': 'string', 'description': 'CSS 选择器（默认 body，提取主要内容区域）', 'required': False},
                },
            },
            'browser_screenshot': {
                'description': '对浏览器当前页面截图',
                'inputs': {
                    'selector': {'type': 'string', 'description': '指定截图区域的 CSS 选择器（可选，默认整个页面）', 'required': False},
                    'fullPage': {'type': 'boolean', 'description': '是否截取完整页面（包括滚动区域）', 'required': False},
                },
            },
            'browser_evaluate': {
                'description': '在浏览器当前页面执行 JavaScript 代码并返回结果',
                'inputs': {
                    'expression': {'type': 'string', 'description': 'JavaScript 表达式（支持 async）', 'required': True},
                },
            },
        }
        tools = []
        for name in self.builtin_tools:
            meta = BUILTIN_META.get(name, {})
            tools.append({'name': name, 'type': 'builtin', **meta})
        for name, spec in self.plugin_tools.items():
            tools.append({
                'name': name,
                'type': 'plugin',
                'description': spec.get('description', ''),
                'inputs': spec.get('inputs', {}),
                'dangerLevel': spec.get('dangerLevel', 'safe'),
                'version': spec.get('version', '1.0.0'),
            })
        for name, spec in self.instruction_tools.items():
            tools.append({
                'name': name,
                'type': 'instruction',
                'description': spec.get('description', ''),
                'inputs': spec.get('inputs', {}),
                'dangerLevel': 'safe',
                'version': spec.get('version', '1.0.0'),
            })
        for name, spec in self.mcp_tools.items():
            tools.append({
                'name': name,
                'type': 'mcp',
                'server': spec.get('server', ''),
                'description': spec.get('description', ''),
                'inputs': spec.get('inputs', {}),
                'dangerLevel': 'safe',
                'version': '1.0.0',
            })
        return tools


# ============================================
# 🤖 子代理管理器 (Quest 模式支持)
# ============================================

class SubagentManager:
    """
    子代理管理器 - 支持并行探索任务
    用于 Quest 模式的探索阶段，可同时运行多个轻量级代理
    """
    MAX_CONCURRENT = 5  # 最大并发数
    AGENT_TIMEOUT = 30  # 单个代理超时(秒)
    
    def __init__(self, tool_registry: ToolRegistry):
        self.registry = tool_registry
        self.agents: dict[str, dict] = {}  # agent_id -> agent_info
        self.executor = ThreadPoolExecutor(max_workers=self.MAX_CONCURRENT)
        self.futures: dict[str, Future] = {}  # agent_id -> Future
        self.lock = threading.Lock()
    
    def spawn(self, agent_type: str, task: str, tools: list[str], context: str = '') -> str:
        """
        启动一个子代理
        
        Args:
            agent_type: 代理类型 ('explore', 'plan', 'execute')
            task: 任务描述
            tools: 可用工具列表
            context: 上下文信息
        
        Returns:
            agent_id: 代理 ID
        """
        agent_id = f"subagent-{int(time.time()*1000)}-{uuid.uuid4().hex[:6]}"
        
        agent_info = {
            'id': agent_id,
            'type': agent_type,
            'task': task,
            'tools': tools,
            'context': context,
            'status': 'pending',
            'result': None,
            'error': None,
            'started_at': time.time(),
            'completed_at': None,
        }
        
        with self.lock:
            # 检查并发限制
            running_count = sum(1 for a in self.agents.values() if a['status'] == 'running')
            if running_count >= self.MAX_CONCURRENT:
                agent_info['status'] = 'queued'
                agent_info['error'] = f'Queue full, max {self.MAX_CONCURRENT} concurrent agents'
                self.agents[agent_id] = agent_info
                return agent_id
            
            self.agents[agent_id] = agent_info
        
        # 异步执行
        future = self.executor.submit(self._run_agent, agent_id, task, tools, context)
        self.futures[agent_id] = future
        
        with self.lock:
            self.agents[agent_id]['status'] = 'running'
        
        print(f"[SubagentManager] Spawned {agent_type} agent: {agent_id}")
        return agent_id
    
    def _run_agent(self, agent_id: str, task: str, tools: list[str], context: str) -> str:
        """
        执行子代理任务 (简化版 - 单工具调用)
        
        对于探索阶段，每个子代理通常只需要调用一个工具
        """
        try:
            result_parts = []
            
            # 根据任务类型选择工具
            for tool_name in tools:
                if not self.registry.is_registered(tool_name):
                    continue
                
                # 构建工具参数
                args = self._build_tool_args(tool_name, task, context)
                
                # 执行工具
                tool_result = self._execute_tool(tool_name, args)
                
                if tool_result.get('status') == 'success':
                    result_parts.append(f"[{tool_name}] {tool_result.get('result', '')[:1000]}")
                else:
                    result_parts.append(f"[{tool_name}] Error: {tool_result.get('result', 'Unknown error')[:200]}")
            
            final_result = '\n\n'.join(result_parts) if result_parts else 'No tools executed'
            
            with self.lock:
                if agent_id in self.agents:
                    self.agents[agent_id]['status'] = 'completed'
                    self.agents[agent_id]['result'] = final_result
                    self.agents[agent_id]['completed_at'] = time.time()
            
            return final_result
            
        except Exception as e:
            error_msg = str(e)
            with self.lock:
                if agent_id in self.agents:
                    self.agents[agent_id]['status'] = 'failed'
                    self.agents[agent_id]['error'] = error_msg
                    self.agents[agent_id]['completed_at'] = time.time()
            return f"Error: {error_msg}"
    
    def _build_tool_args(self, tool_name: str, task: str, context: str) -> dict:
        """根据工具类型构建参数"""
        # MCP quest 工具的特殊处理
        if tool_name == 'mcp__quest__search_codebase':
            # 提取关键词
            keywords = self._extract_keywords(task)
            return {
                'query': task,
                'key_words': ','.join(keywords[:3]),
                'explanation': f'Exploring: {task[:50]}'
            }
        elif tool_name == 'mcp__quest__search_symbol':
            # 从任务中提取符号名
            symbols = self._extract_symbols(task)
            return {
                'queries': [{'symbol': s, 'relation': 'all'} for s in symbols[:2]],
                'explanation': f'Symbol search for: {task[:50]}'
            }
        elif tool_name == 'readFile':
            # 从上下文中提取文件路径
            paths = self._extract_file_paths(context)
            return {'path': paths[0] if paths else ''}
        elif tool_name == 'listDir':
            return {'path': '.', 'recursive': False}
        else:
            return {'query': task}
    
    def _extract_keywords(self, text: str) -> list[str]:
        """从文本中提取关键词"""
        # 简单实现：提取英文单词和中文词组
        words = re.findall(r'[a-zA-Z_][a-zA-Z0-9_]*|[\u4e00-\u9fff]+', text)
        # 过滤常见词
        stopwords = {'the', 'a', 'an', 'is', 'are', 'to', 'for', 'of', 'in', 'on', 'with', '的', '是', '在', '和', '了'}
        return [w for w in words if w.lower() not in stopwords and len(w) > 1][:5]
    
    def _extract_symbols(self, text: str) -> list[str]:
        """从文本中提取可能的符号名（函数名、类名等）"""
        # 匹配驼峰命名和下划线命名
        symbols = re.findall(r'\b([A-Z][a-zA-Z0-9]*|[a-z][a-zA-Z0-9]*(?:_[a-zA-Z0-9]+)+)\b', text)
        return list(set(symbols))[:3]
    
    def _extract_file_paths(self, text: str) -> list[str]:
        """从文本中提取文件路径"""
        paths = re.findall(r'[a-zA-Z0-9_./\\-]+\.[a-zA-Z]+', text)
        return paths[:3]
    
    def _execute_tool(self, tool_name: str, args: dict) -> dict:
        """执行单个工具"""
        # 检查 MCP 工具
        if tool_name.startswith('mcp__') and self.registry.mcp_manager:
            try:
                result = self.registry.mcp_manager.call_tool(tool_name, args)
                return {'status': 'success', 'result': str(result)[:2000]}
            except Exception as e:
                return {'status': 'error', 'result': str(e)}
        
        # 内置工具需要通过 handler 执行，这里返回占位
        return {'status': 'error', 'result': f'Tool {tool_name} not directly executable in subagent'}
    
    def get_status(self, agent_id: str) -> dict | None:
        """获取子代理状态"""
        with self.lock:
            return self.agents.get(agent_id)
    
    def get_all_status(self) -> list[dict]:
        """获取所有子代理状态"""
        with self.lock:
            return list(self.agents.values())
    
    def collect_results(self, agent_ids: list[str], timeout: float = 60.0) -> list[dict]:
        """
        收集多个子代理的结果
        
        Args:
            agent_ids: 要收集的代理 ID 列表
            timeout: 等待超时时间(秒)
        
        Returns:
            结果列表
        """
        results = []
        deadline = time.time() + timeout
        
        for agent_id in agent_ids:
            remaining = deadline - time.time()
            if remaining <= 0:
                break
            
            future = self.futures.get(agent_id)
            if future:
                try:
                    future.result(timeout=remaining)
                except Exception:
                    pass
            
            with self.lock:
                agent = self.agents.get(agent_id)
                if agent:
                    results.append({
                        'id': agent_id,
                        'type': agent['type'],
                        'task': agent['task'],
                        'status': agent['status'],
                        'result': agent.get('result'),
                        'error': agent.get('error'),
                    })
        
        return results
    
    def cleanup_old_agents(self, max_age: float = 300.0):
        """清理超过指定时间的旧代理记录"""
        cutoff = time.time() - max_age
        with self.lock:
            to_remove = [
                aid for aid, agent in self.agents.items()
                if agent.get('completed_at', 0) < cutoff and agent['status'] in ('completed', 'failed')
            ]
            for aid in to_remove:
                del self.agents[aid]
                self.futures.pop(aid, None)
        
        if to_remove:
            print(f"[SubagentManager] Cleaned up {len(to_remove)} old agents")


# ============================================
# 🌐 浏览器自动化管理器 (Playwright)
# ============================================

class BrowserManager:
    """
    Playwright 浏览器管理器 - 懒启动 + 空闲自动回收
    
    提供持久化浏览器会话，支持跨工具调用复用同一个 page。
    首次调用浏览器工具时启动 Chromium，空闲 5 分钟自动关闭。
    """
    IDLE_TIMEOUT = 300  # 空闲超时 (秒)
    
    def __init__(self):
        self._playwright = None
        self._browser = None
        self._context = None
        self._page = None
        self._lock = threading.Lock()
        self._last_used = 0.0
        self._idle_timer = None
        self._available = None  # None = 未检测, True/False
    
    def is_available(self) -> bool:
        """检测 Playwright 是否可用"""
        if self._available is not None:
            return self._available
        try:
            from playwright.sync_api import sync_playwright
            self._available = True
        except ImportError:
            self._available = False
            print("[BrowserManager] playwright not installed. Run: pip install playwright && playwright install chromium")
        return self._available
    
    def _ensure_browser(self):
        """确保浏览器已启动 (线程安全)"""
        if self._page and not self._page.is_closed():
            self._last_used = time.time()
            return
        
        with self._lock:
            # double check
            if self._page and not self._page.is_closed():
                self._last_used = time.time()
                return
            
            self._cleanup_internal()
            
            from playwright.sync_api import sync_playwright
            
            print("[BrowserManager] Launching Chromium...")
            self._playwright = sync_playwright().start()
            self._browser = self._playwright.chromium.launch(
                headless=True,
                args=[
                    '--no-sandbox',
                    '--disable-blink-features=AutomationControlled',
                ]
            )
            self._context = self._browser.new_context(
                viewport={'width': 1280, 'height': 900},
                user_agent='Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/131.0.0.0 Safari/537.36',
                locale='zh-CN',
            )
            self._page = self._context.new_page()
            self._last_used = time.time()
            self._schedule_idle_check()
            print("[BrowserManager] Browser ready")
    
    def _schedule_idle_check(self):
        """定时检测空闲超时"""
        if self._idle_timer:
            self._idle_timer.cancel()
        self._idle_timer = threading.Timer(60, self._check_idle)
        self._idle_timer.daemon = True
        self._idle_timer.start()
    
    def _check_idle(self):
        """检查空闲超时，自动关闭浏览器"""
        if self._page and (time.time() - self._last_used > self.IDLE_TIMEOUT):
            print("[BrowserManager] Idle timeout, shutting down browser")
            self.shutdown()
        elif self._page:
            self._schedule_idle_check()
    
    def _cleanup_internal(self):
        """内部清理 (不加锁)"""
        try:
            if self._page and not self._page.is_closed():
                self._page.close()
        except Exception:
            pass
        try:
            if self._context:
                self._context.close()
        except Exception:
            pass
        try:
            if self._browser:
                self._browser.close()
        except Exception:
            pass
        try:
            if self._playwright:
                self._playwright.stop()
        except Exception:
            pass
        self._page = None
        self._context = None
        self._browser = None
        self._playwright = None
    
    def shutdown(self):
        """关闭浏览器 (线程安全)"""
        with self._lock:
            if self._idle_timer:
                self._idle_timer.cancel()
                self._idle_timer = None
            self._cleanup_internal()
            print("[BrowserManager] Browser shut down")
    
    # ---- 工具方法 ----
    
    def navigate(self, url: str, wait_until: str = 'domcontentloaded') -> str:
        """导航到 URL，返回页面标题和摘要"""
        self._ensure_browser()
        try:
            self._page.goto(url, wait_until=wait_until, timeout=30000)
            title = self._page.title()
            # 提取可见文本摘要
            text = self._page.evaluate('''() => {
                const sel = document.querySelectorAll('article, main, [role="main"], .content, #content, body');
                const el = sel[0] || document.body;
                return el.innerText.slice(0, 6000);
            }''')
            return json.dumps({
                'status': 'ok',
                'url': self._page.url,
                'title': title,
                'text': text.strip()[:4000] if text else '',
            }, ensure_ascii=False)
        except Exception as e:
            return json.dumps({'status': 'error', 'error': str(e)}, ensure_ascii=False)
    
    def click(self, selector: str) -> str:
        """点击页面元素"""
        self._ensure_browser()
        try:
            self._page.click(selector, timeout=10000)
            self._page.wait_for_load_state('domcontentloaded', timeout=10000)
            title = self._page.title()
            return json.dumps({
                'status': 'ok',
                'url': self._page.url,
                'title': title,
                'message': f'Clicked "{selector}" successfully',
            }, ensure_ascii=False)
        except Exception as e:
            return json.dumps({'status': 'error', 'error': str(e)}, ensure_ascii=False)
    
    def fill(self, selector: str, value: str) -> str:
        """填写表单字段"""
        self._ensure_browser()
        try:
            self._page.fill(selector, value, timeout=10000)
            return json.dumps({
                'status': 'ok',
                'message': f'Filled "{selector}" with value',
            }, ensure_ascii=False)
        except Exception as e:
            return json.dumps({'status': 'error', 'error': str(e)}, ensure_ascii=False)
    
    def extract(self, selector: str = 'body') -> str:
        """提取页面元素文本内容"""
        self._ensure_browser()
        try:
            if selector == 'body':
                text = self._page.evaluate('''() => {
                    const sel = document.querySelectorAll('article, main, [role="main"], .content, #content, body');
                    const el = sel[0] || document.body;
                    return el.innerText;
                }''')
            else:
                el = self._page.query_selector(selector)
                text = el.inner_text() if el else ''
            
            return json.dumps({
                'status': 'ok',
                'url': self._page.url,
                'title': self._page.title(),
                'text': (text or '').strip()[:6000],
            }, ensure_ascii=False)
        except Exception as e:
            return json.dumps({'status': 'error', 'error': str(e)}, ensure_ascii=False)
    
    def screenshot(self, selector: str = None, full_page: bool = False) -> str:
        """截图并返回 base64 编码"""
        self._ensure_browser()
        try:
            import base64
            if selector:
                el = self._page.query_selector(selector)
                if el:
                    img_bytes = el.screenshot()
                else:
                    return json.dumps({'status': 'error', 'error': f'Selector "{selector}" not found'}, ensure_ascii=False)
            else:
                img_bytes = self._page.screenshot(full_page=full_page)
            
            b64 = base64.b64encode(img_bytes).decode('ascii')
            return json.dumps({
                'status': 'ok',
                'url': self._page.url,
                'title': self._page.title(),
                'image_base64': b64[:200] + '...(truncated)',
                'image_size': len(img_bytes),
                'message': f'Screenshot taken ({len(img_bytes)} bytes)',
            }, ensure_ascii=False)
        except Exception as e:
            return json.dumps({'status': 'error', 'error': str(e)}, ensure_ascii=False)
    
    def evaluate(self, expression: str) -> str:
        """在页面上执行 JavaScript 表达式"""
        self._ensure_browser()
        try:
            result = self._page.evaluate(expression)
            return json.dumps({
                'status': 'ok',
                'result': result if isinstance(result, (str, int, float, bool, list, dict, type(None))) else str(result),
            }, ensure_ascii=False)
        except Exception as e:
            return json.dumps({'status': 'error', 'error': str(e)}, ensure_ascii=False)


# 全局浏览器管理器单例
_browser_manager = BrowserManager()


class ClawdDataHandler(BaseHTTPRequestHandler):
    clawd_path = None
    project_path = None  # 项目目录，用于加载内置技能
    registry = None  # type: ToolRegistry
    subagent_manager = None  # type: SubagentManager
    tasks = {}
    tasks_lock = threading.Lock()
    _gene_file_lock = threading.Lock()  # 基因文件读写锁，防止并发写入损坏
    
    def log_message(self, format, *args):
        timestamp = datetime.now().strftime('%H:%M:%S')
        print(f"[{timestamp}] {format % args}")
    
    def send_cors_headers(self):
        self.send_header('Access-Control-Allow-Origin', '*')
        self.send_header('Access-Control-Allow-Methods', 'GET, POST, PUT, DELETE, OPTIONS')
        self.send_header('Access-Control-Allow-Headers', 'Content-Type')
    
    def send_json(self, data, status=200):
        self.send_response(status)
        self.send_header('Content-Type', 'application/json')
        self.send_cors_headers()
        self.end_headers()
        try:
            self.wfile.write(json.dumps(data, ensure_ascii=False, indent=2).encode('utf-8'))
        except (ConnectionAbortedError, ConnectionResetError, BrokenPipeError):
            pass  # 客户端已断开，静默忽略
    
    def send_text(self, text, status=200):
        self.send_response(status)
        self.send_header('Content-Type', 'text/plain; charset=utf-8')
        self.send_cors_headers()
        self.end_headers()
        self.wfile.write(text.encode('utf-8'))
    
    def send_error_json(self, message, status=404):
        self.send_json({'error': message, 'status': 'error'}, status)
    
    def do_OPTIONS(self):
        self.send_response(200)
        self.send_cors_headers()
        self.end_headers()
    
    def do_GET(self):
        parsed = urlparse(self.path)
        path = unquote(parsed.path)
        query = parse_qs(parsed.query)
        
        routes = {
            '/status': self.handle_status,
            '/files': self.handle_files,
            '/skills': self.handle_skills,
            '/nexuses': self.handle_nexuses,
            '/memories': self.handle_memories,
            '/tools': self.handle_tools_list,
            '/all': self.handle_all,
            '/': self.handle_index,
            '': self.handle_index,
        }
        
        if path in routes:
            routes[path]()
        elif path.startswith('/file/'):
            self.handle_file(path[6:])
        elif path.startswith('/nexuses/') and '/experience' not in path:
            nexus_name = path[9:]  # strip '/nexuses/'
            if nexus_name == 'health':
                self.handle_nexuses_health()
            elif nexus_name.endswith('/fitness'):
                self.handle_nexus_fitness_get(nexus_name[:-8])  # strip '/fitness'
            elif nexus_name.endswith('/sop-content'):
                self.handle_nexus_sop_content_get(nexus_name[:-12])  # strip '/sop-content'
            elif nexus_name.endswith('/sop-history'):
                self.handle_nexus_sop_history_get(nexus_name[:-12])  # strip '/sop-history'
            else:
                self.handle_nexus_detail(nexus_name)
        elif path.startswith('/task/status/'):
            task_id = path[13:]
            offset = int(query.get('offset', ['0'])[0])
            self.handle_task_status(task_id, offset)
        elif path == '/api/traces/search':
            self.handle_trace_search(query)
        elif path == '/api/traces/recent':
            self.handle_trace_recent(query)
        elif path == '/api/genes/load':
            self.handle_gene_load()
        elif path == '/api/capsules/load':
            self.handle_capsule_load()
        elif path == '/api/registry/skills':
            self.handle_registry_skills_search(query)
        elif path == '/api/registry/mcp':
            self.handle_registry_mcp_search(query)
        elif path.startswith('/skills/') and path.endswith('/raw'):
            skill_name = path[8:-4]  # strip '/skills/' and '/raw'
            self.handle_skill_raw(skill_name)
        elif path == '/mcp/servers':
            self.handle_mcp_servers_list()
        # V2: Session API
        elif path == '/api/sessions':
            self.handle_sessions_list(query)
        elif path.startswith('/api/sessions/') and path.endswith('/messages'):
            session_id = path[14:-9]  # strip '/api/sessions/' and '/messages'
            self.handle_session_messages_get(session_id, query)
        elif path.startswith('/api/sessions/') and path.endswith('/checkpoint'):
            session_id = path[14:-11]  # strip '/api/sessions/' and '/checkpoint'
            self.handle_session_checkpoint_get(session_id)
        elif path.startswith('/api/sessions/') and not path.endswith('/messages') and not path.endswith('/checkpoint'):
            session_id = path[14:]
            self.handle_session_get(session_id)
        # V2: Memory API
        elif path == '/api/memory/search':
            self.handle_memory_search(query)
        elif path == '/api/memory/stats':
            self.handle_memory_stats()
        elif path.startswith('/api/memory/nexus/'):
            nexus_id = path[18:]
            limit = int(query.get('limit', ['20'])[0])
            self.handle_memory_by_nexus(nexus_id, limit)
        # V2: Scoring API
        elif path.startswith('/api/nexus/') and path.endswith('/scoring'):
            nexus_id = path[11:-8]  # strip '/api/nexus/' and '/scoring'
            self.handle_scoring_get(nexus_id)
        elif path.startswith('/data/'):
            # 前端数据读取 API
            key = path[6:]  # strip '/data/'
            self.handle_data_get(key)
        elif path == '/data':
            # 列出所有数据键
            self.handle_data_list()
        else:
            # 静态文件服务 (托管 dist/ 目录)
            self.serve_static_file(path)
    
    def do_POST(self):
        parsed = urlparse(self.path)
        path = unquote(parsed.path)
        content_type = self.headers.get('Content-Type', '')
        
        # 文件上传：multipart/form-data 单独处理（避免大文件 JSON 编码 OOM）
        if path == '/api/files/upload' and 'multipart/form-data' in content_type:
            self.handle_file_upload_multipart()
            return
        
        content_length = int(self.headers.get('Content-Length', 0))
        body = self.rfile.read(content_length).decode('utf-8') if content_length > 0 else '{}'
        
        try:
            data = json.loads(body) if body else {}
        except json.JSONDecodeError:
            self.send_error_json('Invalid JSON', 400)
            return
        
        # 🌟 新增：工具执行接口
        if path == '/api/tools/execute':
            self.handle_tool_execution(data)
        elif path == '/api/files/upload':
            self.handle_file_upload(data)
        elif path == '/tools/reload':
            self.handle_tools_reload(data)
        elif path == '/api/traces/save':
            self.handle_trace_save(data)
        elif path == '/api/genes/save':
            self.handle_gene_save(data)
        elif path == '/api/capsules/save':
            self.handle_capsule_save(data)
        elif path == '/mcp/reload':
            self.handle_mcp_reload(data)
        elif path.startswith('/mcp/servers/') and path.endswith('/reconnect'):
            server_name = path[13:-10]  # Extract server name
            self.handle_mcp_reconnect(server_name)
        elif path == '/mcp/install':
            self.handle_mcp_install(data)
        elif path == '/skills/install':
            self.handle_skill_install(data)
        elif path == '/skills/uninstall':
            self.handle_skill_uninstall(data)
        elif path == '/clawhub/install':
            self.handle_clawhub_install(data)
        elif path == '/clawhub/publish':
            self.handle_clawhub_publish(data)
        elif path.startswith('/nexuses/') and path.endswith('/skills'):
            nexus_name = path[9:-7]  # strip '/nexuses/' and '/skills'
            self.handle_nexus_update_skills(nexus_name, data)
        elif path.startswith('/nexuses/') and path.endswith('/experience'):
            nexus_name = path[9:-11]  # strip '/nexuses/' and '/experience'
            self.handle_add_experience(nexus_name, data)
        elif path.startswith('/nexuses/') and path.endswith('/meta'):
            nexus_name = path[9:-5]  # strip '/nexuses/' and '/meta'
            self.handle_nexus_update_meta(nexus_name, data)
        elif path.startswith('/nexuses/') and path.endswith('/fitness'):
            nexus_name = path[9:-8]  # strip '/nexuses/' and '/fitness'
            self.handle_nexus_fitness_save(nexus_name, data)
        elif path.startswith('/nexuses/') and path.endswith('/sop-content'):
            nexus_name = path[9:-12]  # strip '/nexuses/' and '/sop-content'
            self.handle_nexus_sop_content_save(nexus_name, data)
        elif path.startswith('/nexuses/') and path.endswith('/sop-history'):
            nexus_name = path[9:-12]  # strip '/nexuses/' and '/sop-history'
            self.handle_nexus_sop_history_save(nexus_name, data)
        elif path == '/task/execute':
            self.handle_task_execute(data)
        # V2: Session API (POST)
        elif path == '/api/sessions':
            self.handle_session_create(data)
        elif path.startswith('/api/sessions/') and path.endswith('/messages'):
            session_id = path[14:-9]
            self.handle_session_message_append(session_id, data)
        elif path.startswith('/api/sessions/') and path.endswith('/checkpoint'):
            session_id = path[14:-11]
            self.handle_session_checkpoint_save(session_id, data)
        elif path.startswith('/api/sessions/') and path.endswith('/meta'):
            session_id = path[14:-5]
            self.handle_session_meta_update(session_id, data)
        # V2: Memory API (POST)
        elif path == '/api/memory/write':
            self.handle_memory_write(data)
        elif path == '/api/memory/write-batch':
            self.handle_memory_write_batch(data)
        elif path == '/api/memory/prune':
            self.handle_memory_prune(data)
        elif path == '/api/memory/decay':
            self.handle_memory_decay(data)
        elif path.startswith('/data/'):
            # 前端数据写入 API
            key = path[6:]  # strip '/data/'
            self.handle_data_set(key, data)
        # 🤖 子代理 API (Quest 模式支持)
        elif path == '/api/subagent/spawn':
            self.handle_subagent_spawn(data)
        elif path == '/api/subagent/collect':
            self.handle_subagent_collect(data)
        elif path.startswith('/api/subagent/') and path.endswith('/status'):
            agent_id = path[14:-7]  # strip '/api/subagent/' and '/status'
            self.handle_subagent_status(agent_id)
        # 🌐 EvoMap 代理 (解决 CORS 问题)
        elif path.startswith('/api/evomap/'):
            self.handle_evomap_proxy(path, data)
        # 🌐 LLM API 代理 (解决 CORS 问题: Moonshot 等 API 的 preflight 不返回 CORS 头)
        elif path == '/api/llm/proxy':
            self.handle_llm_proxy(data)
        else:
            self.send_error_json(f'Unknown endpoint: {path}', 404)
    
    def do_PUT(self):
        parsed = urlparse(self.path)
        path = unquote(parsed.path)
        content_length = int(self.headers.get('Content-Length', 0))
        body = self.rfile.read(content_length).decode('utf-8') if content_length > 0 else '{}'
        try:
            data = json.loads(body) if body else {}
        except json.JSONDecodeError:
            self.send_error_json('Invalid JSON', 400)
            return
        
        if path.startswith('/api/nexus/') and path.endswith('/scoring'):
            nexus_id = path[11:-8]
            self.handle_scoring_put(nexus_id, data)
        else:
            self.send_error_json(f'Unknown PUT endpoint: {path}', 404)
    
    def do_DELETE(self):
        parsed = urlparse(self.path)
        path = unquote(parsed.path)
        
        if path.startswith('/api/sessions/') and path.endswith('/checkpoint'):
            session_id = path[14:-11]
            self.handle_session_checkpoint_delete(session_id)
        elif path.startswith('/api/sessions/'):
            session_id = path[14:]
            self.handle_session_delete(session_id)
        else:
            self.send_error_json(f'Unknown DELETE endpoint: {path}', 404)
    
    # ============================================
    # V2: SQLite API Handlers
    # ============================================

    def _get_db(self) -> sqlite3.Connection:
        global _db_conn
        if _db_conn is None:
            db_path = self.clawd_path / 'ddos_v2.db'
            _db_conn = init_sqlite_db(db_path)
        return _db_conn

    # ---- Sessions ----

    def handle_session_create(self, data: dict):
        """POST /api/sessions - 创建新会话"""
        db = self._get_db()
        session_id = data.get('id') or f"sess-{uuid.uuid4().hex[:12]}"
        title = data.get('title', '')
        sess_type = data.get('type', 'general')
        nexus_id = data.get('nexusId')
        now = int(time.time() * 1000)
        with _db_lock:
            db.execute(
                "INSERT OR IGNORE INTO sessions (id, title, type, nexus_id, created_at, updated_at) VALUES (?,?,?,?,?,?)",
                (session_id, title, sess_type, nexus_id, now, now)
            )
            db.commit()
        self.send_json({'id': session_id, 'title': title, 'type': sess_type, 'nexusId': nexus_id, 'createdAt': now, 'updatedAt': now})

    def handle_sessions_list(self, query: dict):
        """GET /api/sessions - 列出会话"""
        db = self._get_db()
        sess_type = query.get('type', [None])[0]
        nexus_id = query.get('nexusId', [None])[0]
        limit = int(query.get('limit', ['50'])[0])
        offset = int(query.get('offset', ['0'])[0])
        
        sql = "SELECT * FROM sessions WHERE 1=1"
        params = []
        if sess_type:
            sql += " AND type = ?"
            params.append(sess_type)
        if nexus_id:
            sql += " AND nexus_id = ?"
            params.append(nexus_id)
        sql += " ORDER BY updated_at DESC LIMIT ? OFFSET ?"
        params.extend([limit, offset])
        
        rows = db.execute(sql, params).fetchall()
        sessions = [{'id': r['id'], 'title': r['title'], 'type': r['type'], 'nexusId': r['nexus_id'],
                      'createdAt': r['created_at'], 'updatedAt': r['updated_at'],
                      'lastMessagePreview': r['last_message_preview']} for r in rows]
        self.send_json(sessions)

    def handle_session_get(self, session_id: str):
        """GET /api/sessions/{id} - 获取会话详情"""
        db = self._get_db()
        row = db.execute("SELECT * FROM sessions WHERE id = ?", (session_id,)).fetchone()
        if not row:
            self.send_error_json('Session not found', 404)
            return
        messages = db.execute("SELECT * FROM messages WHERE session_id = ? ORDER BY timestamp", (session_id,)).fetchall()
        checkpoint_row = db.execute("SELECT data FROM checkpoints WHERE session_id = ?", (session_id,)).fetchone()
        self.send_json({
            'meta': {'id': row['id'], 'title': row['title'], 'type': row['type'], 'nexusId': row['nexus_id'],
                     'createdAt': row['created_at'], 'updatedAt': row['updated_at']},
            'messages': [{'id': m['id'], 'role': m['role'], 'content': m['content'], 'timestamp': m['timestamp']} for m in messages],
            'checkpoint': json.loads(checkpoint_row['data']) if checkpoint_row else None,
        })

    def handle_session_delete(self, session_id: str):
        """DELETE /api/sessions/{id}"""
        db = self._get_db()
        with _db_lock:
            db.execute("DELETE FROM sessions WHERE id = ?", (session_id,))
            db.commit()
        self.send_json({'status': 'ok'})

    def handle_session_messages_get(self, session_id: str, query: dict):
        """GET /api/sessions/{id}/messages"""
        db = self._get_db()
        limit = int(query.get('limit', ['100'])[0])
        offset = int(query.get('offset', ['0'])[0])
        since = query.get('since', [None])[0]
        
        sql = "SELECT * FROM messages WHERE session_id = ?"
        params: list = [session_id]
        if since:
            sql += " AND timestamp > ?"
            params.append(int(since))
        sql += " ORDER BY timestamp LIMIT ? OFFSET ?"
        params.extend([limit, offset])
        
        rows = db.execute(sql, params).fetchall()
        self.send_json([{'id': r['id'], 'role': r['role'], 'content': r['content'], 'timestamp': r['timestamp']} for r in rows])

    def handle_session_message_append(self, session_id: str, data: dict):
        """POST /api/sessions/{id}/messages"""
        db = self._get_db()
        msg = data.get('message', data)
        msg_id = msg.get('id') or f"msg-{uuid.uuid4().hex[:12]}"
        role = msg.get('role', 'user')
        content = msg.get('content', '')
        timestamp = msg.get('timestamp') or int(time.time() * 1000)
        now = int(time.time() * 1000)
        
        with _db_lock:
            db.execute("INSERT OR IGNORE INTO messages (id, session_id, role, content, timestamp) VALUES (?,?,?,?,?)",
                       (msg_id, session_id, role, content, timestamp))
            db.execute("UPDATE sessions SET updated_at = ?, last_message_preview = ? WHERE id = ?",
                       (now, content[:100], session_id))
            db.commit()
        self.send_json({'status': 'ok', 'id': msg_id})

    def handle_session_meta_update(self, session_id: str, data: dict):
        """POST /api/sessions/{id}/meta"""
        db = self._get_db()
        updates = []
        params = []
        if 'title' in data:
            updates.append("title = ?")
            params.append(data['title'])
        if 'lastMessagePreview' in data:
            updates.append("last_message_preview = ?")
            params.append(data['lastMessagePreview'])
        if updates:
            params.append(int(time.time() * 1000))
            params.append(session_id)
            with _db_lock:
                db.execute(f"UPDATE sessions SET {', '.join(updates)}, updated_at = ? WHERE id = ?", params)
                db.commit()
        self.send_json({'status': 'ok'})

    def handle_session_checkpoint_get(self, session_id: str):
        """GET /api/sessions/{id}/checkpoint"""
        db = self._get_db()
        row = db.execute("SELECT data FROM checkpoints WHERE session_id = ?", (session_id,)).fetchone()
        self.send_json(json.loads(row['data']) if row else None)

    def handle_session_checkpoint_save(self, session_id: str, data: dict):
        """POST /api/sessions/{id}/checkpoint"""
        db = self._get_db()
        now = int(time.time() * 1000)
        with _db_lock:
            db.execute("INSERT OR REPLACE INTO checkpoints (session_id, data, created_at) VALUES (?,?,?)",
                       (session_id, json.dumps(data, ensure_ascii=False), now))
            db.commit()
        self.send_json({'status': 'ok'})

    def handle_session_checkpoint_delete(self, session_id: str):
        """DELETE /api/sessions/{id}/checkpoint"""
        db = self._get_db()
        with _db_lock:
            db.execute("DELETE FROM checkpoints WHERE session_id = ?", (session_id,))
            db.commit()
        self.send_json({'status': 'ok'})

    # ---- Memory ----

    def handle_memory_write(self, data: dict):
        """POST /api/memory/write - 写入单条记忆"""
        db = self._get_db()
        mem_id = f"mem-{uuid.uuid4().hex[:12]}"
        source = data.get('source', 'ephemeral')
        content = data.get('content', '')
        nexus_id = data.get('nexusId')
        tags = json.dumps(data.get('tags', []), ensure_ascii=False)
        metadata = json.dumps(data.get('metadata', {}), ensure_ascii=False)
        confidence = data.get('confidence', 0.5)
        now = int(time.time() * 1000)
        
        with _db_lock:
            db.execute("INSERT INTO memory (id, source, content, nexus_id, tags, metadata, created_at, confidence) VALUES (?,?,?,?,?,?,?,?)",
                       (mem_id, source, content, nexus_id, tags, metadata, now, confidence))
            db.commit()
        self.send_json({'status': 'ok', 'id': mem_id})

    def handle_memory_write_batch(self, data: dict):
        """POST /api/memory/write-batch - 批量写入记忆"""
        db = self._get_db()
        entries = data.get('entries', [])
        count = 0
        now = int(time.time() * 1000)
        with _db_lock:
            for entry in entries:
                mem_id = f"mem-{uuid.uuid4().hex[:12]}"
                db.execute("INSERT INTO memory (id, source, content, nexus_id, tags, metadata, created_at) VALUES (?,?,?,?,?,?,?)",
                           (mem_id, entry.get('source', 'ephemeral'), entry.get('content', ''),
                            entry.get('nexusId'), json.dumps(entry.get('tags', []), ensure_ascii=False),
                            json.dumps(entry.get('metadata', {}), ensure_ascii=False), now))
                count += 1
            db.commit()
        self.send_json({'status': 'ok', 'count': count})

    def handle_memory_search(self, query: dict):
        """GET /api/memory/search?q=xxx&source=xxx&nexusId=xxx&limit=20"""
        db = self._get_db()
        q = query.get('q', [''])[0]
        source = query.get('source', [None])[0]
        nexus_id = query.get('nexusId', [None])[0]
        limit = int(query.get('limit', ['20'])[0])
        
        if q:
            # FTS5 搜索
            fts_sql = """
                SELECT m.*, rank
                FROM memory_fts fts
                JOIN memory m ON m.rowid = fts.rowid
                WHERE memory_fts MATCH ?
            """
            params: list = [q]
            if source:
                fts_sql += " AND m.source = ?"
                params.append(source)
            if nexus_id:
                fts_sql += " AND m.nexus_id = ?"
                params.append(nexus_id)
            fts_sql += " ORDER BY rank LIMIT ?"
            params.append(limit)
            
            try:
                rows = db.execute(fts_sql, params).fetchall()
            except Exception:
                # FTS 查询失败时降级到 LIKE 搜索
                rows = db.execute(
                    "SELECT * FROM memory WHERE content LIKE ? ORDER BY created_at DESC LIMIT ?",
                    (f"%{q}%", limit)
                ).fetchall()
        else:
            sql = "SELECT * FROM memory WHERE 1=1"
            params = []
            if source:
                sql += " AND source = ?"
                params.append(source)
            if nexus_id:
                sql += " AND nexus_id = ?"
                params.append(nexus_id)
            sql += " ORDER BY created_at DESC LIMIT ?"
            params.append(limit)
            rows = db.execute(sql, params).fetchall()
        
        results = [{
            'id': r['id'], 'source': r['source'], 'content': r['content'],
            'nexusId': r['nexus_id'], 'tags': json.loads(r['tags'] or '[]'),
            'metadata': json.loads(r['metadata'] or '{}'), 'createdAt': r['created_at'],
            'score': 1.0,
        } for r in rows]
        self.send_json(results)

    def handle_memory_stats(self):
        """GET /api/memory/stats"""
        db = self._get_db()
        total = db.execute("SELECT COUNT(*) as cnt FROM memory").fetchone()['cnt']
        by_source = {}
        for row in db.execute("SELECT source, COUNT(*) as cnt FROM memory GROUP BY source").fetchall():
            by_source[row['source']] = row['cnt']
        oldest = db.execute("SELECT MIN(created_at) as ts FROM memory").fetchone()['ts']
        newest = db.execute("SELECT MAX(created_at) as ts FROM memory").fetchone()['ts']
        self.send_json({'totalEntries': total, 'bySource': by_source, 'oldestEntry': oldest, 'newestEntry': newest})

    def handle_memory_by_nexus(self, nexus_id: str, limit: int):
        """GET /api/memory/nexus/{nexusId}?limit=20"""
        db = self._get_db()
        rows = db.execute("SELECT * FROM memory WHERE nexus_id = ? ORDER BY created_at DESC LIMIT ?",
                          (nexus_id, limit)).fetchall()
        results = [{
            'id': r['id'], 'source': r['source'], 'content': r['content'],
            'nexusId': r['nexus_id'], 'tags': json.loads(r['tags'] or '[]'),
            'createdAt': r['created_at'], 'score': 1.0,
        } for r in rows]
        self.send_json(results)

    def handle_memory_prune(self, data: dict):
        """POST /api/memory/prune - 清理过期记忆"""
        db = self._get_db()
        older_than_days = data.get('olderThanDays', 30)
        cutoff = int((time.time() - older_than_days * 86400) * 1000)
        with _db_lock:
            cursor = db.execute("DELETE FROM memory WHERE created_at < ?", (cutoff,))
            db.commit()
        self.send_json({'status': 'ok', 'deleted': cursor.rowcount})

    def handle_memory_decay(self, data: dict):
        """POST /api/memory/decay - 批量衰减 L0 记忆置信度"""
        import math
        db = self._get_db()
        half_life_days = data.get('halfLifeDays', 30)
        min_confidence = data.get('minConfidence', 0.05)
        now_ms = int(time.time() * 1000)
        half_life_ms = half_life_days * 86400 * 1000

        with _db_lock:
            rows = db.execute(
                "SELECT id, confidence, created_at FROM memory WHERE source = 'memory' AND confidence > ?",
                (min_confidence,)
            ).fetchall()

            updated = 0
            cleaned = 0
            for row in rows:
                age_ms = now_ms - row['created_at']
                if age_ms <= 0:
                    continue
                decay_factor = math.pow(0.5, age_ms / half_life_ms)
                new_confidence = row['confidence'] * decay_factor

                if new_confidence < min_confidence:
                    db.execute("DELETE FROM memory WHERE id = ?", (row['id'],))
                    cleaned += 1
                else:
                    db.execute("UPDATE memory SET confidence = ? WHERE id = ?",
                               (round(new_confidence, 4), row['id']))
                    updated += 1

            db.commit()

        self.send_json({'status': 'ok', 'updated': updated, 'cleaned': cleaned})

    # ---- Scoring ----

    def handle_scoring_get(self, nexus_id: str):
        """GET /api/nexus/{nexusId}/scoring"""
        db = self._get_db()
        row = db.execute("SELECT scoring_data FROM nexus_scoring WHERE nexus_id = ?", (nexus_id,)).fetchone()
        if row:
            self.send_json(json.loads(row['scoring_data']))
        else:
            # 同时尝试从旧 fitness 文件迁移
            nexus_dir = self._resolve_nexus_dir(nexus_id)
            if nexus_dir:
                fitness_file = nexus_dir / 'sop-fitness.json'
                if fitness_file.exists():
                    try:
                        with fitness_file.open('r', encoding='utf-8') as f:
                            legacy_data = json.load(f)
                        # 迁移到 SQLite
                        now = int(time.time() * 1000)
                        with _db_lock:
                            db.execute("INSERT OR REPLACE INTO nexus_scoring (nexus_id, scoring_data, updated_at) VALUES (?,?,?)",
                                       (nexus_id, json.dumps(legacy_data, ensure_ascii=False), now))
                            db.commit()
                        self.send_json(legacy_data)
                        return
                    except Exception:
                        pass
            self.send_json(None)

    def handle_scoring_put(self, nexus_id: str, data: dict):
        """PUT /api/nexus/{nexusId}/scoring"""
        db = self._get_db()
        now = int(time.time() * 1000)
        scoring_json = json.dumps(data, ensure_ascii=False)
        with _db_lock:
            db.execute("INSERT OR REPLACE INTO nexus_scoring (nexus_id, scoring_data, updated_at) VALUES (?,?,?)",
                       (nexus_id, scoring_json, now))
            db.commit()
        self.send_json({'status': 'ok'})
    
    def serve_static_file(self, path: str):
        """托管 dist/ 目录的前端构建产物，支持 SPA 路由"""
        # 静态文件目录 (与服务器脚本/exe 同级的 dist/)
        static_dir = APP_DIR / 'dist'
        
        if not static_dir.exists():
            # dist/ 不存在时返回提示
            self.send_response(200)
            self.send_header('Content-Type', 'text/html; charset=utf-8')
            self.end_headers()
            self.wfile.write(b'''<!DOCTYPE html>
<html>
<head><title>DD-OS Server</title></head>
<body style="font-family: system-ui; padding: 40px; background: #1a1a2e; color: #eee;">
<h1>DD-OS Native Server</h1>
<p>Frontend not built. Run <code>npm run build</code> to generate dist/</p>
<p>Or access dev server at <a href="http://localhost:5173">http://localhost:5173</a></p>
<hr>
<p>API Endpoints:</p>
<ul>
<li>GET /status - Server status</li>
<li>GET /skills - List skills</li>
<li>POST /api/tools/execute - Execute tool</li>
</ul>
</body>
</html>''')
            return
        
        # 确定文件路径
        if path == '/' or path == '':
            file_path = static_dir / 'index.html'
        else:
            # 去掉开头的 /
            clean_path = path.lstrip('/')
            file_path = static_dir / clean_path
        
        # SPA 路由支持：如果不是文件（没有扩展名），返回 index.html
        if not file_path.exists():
            if '.' not in file_path.name:
                file_path = static_dir / 'index.html'
        
        if not file_path.exists():
            self.send_error_json(f'File not found: {path}', 404)
            return
        
        # 安全检查：确保路径在 static_dir 内
        try:
            file_path.resolve().relative_to(static_dir.resolve())
        except ValueError:
            self.send_error_json('Access denied', 403)
            return
        
        # 获取 MIME 类型
        suffix = file_path.suffix.lower()
        content_type = MIME_TYPES.get(suffix, 'application/octet-stream')
        
        try:
            with open(file_path, 'rb') as f:
                content = f.read()
            
            self.send_response(200)
            self.send_header('Content-Type', content_type)
            self.send_header('Content-Length', str(len(content)))
            # 缓存控制：静态资源长期缓存
            if '/assets/' in str(file_path):
                self.send_header('Cache-Control', 'public, max-age=31536000')
            self.end_headers()
            self.wfile.write(content)
        except Exception as e:
            self.send_error_json(f'Failed to read file: {str(e)}', 500)
    
    # ============================================
    # 📦 前端数据持久化 API (/data)
    # ============================================
    
    def handle_data_get(self, key: str):
        """读取前端数据"""
        data_dir = self.clawd_path / 'data'
        data_dir.mkdir(exist_ok=True)
        
        # 安全检查：允许字母数字下划线中文等 Unicode 字符，禁止路径穿越
        if not key or '..' in key or '/' in key or '\\' in key or len(key) > 200:
            self.send_error_json('Invalid key format', 400)
            return
        
        file_path = data_dir / f'{key}.json'
        
        if not file_path.exists():
            self.send_json({'key': key, 'value': None, 'exists': False})
            return
        
        try:
            content = file_path.read_text(encoding='utf-8')
            self.send_json({'key': key, 'value': json.loads(content), 'exists': True})
        except Exception as e:
            self.send_error_json(f'Failed to read data: {str(e)}', 500)
    
    def handle_data_set(self, key: str, data: dict):
        """写入前端数据"""
        data_dir = self.clawd_path / 'data'
        data_dir.mkdir(exist_ok=True)
        
        # 安全检查：允许字母数字下划线中文等 Unicode 字符，禁止路径穿越
        if not key or '..' in key or '/' in key or '\\' in key or len(key) > 200:
            self.send_error_json('Invalid key format', 400)
            return
        
        file_path = data_dir / f'{key}.json'
        value = data.get('value')
        
        try:
            if value is None:
                # 删除数据
                if file_path.exists():
                    file_path.unlink()
                self.send_json({'key': key, 'deleted': True})
            else:
                # 写入数据
                file_path.write_text(json.dumps(value, ensure_ascii=False, indent=2), encoding='utf-8')
                self.send_json({'key': key, 'saved': True})
        except Exception as e:
            self.send_error_json(f'Failed to save data: {str(e)}', 500)
    
    def handle_data_list(self):
        """列出所有数据键"""
        data_dir = self.clawd_path / 'data'
        data_dir.mkdir(exist_ok=True)
        
        keys = []
        for f in data_dir.glob('*.json'):
            keys.append(f.stem)
        
        self.send_json({'keys': keys})
    
    # ============================================
    # 🛠️ 工具执行 (核心新功能)
    # ============================================
    
    # ---- Layer 1: 前置检查 ----
    
    def _precheck_tool_args(self, tool_name: str, args: dict) -> tuple:
        """Layer 1: 执行前参数校验，返回 (is_valid, error_message)"""
        
        if tool_name in ('writeFile', 'appendFile'):
            path = args.get('path', '')
            content = args.get('content', '')
            if not path:
                return False, f"{tool_name} 缺少 path 参数"
            if not content and content != '':
                return False, f"{tool_name} 缺少 content 参数"
        
        elif tool_name == 'readFile':
            path = args.get('path', '')
            if not path:
                return False, "readFile 缺少 path 参数"
            try:
                file_path = self._resolve_path(path, allow_outside=args.get('allowOutside', False))
                if not file_path.exists():
                    return False, f"文件不存在: {path}。建议: 先用 listDir 确认路径"
                if not file_path.is_file():
                    return False, f"路径不是文件: {path}。建议: 使用 listDir 查看目录内容"
                if file_path.stat().st_size > MAX_FILE_SIZE:
                    return False, f"文件过大 (>{MAX_FILE_SIZE // 1024 // 1024}MB): {path}"
            except PermissionError:
                return False, f"路径越权: {path}。只允许访问工作目录内的文件"
            except ValueError as e:
                return False, str(e)
        
        elif tool_name == 'runCmd':
            command = args.get('command', '')
            if not command:
                return False, "runCmd 缺少 command 参数"
        
        elif tool_name == 'listDir':
            # listDir 允许空 path (默认 .)，无需校验
            pass
        
        elif tool_name == 'generateSkill':
            name = args.get('name', '')
            python_code = args.get('pythonCode', '')
            if not name:
                return False, "generateSkill 缺少 name 参数"
            if not python_code:
                return False, "generateSkill 缺少 pythonCode 参数"
            if 'def main(' not in python_code and 'async def main(' not in python_code:
                return False, "pythonCode 必须包含 main() 函数入口"
        
        elif tool_name == 'parseFile':
            fp = args.get('filePath') or args.get('path', '')
            if not fp:
                return False, "parseFile 缺少 filePath 参数"
        
        return True, ''
    
    # ---- Layer 2: 错误分类 ----
    
    def _classify_error(self, exception: Exception) -> str:
        """Layer 2: 错误类型分类"""
        if isinstance(exception, (UnicodeDecodeError, UnicodeEncodeError)):
            return 'encoding'
        if isinstance(exception, subprocess.TimeoutExpired):
            return 'timeout'
        if isinstance(exception, PermissionError):
            return 'permission'
        if isinstance(exception, (FileNotFoundError, NotADirectoryError)):
            return 'path'
        error_msg = str(exception).lower()
        if 'codec' in error_msg or 'encode' in error_msg or 'decode' in error_msg:
            return 'encoding'
        if 'timeout' in error_msg or 'timed out' in error_msg:
            return 'timeout'
        if 'permission' in error_msg or 'denied' in error_msg:
            return 'permission'
        if 'not found' in error_msg or 'no such file' in error_msg:
            return 'path'
        return 'unknown'
    
    _ERROR_SUGGESTIONS = {
        'encoding': '建议: 检查文件编码是否为 UTF-8，或命令输出是否包含特殊字符',
        'timeout': '建议: 增加 timeout 参数值，或简化命令/操作',
        'permission': '建议: 检查路径权限，避免访问系统目录',
        'path': '建议: 先用 listDir 确认路径存在，检查拼写是否正确',
    }
    
    # ---- Layer 3: 结果验证 ----
    
    def _verify_tool_result(self, tool_name: str, args: dict, result: str, status: str) -> dict:
        """Layer 3: 工具结果的代码验证"""
        if status == 'error':
            return {'verified': False, 'checks': [], 'confidence': 0.0}
        
        checks = []
        
        if tool_name == 'writeFile':
            path = args.get('path', '')
            content = args.get('content', '')
            try:
                file_path = self._resolve_path(path)
                # Check 1: 文件存在性
                exists = file_path.exists()
                checks.append({
                    'name': '文件存在性',
                    'passed': exists,
                    'details': f'{file_path.name} {"存在" if exists else "不存在"}'
                })
                if exists:
                    # Check 2: 大小匹配
                    actual_size = file_path.stat().st_size
                    expected_size = len(content.encode('utf-8'))
                    size_match = abs(actual_size - expected_size) <= 10  # 允许微小差异
                    checks.append({
                        'name': '大小匹配',
                        'passed': size_match,
                        'details': f'实际 {actual_size}B vs 预期 {expected_size}B'
                    })
            except Exception:
                checks.append({'name': '验证异常', 'passed': False, 'details': '验证过程出错'})
        
        elif tool_name == 'generateSkill':
            name = args.get('name', '')
            nexus_id = args.get('nexusId', '')
            safe_name = re.sub(r'[^\w-]', '-', name.lower()).strip('-')
            safe_name = re.sub(r'-+', '-', safe_name)
            
            if nexus_id:
                skill_dir = self.clawd_path / 'nexuses' / nexus_id / 'skills' / safe_name
            else:
                skill_dir = self.clawd_path / 'skills' / safe_name
            
            # Check 1: SKILL.md 存在
            skill_md = skill_dir / 'SKILL.md'
            checks.append({
                'name': 'SKILL.md 存在',
                'passed': skill_md.exists(),
                'details': f'{skill_md.name} {"已创建" if skill_md.exists() else "未找到"}'
            })
            # Check 2: Python 文件存在
            py_file = skill_dir / f'{safe_name}.py'
            checks.append({
                'name': 'Python 文件存在',
                'passed': py_file.exists(),
                'details': f'{py_file.name} {"已创建" if py_file.exists() else "未找到"}'
            })
        
        elif tool_name == 'runCmd':
            # Check: 输出中是否有替代字符 (编码问题指标)
            replace_count = result.count('\ufffd')
            total_chars = max(len(result), 1)
            replace_ratio = replace_count / total_chars
            encoding_ok = replace_ratio < 0.05
            checks.append({
                'name': '输出编码质量',
                'passed': encoding_ok,
                'details': f'替代字符占比 {replace_ratio:.1%}' if not encoding_ok else '编码正常'
            })
        
        elif tool_name == 'readFile':
            # Check: 返回内容非空
            has_content = bool(result and result.strip())
            checks.append({
                'name': '内容非空',
                'passed': has_content,
                'details': f'{len(result)} 字符' if has_content else '文件内容为空'
            })
        
        elif tool_name == 'appendFile':
            path = args.get('path', '')
            try:
                file_path = self._resolve_path(path)
                exists = file_path.exists()
                checks.append({
                    'name': '文件存在性',
                    'passed': exists,
                    'details': f'{file_path.name} {"存在" if exists else "不存在"}'
                })
            except Exception:
                checks.append({'name': '验证异常', 'passed': False, 'details': '验证过程出错'})
        
        # 计算 confidence
        if not checks:
            return {'verified': True, 'checks': [], 'confidence': 0.95}
        
        passed_count = sum(1 for c in checks if c['passed'])
        confidence = passed_count / len(checks)
        
        return {
            'verified': all(c['passed'] for c in checks),
            'checks': checks,
            'confidence': round(confidence, 2)
        }
    
    # ---- 工具执行主入口 ----
    
    def handle_tool_execution(self, data):
        """处理工具调用请求 - 支持内置工具、插件工具、指令型工具和MCP工具"""
        tool_name = data.get('name', '')
        args = data.get('args', {})

        if not self.registry.is_registered(tool_name):
            all_tools = [t['name'] for t in self.registry.list_all()]
            self.send_json({
                'tool': tool_name,
                'status': 'error',
                'result': f'Tool not registered: {tool_name}. Available: {", ".join(all_tools)}'
            }, 403)
            return

        # Layer 1: 前置检查
        is_valid, precheck_error = self._precheck_tool_args(tool_name, args)
        if not is_valid:
            self.send_json({
                'tool': tool_name,
                'status': 'error',
                'result': f'[前置检查失败] {precheck_error}',
                'error_type': 'precheck_failure',
                'timestamp': datetime.now().isoformat()
            })
            return

        result = ""
        status = "success"
        error_type = None
        start_time = time.time()

        try:
            # 1. 指令型工具 -> 路由到 skill-executor
            instruction_spec = self.registry.get_instruction(tool_name)
            if instruction_spec:
                result = self._execute_instruction_tool(instruction_spec, tool_name, args)
            # 2. 插件工具 -> subprocess 执行
            elif self.registry.get_plugin(tool_name):
                plugin_spec = self.registry.get_plugin(tool_name)
                result = self._execute_plugin_tool(plugin_spec, tool_name, args)
            # 3. MCP 工具 -> 通过 MCPManager 调用
            elif self.registry.get_mcp_tool(tool_name):
                result = self._execute_mcp_tool(tool_name, args)
            # 4. 内置工具 -> 直接调度
            else:
                builtin_handlers = {
                    'readFile': self._tool_read_file,
                    'writeFile': self._tool_write_file,
                    'appendFile': self._tool_append_file,
                    'listDir': self._tool_list_dir,
                    'runCmd': self._tool_run_cmd,
                    'weather': self._tool_weather,
                    'webSearch': self._tool_web_search,
                    'webFetch': self._tool_web_fetch,
                    'saveMemory': self._tool_save_memory,
                    'searchMemory': self._tool_search_memory,
                    'nexusBindSkill': self._tool_nexus_bind_skill,
                    'nexusUnbindSkill': self._tool_nexus_unbind_skill,
                    'openInExplorer': self._tool_open_in_explorer,
                    'parseFile': self._tool_parse_file,
                    'generateSkill': self._tool_generate_skill,
                    'browser_navigate': self._tool_browser_navigate,
                    'browser_click': self._tool_browser_click,
                    'browser_fill': self._tool_browser_fill,
                    'browser_extract': self._tool_browser_extract,
                    'browser_screenshot': self._tool_browser_screenshot,
                    'browser_evaluate': self._tool_browser_evaluate,
                }
                handler = builtin_handlers.get(tool_name)
                if handler:
                    result = handler(args)
                else:
                    raise ValueError(f"No handler for builtin tool: {tool_name}")

        except Exception as e:
            status = "error"
            # Layer 2: 错误分类 + 增强信息
            error_type = self._classify_error(e)
            result = f"Tool execution failed: {str(e)}"
            suggestion = self._ERROR_SUGGESTIONS.get(error_type)
            if suggestion:
                result += f'\n{suggestion}'

        execution_time_ms = int((time.time() - start_time) * 1000)

        # Layer 3: 结果验证
        verification = self._verify_tool_result(tool_name, args, result, status)

        response = {
            'tool': tool_name,
            'status': status,
            'result': result,
            'timestamp': datetime.now().isoformat(),
            'verification': verification,
            'execution_time_ms': execution_time_ms,
        }
        if error_type:
            response['error_type'] = error_type

        self.send_json(response)

    # ============================================
    # 📎 文件上传 + 自动解析
    # ============================================

    UPLOAD_ALLOWED_EXT = {'.pdf', '.docx', '.pptx', '.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.webp', '.txt', '.md', '.csv'}

    def handle_file_upload_multipart(self):
        """接收 FormData multipart 上传，保存并自动解析
        
        使用手动 multipart boundary 解析，不依赖已废弃的 cgi.FieldStorage
        （cgi 在 Python 3.11+ deprecated，3.13 removed）
        """
        content_type = self.headers.get('Content-Type', '')
        content_length = int(self.headers.get('Content-Length', 0))

        if content_length > MAX_FILE_SIZE:
            # 消耗请求体避免连接异常
            remaining = content_length
            while remaining > 0:
                chunk = min(remaining, 65536)
                self.rfile.read(chunk)
                remaining -= chunk
            self.send_error_json(f'文件过大 (>{MAX_FILE_SIZE // 1024 // 1024}MB)', 413)
            return

        # 从 Content-Type 提取 boundary
        boundary = None
        for part in content_type.split(';'):
            part = part.strip()
            if part.startswith('boundary='):
                boundary = part[len('boundary='):].strip('"')
        if not boundary:
            self.send_error_json('无效的 multipart 请求：缺少 boundary', 400)
            return

        # 读取整个请求体（已验证 content_length <= 10MB，内存安全）
        try:
            raw_body = self.rfile.read(content_length)
        except Exception as e:
            print(f"[ERROR] 读取请求体失败: {e}", file=sys.stderr)
            self.send_error_json('读取上传数据失败', 400)
            return

        # 按 boundary 分割，提取包含 filename 的 part
        boundary_bytes = ('--' + boundary).encode()
        parts = raw_body.split(boundary_bytes)

        file_bytes = None
        file_name = 'unknown'
        for part_data in parts:
            if b'filename=' not in part_data:
                continue
            # headers 和 body 以空行 (\r\n\r\n) 分隔
            header_end = part_data.find(b'\r\n\r\n')
            if header_end == -1:
                continue
            headers_raw_bytes = part_data[:header_end]
            # 尝试 UTF-8（浏览器 FormData），回退到 GBK（Windows curl/工具）
            try:
                headers_raw = headers_raw_bytes.decode('utf-8')
            except UnicodeDecodeError:
                headers_raw = headers_raw_bytes.decode('gbk', errors='replace')
            file_bytes = part_data[header_end + 4:]
            # 去掉尾部的 \r\n（multipart 格式约定）
            if file_bytes.endswith(b'\r\n'):
                file_bytes = file_bytes[:-2]
            # 从 Content-Disposition 提取 filename
            for line in headers_raw.split('\r\n'):
                if 'filename=' in line:
                    # 支持: filename="中文.pptx" 和 filename=file.pdf
                    match = re.search(r'filename="?([^";\r\n]+)"?', line)
                    if match:
                        file_name = match.group(1).strip()
            break  # 只取第一个文件

        if file_bytes is None:
            self.send_error_json('未找到上传文件', 400)
            return

        # 清理文件名（保留中文、字母、数字、点、横线）
        safe_name = re.sub(r'[^\w.\-\u4e00-\u9fff]', '_', file_name)
        # 文件名长度限制（NTFS/ext4 最大 255 字符）
        stem, ext = os.path.splitext(safe_name)
        ext = ext.lower()
        if len(safe_name) > 200:
            safe_name = stem[:200 - len(ext)] + ext

        if ext not in self.UPLOAD_ALLOWED_EXT:
            self.send_error_json(f'不支持的文件类型: {ext}，支持: {", ".join(sorted(self.UPLOAD_ALLOWED_EXT))}', 400)
            return

        if len(file_bytes) > MAX_FILE_SIZE:
            self.send_error_json(f'文件过大 (>{MAX_FILE_SIZE // 1024 // 1024}MB)', 413)
            return

        # 保存到临时目录
        upload_dir = self.clawd_path / 'temp' / 'uploads'
        upload_dir.mkdir(parents=True, exist_ok=True)
        unique_name = f"{uuid.uuid4().hex[:8]}_{safe_name}"
        file_path = upload_dir / unique_name

        try:
            file_path.write_bytes(file_bytes)
        except Exception as e:
            print(f"[ERROR] 文件保存失败: {e}", file=sys.stderr)
            self.send_error_json('文件保存失败', 500)
            return

        # 自动解析
        parsed_text = ''
        try:
            parsed_text = self._tool_parse_file({'filePath': str(file_path)})
        except Exception as e:
            print(f"[ERROR] 文件解析失败: {e}", file=sys.stderr)
            err_msg = str(e)
            if '未安装' in err_msg or 'pip install' in err_msg:
                parsed_text = f'[解析失败: {err_msg}]'
            else:
                parsed_text = f'[解析失败: 请检查文件格式是否正确]'

        file_size = len(file_bytes)
        self.send_json({
            'success': True,
            'filePath': str(file_path),
            'originalName': file_name,
            'fileSize': file_size,
            'parsedText': parsed_text,
            'timestamp': datetime.now().isoformat()
        })

    def handle_file_upload(self, data: dict):
        """接收前端上传的文件（Base64），保存到临时目录并自动解析"""
        file_name = data.get('fileName', '')
        data_base64 = data.get('dataBase64', '')

        if not file_name or not data_base64:
            self.send_error_json('fileName and dataBase64 are required', 400)
            return

        # 清理文件名
        safe_name = re.sub(r'[^\w.\-\u4e00-\u9fff]', '_', file_name)
        ext = os.path.splitext(safe_name)[1].lower()

        if ext not in self.UPLOAD_ALLOWED_EXT:
            self.send_error_json(f'不支持的文件类型: {ext}，支持: {", ".join(sorted(self.UPLOAD_ALLOWED_EXT))}', 400)
            return

        # 解码 Base64 (去掉 data:xxx;base64, 前缀)
        try:
            if ';base64,' in data_base64:
                data_base64 = data_base64.split(';base64,')[1]
            file_bytes = base64.b64decode(data_base64)
        except Exception as e:
            self.send_error_json(f'Base64 解码失败: {str(e)}', 400)
            return

        if len(file_bytes) > MAX_FILE_SIZE:
            self.send_error_json(f'文件过大 (>{MAX_FILE_SIZE // 1024 // 1024}MB)', 413)
            return

        # 保存到临时目录
        upload_dir = self.clawd_path / 'temp' / 'uploads'
        upload_dir.mkdir(parents=True, exist_ok=True)
        unique_name = f"{uuid.uuid4().hex[:8]}_{safe_name}"
        file_path = upload_dir / unique_name

        try:
            file_path.write_bytes(file_bytes)
        except Exception as e:
            self.send_error_json(f'文件保存失败: {str(e)}', 500)
            return

        # 自动解析
        parsed_text = ''
        try:
            parsed_text = self._tool_parse_file({'filePath': str(file_path)})
        except Exception as e:
            parsed_text = f'[解析失败: {str(e)}]'

        self.send_json({
            'success': True,
            'filePath': str(file_path),
            'originalName': file_name,
            'parsedText': parsed_text,
            'timestamp': datetime.now().isoformat()
        })

    def _execute_plugin_tool(self, spec: dict, tool_name: str, args: dict) -> str:
        """执行插件工具 - subprocess 隔离执行"""
        exe_path = spec['exe_path']
        runtime = spec.get('runtime', 'python')

        # 确定运行时命令
        if runtime == 'python':
            cmd = [sys.executable, exe_path]
        elif runtime == 'node':
            cmd = ['node', exe_path]
        else:
            raise ValueError(f"Unsupported runtime: {runtime}")

        # 构建输入：包含工具名和参数（支持多工具 manifest）
        input_data = json.dumps({
            'tool': tool_name,
            'args': args
        }, ensure_ascii=False)

        try:
            process = subprocess.run(
                cmd,
                input=input_data,
                capture_output=True,
                text=True,
                timeout=PLUGIN_TIMEOUT,
                cwd=spec.get('skill_dir', str(self.clawd_path)),
            )

            if process.returncode != 0:
                stderr = process.stderr[:MAX_OUTPUT_SIZE] if process.stderr else ''
                raise RuntimeError(f"Plugin exited with code {process.returncode}: {stderr}")

            return process.stdout[:MAX_OUTPUT_SIZE] if process.stdout else ''

        except subprocess.TimeoutExpired:
            raise RuntimeError(f"Plugin timed out after {PLUGIN_TIMEOUT}s")

    def _execute_instruction_tool(self, spec: dict, tool_name: str, args: dict) -> str:
        """执行指令型工具 - 通过 skill-executor 解析 SKILL.md 并返回指令"""
        skill_executor = self.clawd_path / 'skills' / 'skill-executor' / 'execute.py'

        if not skill_executor.exists():
            raise RuntimeError(f"skill-executor not found at {skill_executor}")

        # 使用 original_name (kebab-case) 让 SkillDiscovery 能找到目录
        original_name = spec.get('original_name', tool_name)

        input_data = json.dumps({
            'tool': 'run_skill',
            'args': {
                'skill_name': original_name,
                'args': args,
                'project_root': str(self.clawd_path),
            }
        }, ensure_ascii=False)

        try:
            process = subprocess.run(
                [sys.executable, str(skill_executor)],
                input=input_data,
                capture_output=True,
                text=True,
                timeout=PLUGIN_TIMEOUT,
                cwd=str(skill_executor.parent),
            )

            if process.returncode != 0:
                stderr = process.stderr[:MAX_OUTPUT_SIZE] if process.stderr else ''
                raise RuntimeError(f"Instruction skill error: {stderr}")

            result = json.loads(process.stdout)
            if not result.get('success'):
                raise RuntimeError(result.get('error', 'Unknown error'))

            return result.get('instructions', result.get('output', ''))

        except subprocess.TimeoutExpired:
            raise RuntimeError(f"Instruction skill timed out after {PLUGIN_TIMEOUT}s")
        except json.JSONDecodeError:
            # skill-executor 返回非 JSON 时，直接返回原文
            return process.stdout[:MAX_OUTPUT_SIZE] if process.stdout else ''

    def _execute_mcp_tool(self, tool_name: str, args: dict) -> str:
        """执行 MCP 工具 - 通过 MCPManager 调用远程 MCP 服务器"""
        if not self.registry.mcp_manager:
            raise RuntimeError("MCP manager not initialized")

        try:
            result = self.registry.mcp_manager.call_tool(tool_name, args, timeout=PLUGIN_TIMEOUT)
            if result is None:
                return ""
            return str(result)
        except Exception as e:
            raise RuntimeError(f"MCP tool execution failed: {e}")
    
    def _resolve_path(self, relative_path: str, allow_outside: bool = False) -> Path:
        """解析并验证路径安全性"""
        if not relative_path:
            raise ValueError("Path cannot be empty")
        
        # 移除开头的斜杠
        clean_path = relative_path.lstrip('/')
        
        # 默认在 clawd 目录下操作
        if allow_outside and os.path.isabs(relative_path):
            file_path = Path(relative_path)
        else:
            file_path = self.clawd_path / clean_path
        
        # 安全检查：防止路径遍历
        try:
            resolved = file_path.resolve()
            if not allow_outside:
                resolved.relative_to(self.clawd_path.resolve())
        except ValueError:
            raise PermissionError(f"Access denied: path outside allowed directory")
        
        return resolved
    
    def _tool_read_file(self, args: dict) -> str:
        """读取文件内容"""
        path = args.get('path', '')
        file_path = self._resolve_path(path, allow_outside=args.get('allowOutside', False))
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {path}")
        if not file_path.is_file():
            raise ValueError(f"Not a file: {path}")
        if file_path.stat().st_size > MAX_FILE_SIZE:
            raise ValueError(f"File too large (>{MAX_FILE_SIZE} bytes)")
        
        try:
            return file_path.read_text(encoding='utf-8')
        except UnicodeDecodeError:
            content = file_path.read_text(encoding='utf-8', errors='replace')
            return f"[注意: 文件包含非UTF-8字符，已用替代字符显示]\n{content}"
    
    def _tool_parse_file(self, args: dict) -> str:
        """解析文档或图像文件，返回提取的文本内容"""
        file_path_str = args.get('filePath') or args.get('path', '')
        if not file_path_str:
            raise ValueError("filePath is required")
        
        # 支持绝对路径（上传的临时文件）和相对路径
        if os.path.isabs(file_path_str):
            file_path = Path(file_path_str).resolve()
            # 安全检查：只允许访问 clawd 工作目录下的文件
            allowed_root = self.clawd_path.resolve()
            try:
                file_path.relative_to(allowed_root)
            except ValueError:
                raise PermissionError(f"Access denied: path outside allowed directory")
        else:
            file_path = self._resolve_path(file_path_str, allow_outside=True)
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path_str}")
        if file_path.stat().st_size > MAX_FILE_SIZE:
            raise ValueError(f"File too large (>{MAX_FILE_SIZE // 1024 // 1024}MB)")
        
        ext = file_path.suffix.lower()
        text = ""
        
        if ext == '.pdf':
            if not HAS_PDF:
                raise RuntimeError("pdfplumber 未安装，请运行 pip install pdfplumber")
            with pdfplumber.open(str(file_path)) as pdf:
                pages = []
                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text() or ''
                    if page_text.strip():
                        pages.append(f"--- 第{i+1}页 ---\n{page_text}")
                text = "\n\n".join(pages)
        
        elif ext == '.docx':
            if not HAS_DOCX:
                raise RuntimeError("python-docx 未安装，请运行 pip install python-docx")
            doc = DocxDocument(str(file_path))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    paragraphs.append(" | ".join(cells))
            text = "\n".join(paragraphs)
        
        elif ext == '.pptx':
            if not HAS_PPTX:
                raise RuntimeError("python-pptx 未安装，请运行 pip install python-pptx")
            prs = PptxPresentation(str(file_path))
            slides = []
            for i, slide in enumerate(prs.slides):
                parts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                parts.append(t)
                if parts:
                    slides.append(f"--- 幻灯片{i+1} ---\n" + "\n".join(parts))
            text = "\n\n".join(slides)
        
        elif ext in ('.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.webp'):
            if not HAS_OCR:
                raise RuntimeError("pytesseract/Pillow 未安装，请运行 pip install pytesseract Pillow")
            img = Image.open(str(file_path))
            lang = args.get('language', 'eng+chi_sim')
            text = pytesseract.image_to_string(img, lang=lang)
        
        else:
            # 回退：尝试当纯文本读取
            try:
                text = file_path.read_text(encoding='utf-8')
            except Exception:
                raise ValueError(f"不支持的文件类型: {ext}")
        
        if not text.strip():
            return f"[文件 {file_path.name} 无可提取的文本内容]"
        
        # 截断到 MAX_OUTPUT_SIZE（安全 UTF-8 边界截断）
        encoded = text.encode('utf-8')
        if len(encoded) > MAX_OUTPUT_SIZE:
            safe_idx = MAX_OUTPUT_SIZE
            # 回退到 UTF-8 字符边界，避免截断多字节字符导致乱码
            while safe_idx > 0 and (encoded[safe_idx] & 0xC0) == 0x80:
                safe_idx -= 1
            text = encoded[:safe_idx].decode('utf-8')
            text += f"\n\n[内容过长，已截断至约 {MAX_OUTPUT_SIZE // 1024}KB]"
        
        return text
    
    def _tool_write_file(self, args: dict) -> str:
        """写入文件"""
        path = args.get('path', '')
        content = args.get('content', '')
        
        file_path = self._resolve_path(path)
        
        # === Nexus 涌现去重网关 ===
        if 'nexuses/' in path and path.endswith('NEXUS.md'):
            # 仅在文件不存在时（即新建操作）进行去重检查
            if not file_path.exists():
                duplicate_id = self._check_nexus_duplication(content)
                if duplicate_id:
                    return (f"【系统拦截】创建失败！\n"
                            f"检测到高度相似的 Nexus 节点已存在 (节点 ID: {duplicate_id})。\n"
                            f"为避免知识图谱碎片化，请不要创建新目录，请直接使用 'readFile' 和 'writeFile' "
                            f"读取并更新原有的 nexuses/{duplicate_id}/NEXUS.md，或者向其追加 experience。")
        
        # === Nexus 格式引导 ===
        # 检测写入 nexuses/ 目录但不是 NEXUS.md 的情况，提供格式纠正提示
        if 'nexuses/' in path and not path.endswith('NEXUS.md'):
            # 提取可能的 nexus id
            import re
            nexus_match = re.search(r'nexuses/([^/]+)', path)
            nexus_id = nexus_match.group(1) if nexus_match else 'your-nexus-id'
            
            # 如果是写入 .json 或其他配置文件，返回警告并引导正确格式
            if path.endswith('.json') or (path.endswith('.md') and 'NEXUS.md' not in path):
                return (f"【格式提示】检测到你正在向 nexuses/ 目录写入非标准文件。\n\n"
                        f"⚠️ Nexus 只能通过 NEXUS.md 文件定义，系统不会识别 .json 或其他 .md 文件！\n\n"
                        f"📝 正确做法：请创建 nexuses/{nexus_id}/NEXUS.md 文件，格式如下：\n"
                        f"```markdown\n"
                        f"---\n"
                        f"name: Nexus名称\n"
                        f"description: 功能描述\n"
                        f"version: 1.0.0\n"
                        f"skill_dependencies:\n"
                        f"  - 技能ID\n"
                        f"tags:\n"
                        f"  - 标签\n"
                        f"triggers:\n"
                        f"  - 触发词\n"
                        f"objective: 核心目标\n"
                        f"metrics:\n"
                        f"  - 质量指标\n"
                        f"strategy: 执行策略\n"
                        f"---\n\n"
                        f"# Nexus名称 SOP\n\n"
                        f"（详细的标准作业程序）\n"
                        f"```\n\n"
                        f"请使用正确格式重新创建 nexuses/{nexus_id}/NEXUS.md")
        
        # 确保父目录存在
        file_path.parent.mkdir(parents=True, exist_ok=True)
        
        file_path.write_text(content, encoding='utf-8')
        
        # 返回结构化数据，包含完整路径以便前端快速访问
        return json.dumps({
            'action': 'file_created',
            'message': f'已成功写入 {len(content)} 字节',
            'fileName': file_path.name,
            'filePath': str(file_path.resolve()),
            'fileSize': len(content),
        }, ensure_ascii=False)
    
    def _tool_open_in_explorer(self, args: dict) -> str:
        """在文件管理器中打开指定路径并高亮文件"""
        path = args.get('path', '')
        if not path:
            raise ValueError("路径参数不能为空")
        
        file_path = Path(path)
        if not file_path.exists():
            raise FileNotFoundError(f"文件不存在: {path}")
        
        import platform
        import subprocess
        system = platform.system()
        
        try:
            if system == 'Windows':
                # Windows: 使用 explorer /select 高亮文件
                subprocess.run(['explorer', '/select,', str(file_path.resolve())], check=False)
            elif system == 'Darwin':  # macOS
                subprocess.run(['open', '-R', str(file_path.resolve())], check=True)
            else:  # Linux
                # 打开父目录
                subprocess.run(['xdg-open', str(file_path.parent.resolve())], check=True)
            
            return f"已在文件管理器中打开: {file_path.name}"
        except Exception as e:
            raise RuntimeError(f"无法打开文件管理器: {str(e)}")
    
    def _check_nexus_duplication(self, new_content: str) -> str | None:
        """检查新建的 Nexus 是否与现存 Nexus 重复，返回重复的 Nexus ID"""
        # 1. 提取新 Nexus 的 frontmatter
        match = re.match(r'^---\s*\n(.*?)\n---\s*\n', new_content, re.DOTALL)
        if not match:
            return None
        
        new_meta = {}
        if HAS_YAML:
            try:
                new_meta = yaml.safe_load(match.group(1)) or {}
            except Exception:
                pass
        else:
            for line in match.group(1).split('\n'):
                m = re.match(r'^(\w+)\s*:\s*(.+)$', line.strip())
                if m:
                    new_meta[m.group(1)] = m.group(2).strip()
        
        new_name = str(new_meta.get('name', ''))
        new_desc = str(new_meta.get('description', ''))
        if not new_name and not new_desc:
            return None
        
        new_text = f"{new_name} {new_desc}"
        
        # 2. 遍历现有 Nexus 进行对比
        nexuses_dir = self.clawd_path / 'nexuses'
        if not nexuses_dir.exists():
            return None
        
        best_match = None
        highest_score = 0.0
        
        for nexus_md in nexuses_dir.rglob('NEXUS.md'):
            existing_meta = parse_nexus_frontmatter(nexus_md)
            ext_name = str(existing_meta.get('name', ''))
            ext_desc = str(existing_meta.get('description', ''))
            
            ext_text = f"{ext_name} {ext_desc}"
            score = calculate_text_similarity(new_text, ext_text)
            
            if score > highest_score:
                highest_score = score
                best_match = nexus_md.parent.name
        
        # 阈值：超过 55% 的特征重合即判定为重复
        if highest_score >= 0.55:
            return best_match
        
        return None
    
    def _tool_append_file(self, args: dict) -> str:
        """追加内容到文件"""
        path = args.get('path', '')
        content = args.get('content', '')
        
        file_path = self._resolve_path(path)
        file_path.parent.mkdir(parents=True, exist_ok=True)
        
        with open(file_path, 'a', encoding='utf-8') as f:
            f.write(content)
        
        return f"Appended {len(content)} bytes to {file_path.name}"
    
    def _tool_list_dir(self, args: dict) -> str:
        """列出目录内容"""
        path = args.get('path', '.')
        dir_path = self._resolve_path(path)
        
        if not dir_path.exists():
            raise FileNotFoundError(f"Directory not found: {path}")
        if not dir_path.is_dir():
            raise ValueError(f"Not a directory: {path}")
        
        items = []
        for item in sorted(dir_path.iterdir()):
            item_type = 'dir' if item.is_dir() else 'file'
            size = item.stat().st_size if item.is_file() else 0
            items.append({
                'name': item.name,
                'type': item_type,
                'size': size
            })
        
        return json.dumps(items, ensure_ascii=False)
    
    def _tool_run_cmd(self, args: dict) -> str:
        """执行 Shell 命令 (⚠️ 高危操作)"""
        command = args.get('command', '')
        cwd = args.get('cwd', str(self.clawd_path))
        timeout = min(args.get('timeout', 60), 300)  # 最大 5 分钟
        
        if not command:
            raise ValueError("Command cannot be empty")
        
        # 安全检查
        cmd_lower = command.lower()
        for dangerous in DANGEROUS_COMMANDS:
            if dangerous in cmd_lower:
                raise PermissionError(f"Dangerous command blocked: {command}")
        
        try:
            process = subprocess.run(
                command,
                shell=True,
                cwd=cwd,
                capture_output=True,
                text=True,
                encoding='utf-8',
                errors='replace',
                timeout=timeout
            )
            
            stdout = safe_utf8_truncate(process.stdout, MAX_OUTPUT_SIZE) if process.stdout else ''
            stderr = safe_utf8_truncate(process.stderr, MAX_OUTPUT_SIZE) if process.stderr else ''
            
            result_parts = []
            if stdout:
                result_parts.append(f"STDOUT:\n{stdout}")
            if stderr:
                result_parts.append(f"STDERR:\n{stderr}")
            
            rc = process.returncode
            if rc == 0:
                result_parts.append(f"Exit Code: 0 (成功)")
            else:
                # 提供常见 exit code 的可读解释
                code_hints = {
                    1: "通用错误",
                    2: "参数错误或命令误用",
                    3: "URL 格式错误 (curl)",
                    6: "无法解析主机名 (DNS 失败)",
                    7: "无法连接到服务器",
                    28: "操作超时",
                    35: "SSL/TLS 连接错误",
                    56: "网络数据接收失败",
                    60: "SSL 证书验证失败",
                    127: "命令未找到",
                    128: "无效的退出参数",
                }
                hint = code_hints.get(rc, "未知错误")
                result_parts.append(f"Exit Code: {rc} ({hint})")
                # 当没有任何输出时，补充提示帮助 LLM 理解错误
                if not stdout and not stderr:
                    result_parts.append(f"注意: 命令 '{command[:80]}' 执行失败且无输出。建议换用其他工具或方式完成任务。")
            
            return '\n'.join(result_parts)
        
        except subprocess.TimeoutExpired:
            return f"Command timed out after {timeout}s"
    
    def _tool_weather(self, args: dict) -> str:
        """查询天气 (基于 OpenClaw weather skill)"""
        import urllib.request
        import urllib.parse
        
        location = args.get('location', args.get('city', ''))
        if not location:
            raise ValueError("Location/city is required")
        
        # 使用 wttr.in API (无需 API Key)
        encoded_location = urllib.parse.quote(location)
        
        try:
            # 获取详细天气信息
            url = f"https://wttr.in/{encoded_location}?format=j1"
            req = urllib.request.Request(url, headers={'User-Agent': 'curl/7.68.0'})
            
            with urllib.request.urlopen(req, timeout=10) as response:
                data = json.loads(response.read().decode('utf-8'))
            
            current = data.get('current_condition', [{}])[0]
            area = data.get('nearest_area', [{}])[0]
            
            # 格式化输出
            city_name = area.get('areaName', [{}])[0].get('value', location)
            country = area.get('country', [{}])[0].get('value', '')
            
            result = f"""天气查询结果 - {city_name}, {country}

当前温度: {current.get('temp_C', 'N/A')}°C (体感: {current.get('FeelsLikeC', 'N/A')}°C)
天气状况: {current.get('weatherDesc', [{}])[0].get('value', 'N/A')}
湿度: {current.get('humidity', 'N/A')}%
风速: {current.get('windspeedKmph', 'N/A')} km/h ({current.get('winddir16Point', '')})
能见度: {current.get('visibility', 'N/A')} km
紫外线指数: {current.get('uvIndex', 'N/A')}
"""
            return result
            
        except Exception as e:
            # 降级方案：使用简单格式
            try:
                simple_url = f"https://wttr.in/{encoded_location}?format=%l:+%c+%t+(%f)+%h+%w"
                req = urllib.request.Request(simple_url, headers={'User-Agent': 'curl/7.68.0'})
                with urllib.request.urlopen(req, timeout=10) as response:
                    return response.read().decode('utf-8')
            except:
                return f"无法查询 {location} 的天气: {str(e)}"
    
    def _tool_web_search(self, args: dict) -> str:
        """网页搜索 (多源: Bing CN → DuckDuckGo，自动切换)"""
        import urllib.request
        import urllib.parse
        import re
        
        query = args.get('query', args.get('q', ''))
        if not query:
            raise ValueError("Search query is required")
        
        encoded_query = urllib.parse.quote(query)
        ua = 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/131.0.0.0 Safari/537.36'
        
        # ---- 搜索源1: Bing CN (国内直连稳定) ----
        try:
            bing_url = f"https://cn.bing.com/search?q={encoded_query}&ensearch=0"
            req = urllib.request.Request(bing_url, headers={
                'User-Agent': ua,
                'Accept-Language': 'zh-CN,zh;q=0.9',
            })
            with urllib.request.urlopen(req, timeout=12) as response:
                html = response.read().decode('utf-8', errors='ignore')
            
            results = []
            # Bing 搜索结果: <li class="b_algo"><h2><a href="URL">TITLE</a></h2>
            blocks = re.findall(r'<li class="b_algo">([\s\S]*?)</li>', html)
            for block in blocks[:8]:
                link_match = re.search(r'<a[^>]+href="(https?://[^"]+)"[^>]*>([\s\S]*?)</a>', block)
                if not link_match:
                    continue
                link = link_match.group(1)
                title = re.sub(r'<[^>]+>', '', link_match.group(2)).strip()
                if not title:
                    continue
                # 提取摘要
                snippet = ''
                snippet_match = re.search(r'<p[^>]*>([\s\S]*?)</p>', block)
                if snippet_match:
                    snippet = re.sub(r'<[^>]+>', '', snippet_match.group(1)).strip()[:200]
                results.append(f"{len(results)+1}. {title}\n   {link}" + (f"\n   {snippet}" if snippet else ''))
            
            if results:
                return f"搜索 '{query}' 的结果 (Bing):\n\n" + "\n\n".join(results[:6])
        except Exception as e:
            print(f"[webSearch] Bing failed: {e}", file=sys.stderr)
        
        # ---- 搜索源2: DuckDuckGo (备用) ----
        try:
            ddg_url = f"https://html.duckduckgo.com/html/?q={encoded_query}"
            req = urllib.request.Request(ddg_url, headers={'User-Agent': ua})
            with urllib.request.urlopen(req, timeout=12) as response:
                html = response.read().decode('utf-8', errors='ignore')
            
            results = []
            pattern = r'<a[^>]+class="result__a"[^>]*href="([^"]*)"[^>]*>([\s\S]*?)</a>'
            matches = re.findall(pattern, html)
            
            for link, title_raw in matches[:6]:
                title = re.sub(r'<[^>]+>', '', title_raw).strip()
                if 'uddg=' in link:
                    link = urllib.parse.unquote(link.split('uddg=')[-1].split('&')[0])
                if title:
                    results.append(f"{len(results)+1}. {title}\n   {link}")
            
            if results:
                return f"搜索 '{query}' 的结果 (DuckDuckGo):\n\n" + "\n\n".join(results)
        except Exception as e:
            print(f"[webSearch] DuckDuckGo failed: {e}", file=sys.stderr)
        
        # ---- 搜索源3: 使用浏览器搜索 (最终兜底) ----
        if _browser_manager.is_available():
            try:
                bing_url = f"https://cn.bing.com/search?q={encoded_query}&ensearch=0"
                result_json = _browser_manager.navigate(bing_url, wait_until='networkidle')
                parsed = json.loads(result_json)
                if parsed.get('status') == 'ok' and parsed.get('text'):
                    text = parsed['text'][:3000]
                    return f"搜索 '{query}' 的结果 (浏览器 Bing):\n\n{text}"
            except Exception as e:
                print(f"[webSearch] Browser fallback failed: {e}", file=sys.stderr)
        
        return f"搜索 '{query}' 失败: 所有搜索源均不可用（Bing/DuckDuckGo/Browser）"
    
    def _tool_web_fetch(self, args: dict) -> str:
        """获取网页内容 (简化版，提取主要文本)"""
        import urllib.request
        import urllib.parse
        import re
        from html.parser import HTMLParser
        
        url = args.get('url', '')
        if not url:
            raise ValueError("URL is required")
        
        # 确保 URL 有协议
        if not url.startswith(('http://', 'https://')):
            url = 'https://' + url
        
        try:
            req = urllib.request.Request(url, headers={
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36',
                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8',
            })
            
            with urllib.request.urlopen(req, timeout=15) as response:
                # 检查内容类型
                content_type = response.headers.get('Content-Type', '')
                if 'text/html' not in content_type and 'text/plain' not in content_type:
                    return f"无法读取此类型的内容: {content_type}"
                
                html = response.read().decode('utf-8', errors='ignore')
            
            # 简单的 HTML 文本提取
            # 移除 script 和 style 标签
            html = re.sub(r'<script[^>]*>[\s\S]*?</script>', '', html, flags=re.IGNORECASE)
            html = re.sub(r'<style[^>]*>[\s\S]*?</style>', '', html, flags=re.IGNORECASE)
            html = re.sub(r'<head[^>]*>[\s\S]*?</head>', '', html, flags=re.IGNORECASE)
            
            # 提取 title
            title_match = re.search(r'<title[^>]*>([^<]*)</title>', html, re.IGNORECASE)
            title = title_match.group(1).strip() if title_match else ''
            
            # 移除所有 HTML 标签
            text = re.sub(r'<[^>]+>', ' ', html)
            # 清理多余空白
            text = re.sub(r'\s+', ' ', text).strip()
            # 限制长度
            text = text[:4000]
            
            result = f"URL: {url}\n"
            if title:
                result += f"标题: {title}\n"
            result += f"\n内容摘要:\n{text}"
            
            return result
            
        except urllib.error.HTTPError as e:
            return f"HTTP 错误 {e.code}: {e.reason}"
        except urllib.error.URLError as e:
            return f"无法访问 URL: {e.reason}"
        except Exception as e:
            return f"获取网页失败: {str(e)}"
    
    def _tool_save_memory(self, args: dict) -> str:
        """保存记忆到文件"""
        key = args.get('key', '')
        content = args.get('content', '')
        memory_type = args.get('type', 'general')
        
        if not key or not content:
            raise ValueError("key 和 content 参数必填")
        
        # 记忆存储在 memory 目录下
        memory_dir = self.clawd_path / 'memory'
        memory_dir.mkdir(parents=True, exist_ok=True)
        
        # 按日期组织记忆文件
        today = datetime.now().strftime('%Y-%m-%d')
        memory_file = memory_dir / f'{today}.md'
        
        # 格式化记忆条目
        timestamp = datetime.now().strftime('%H:%M:%S')
        entry = f"\n## [{timestamp}] {key}\n- **类型**: {memory_type}\n- **内容**: {content}\n"
        
        # 追加到记忆文件
        with open(memory_file, 'a', encoding='utf-8') as f:
            f.write(entry)
        
        return f"记忆已保存: {key} (类型: {memory_type})"
    
    def _tool_search_memory(self, args: dict) -> str:
        """检索历史记忆"""
        query = args.get('query', '')
        
        if not query:
            raise ValueError("query 参数必填")
        
        memory_dir = self.clawd_path / 'memory'
        if not memory_dir.exists():
            return "记忆库为空，暂无历史记忆。"
        
        results = []
        query_lower = query.lower()
        
        # 遍历所有记忆文件
        for memory_file in sorted(memory_dir.glob('*.md'), reverse=True)[:7]:  # 最近7天
            try:
                content = memory_file.read_text(encoding='utf-8')
                
                # 按条目分割
                entries = content.split('\n## ')
                for entry in entries:
                    if query_lower in entry.lower():
                        # 提取日期和内容
                        date = memory_file.stem
                        results.append(f"[{date}] {entry.strip()[:200]}")
                        
                        if len(results) >= 5:  # 最多返回5条
                            break
            except Exception:
                continue
            
            if len(results) >= 5:
                break
        
        if results:
            return f"找到 {len(results)} 条相关记忆:\n\n" + "\n\n---\n\n".join(results)
        else:
            return f"未找到与 '{query}' 相关的记忆。"
    
    def _tool_nexus_bind_skill(self, args: dict) -> str:
        """为 Nexus 绑定新技能"""
        nexus_id = args.get('nexusId', '')
        skill_id = args.get('skillId', '')
        if not nexus_id or not skill_id:
            raise ValueError('Missing nexusId or skillId')

        nexus_md = self.clawd_path / 'nexuses' / nexus_id / 'NEXUS.md'
        if not nexus_md.exists():
            raise ValueError(f"Nexus '{nexus_id}' not found")

        # 验证技能存在 (skills/ 目录中有对应目录)
        skill_dir = self.clawd_path / 'skills' / skill_id
        if not skill_dir.exists():
            raise ValueError(f"Skill '{skill_id}' not found in skills/")

        frontmatter = parse_nexus_frontmatter(nexus_md)
        deps = list(frontmatter.get('skill_dependencies', []))

        if skill_id in deps:
            return f"Skill '{skill_id}' already bound to Nexus '{nexus_id}'"

        deps.append(skill_id)
        update_nexus_frontmatter(nexus_md, {'skill_dependencies': deps})
        return f"Skill '{skill_id}' bound to Nexus '{nexus_id}'. Dependencies: {deps}"

    def _tool_nexus_unbind_skill(self, args: dict) -> str:
        """从 Nexus 解绑技能"""
        nexus_id = args.get('nexusId', '')
        skill_id = args.get('skillId', '')
        if not nexus_id or not skill_id:
            raise ValueError('Missing nexusId or skillId')

        nexus_md = self.clawd_path / 'nexuses' / nexus_id / 'NEXUS.md'
        if not nexus_md.exists():
            raise ValueError(f"Nexus '{nexus_id}' not found")

        frontmatter = parse_nexus_frontmatter(nexus_md)
        deps = list(frontmatter.get('skill_dependencies', []))

        if skill_id not in deps:
            return f"Skill '{skill_id}' not bound to Nexus '{nexus_id}'"

        if len(deps) <= 1:
            return f"Cannot remove last skill from Nexus '{nexus_id}'. At least 1 skill required."

        deps.remove(skill_id)
        update_nexus_frontmatter(nexus_md, {'skill_dependencies': deps})
        return f"Skill '{skill_id}' unbound from Nexus '{nexus_id}'. Remaining: {deps}"

    def _tool_generate_skill(self, args: dict) -> str:
        """动态生成 Python SKILL 并保存
        
        当遇到无法完成的任务时，Agent 可以调用此工具生成新的 Python 技能来解决问题。
        生成的技能会保存到 skills/ 目录（或 nexuses/{nexusId}/ 目录）并自动热加载。
        
        参数:
        - name: 技能名称 (kebab-case, 如 "pdf-merger")
        - description: 技能描述
        - pythonCode: Python 实现代码 (必须包含 main() 函数)
        - nexusId: 可选，如果指定则保存到对应 Nexus 目录
        - triggers: 可选，触发关键词列表
        """
        name = args.get('name', '')
        description = args.get('description', '')
        python_code = args.get('pythonCode', '')
        nexus_id = args.get('nexusId', '')
        triggers = args.get('triggers', [])
        
        if not name or not description or not python_code:
            raise ValueError("Missing required parameters: name, description, pythonCode")
        
        # 规范化技能名称 (kebab-case)
        safe_name = re.sub(r'[^\w-]', '-', name.lower()).strip('-')
        safe_name = re.sub(r'-+', '-', safe_name)
        
        if not safe_name:
            raise ValueError("Invalid skill name")
        
        # 验证 Python 代码包含 main() 函数
        if 'def main(' not in python_code and 'async def main(' not in python_code:
            raise ValueError("Python code must contain a main() function")
        
        # 确定保存路径
        if nexus_id:
            # 保存到 Nexus 专属目录
            skill_dir = self.clawd_path / 'nexuses' / nexus_id / 'skills' / safe_name
        else:
            # 保存到全局 skills 目录
            skill_dir = self.clawd_path / 'skills' / safe_name
        
        skill_dir.mkdir(parents=True, exist_ok=True)
        
        # 生成 SKILL.md
        trigger_list = '\n'.join(f'- {t}' for t in triggers) if triggers else f'- {safe_name}'
        skill_md_content = f'''---
name: {safe_name}
description: {description}
version: "1.0.0"
author: auto-generated
triggers:
{trigger_list}
---

# {name}

{description}

## 使用方法

此技能由 DD-OS Agent 自动生成，用于解决特定任务。

### 执行

```bash
python {safe_name}.py
```

### 参数

请参考 Python 代码中的 `main()` 函数签名。

## 实现

参见 `{safe_name}.py`
'''
        
        # 写入文件
        skill_md_path = skill_dir / 'SKILL.md'
        skill_md_path.write_text(skill_md_content, encoding='utf-8')
        
        python_file_path = skill_dir / f'{safe_name}.py'
        python_file_path.write_text(python_code, encoding='utf-8')
        
        # 热加载: 重新注册工具
        try:
            tool_registry.refresh_skills()
            loaded_msg = "并已热加载到工具列表"
        except Exception as e:
            loaded_msg = f"但热加载失败: {e}"
        
        return json.dumps({
            'action': 'skill_created',
            'message': f'技能 "{safe_name}" 已成功创建{loaded_msg}',
            'skillName': safe_name,
            'skillDir': str(skill_dir),
            'files': [str(skill_md_path), str(python_file_path)],
            'nexusId': nexus_id or None,
        }, ensure_ascii=False)

    # ============================================
    # 🌐 浏览器自动化工具 (Playwright)
    # ============================================

    def _tool_browser_navigate(self, args: dict) -> str:
        """浏览器导航到指定 URL"""
        if not _browser_manager.is_available():
            raise RuntimeError("Playwright 未安装。请运行: pip install playwright && playwright install chromium")
        url = args.get('url', '')
        if not url:
            raise ValueError("url is required")
        if not url.startswith(('http://', 'https://')):
            url = 'https://' + url
        wait_until = args.get('waitUntil', 'domcontentloaded')
        return _browser_manager.navigate(url, wait_until=wait_until)

    def _tool_browser_click(self, args: dict) -> str:
        """点击页面元素"""
        if not _browser_manager.is_available():
            raise RuntimeError("Playwright 未安装")
        selector = args.get('selector', '')
        if not selector:
            raise ValueError("selector is required")
        return _browser_manager.click(selector)

    def _tool_browser_fill(self, args: dict) -> str:
        """填写表单字段"""
        if not _browser_manager.is_available():
            raise RuntimeError("Playwright 未安装")
        selector = args.get('selector', '')
        value = args.get('value', '')
        if not selector:
            raise ValueError("selector is required")
        return _browser_manager.fill(selector, value)

    def _tool_browser_extract(self, args: dict) -> str:
        """提取页面文本内容"""
        if not _browser_manager.is_available():
            raise RuntimeError("Playwright 未安装")
        selector = args.get('selector', 'body')
        return _browser_manager.extract(selector)

    def _tool_browser_screenshot(self, args: dict) -> str:
        """页面截图"""
        if not _browser_manager.is_available():
            raise RuntimeError("Playwright 未安装")
        selector = args.get('selector')
        full_page = args.get('fullPage', False)
        return _browser_manager.screenshot(selector=selector, full_page=full_page)

    def _tool_browser_evaluate(self, args: dict) -> str:
        """执行 JavaScript"""
        if not _browser_manager.is_available():
            raise RuntimeError("Playwright 未安装")
        expression = args.get('expression', '')
        if not expression:
            raise ValueError("expression is required")
        return _browser_manager.evaluate(expression)

    # ============================================
    # 原有处理器 (保持兼容)
    # ============================================
    
    def handle_index(self):
        # 优先托管前端 dist/index.html (便携式分发模式)
        index_file = APP_DIR / 'dist' / 'index.html'
        if index_file.exists():
            self.serve_static_file('/')
            return

        # dist 不存在时显示 API 文档页
        html = f"""<!DOCTYPE html>
<html>
<head><title>DD-OS Native Server</title></head>
<body style="font-family: monospace; background: #0f172a; color: #e2e8f0; padding: 30px;">
<h1>DD-OS Native Server v{VERSION}</h1>
<p style="color: #94a3b8;">独立运行的本地 AI 操作系统后端</p>
<p>Clawd Path: <code style="color: #22d3ee;">{self.clawd_path}</code></p>

<h2>📡 API Endpoints</h2>
<div style="background: #1e293b; padding: 15px; border-radius: 8px;">
<h3 style="color: #f59e0b;">数据读取</h3>
<ul>
<li><a href="/status" style="color: #60a5fa;">/status</a> - 服务状态</li>
<li><a href="/files" style="color: #60a5fa;">/files</a> - 文件列表</li>
<li><a href="/file/SOUL.md" style="color: #60a5fa;">/file/SOUL.md</a> - 读取 SOUL</li>
<li><a href="/skills" style="color: #60a5fa;">/skills</a> - 技能列表</li>
<li><a href="/all" style="color: #60a5fa;">/all</a> - 所有数据</li>
</ul>

<h3 style="color: #10b981;">🛠️ 工具执行 (POST)</h3>
<ul>
<li><code>/api/tools/execute</code> - 执行工具</li>
<li>支持: readFile, writeFile, listDir, runCmd, appendFile</li>
</ul>
</div>

<h2>🧪 测试</h2>
<pre style="background: #1e293b; padding: 15px; border-radius: 8px; overflow-x: auto;">
curl -X POST http://localhost:3001/api/tools/execute \\
  -H "Content-Type: application/json" \\
  -d '{{"name": "listDir", "args": {{"path": "."}}}}'
</pre>
</body>
</html>"""
        
        self.send_response(200)
        self.send_header('Content-Type', 'text/html; charset=utf-8')
        self.send_cors_headers()
        self.end_headers()
        self.wfile.write(html.encode('utf-8'))
    
    def handle_status(self):
        files = list_files(self.clawd_path)
        skills_dir = self.clawd_path / 'skills'
        skill_count = len(list(skills_dir.iterdir())) if skills_dir.exists() else 0
        
        self.send_json({
            'status': 'ok',
            'version': VERSION,
            'mode': 'native',
            'clawdPath': str(self.clawd_path),
            'fileCount': len(files),
            'skillCount': skill_count,
            'tools': [t['name'] for t in self.registry.list_all()],
            'toolCount': len(self.registry.list_all()),
            'timestamp': datetime.now().isoformat()
        })
    
    def handle_files(self):
        files = list_files(self.clawd_path)
        self.send_json(files)
    
    def handle_file(self, filename):
        filepath = self.clawd_path / filename
        if not filepath.exists():
            self.send_error_json(f'File not found: {filename}', 404)
            return
        
        if not filepath.is_file():
            self.send_error_json(f'Not a file: {filename}', 400)
            return
        
        try:
            filepath.resolve().relative_to(self.clawd_path.resolve())
        except ValueError:
            self.send_error_json('Access denied', 403)
            return
        
        try:
            content = filepath.read_text(encoding='utf-8')
            self.send_text(content)
        except Exception as e:
            self.send_error_json(f'Read error: {str(e)}', 500)
    
    def handle_skills(self):
        """GET /skills - 统一从 SKILL.md frontmatter 扫描所有技能，支持用户目录 + 项目目录"""
        skills = []
        seen = set()
        seen_ids = set()  # 防止重复技能 (用户目录优先)
        
        # 获取技能目录列表：用户目录优先，项目目录作为后备
        skills_dirs = []
        user_skills_dir = self.clawd_path / 'skills'
        if user_skills_dir.exists() and user_skills_dir.is_dir():
            skills_dirs.append(('user', user_skills_dir))
        
        project_path = self.project_path or APP_DIR
        project_skills_dir = project_path / 'skills'
        if project_skills_dir.exists() and project_skills_dir.is_dir() and project_skills_dir != user_skills_dir:
            skills_dirs.append(('bundled', project_skills_dir))
        
        if not skills_dirs:
            self.send_json([])
            return

        for source, skills_dir in skills_dirs:
            # ── 统一扫描 SKILL.md ──
            for skill_md in skills_dir.rglob('SKILL.md'):
                skill_dir = skill_md.parent
                dir_key = str(skill_dir.resolve())
                skill_id = skill_dir.name
                
                if dir_key in seen or skill_id in seen_ids:
                    continue
                seen.add(dir_key)
                seen_ids.add(skill_id)

                frontmatter = parse_skill_frontmatter(skill_md)

                skill_data = {
                    'id': skill_id,
                    'name': frontmatter.get('name', skill_dir.name),
                    'description': frontmatter.get('description', ''),
                    'location': source,  # 'user' 或 'bundled'
                    'path': str(skill_dir),
                    'status': 'active',
                    'enabled': True,
                    'keywords': frontmatter.get('tags', frontmatter.get('keywords', [])),
                }

                # OpenClaw 生态字段
                if frontmatter.get('emoji'):
                    skill_data['emoji'] = frontmatter['emoji']
                if frontmatter.get('author'):
                    skill_data['author'] = frontmatter['author']
                if frontmatter.get('primaryEnv'):
                    skill_data['primaryEnv'] = frontmatter['primaryEnv']
                if frontmatter.get('requires'):
                    skill_data['requires'] = frontmatter['requires']
                if frontmatter.get('install'):
                    skill_data['install'] = frontmatter['install']
                if frontmatter.get('tags'):
                    skill_data['tags'] = frontmatter['tags']
                if frontmatter.get('version'):
                    skill_data['version'] = frontmatter['version']
                if frontmatter.get('dangerLevel'):
                    skill_data['dangerLevel'] = frontmatter['dangerLevel']

                # 无 frontmatter description 时提取正文首段
                if not skill_data['description']:
                    try:
                        content = skill_md.read_text(encoding='utf-8')
                        for line in content.split('\n'):
                            line = line.strip()
                            if line and not line.startswith('#') and not line.startswith('---'):
                                skill_data['description'] = line[:200]
                                break
                    except Exception:
                        pass

                # ── 从 frontmatter 提取工具信息 (替代 _enrich_skill_from_manifest) ──
                executable = frontmatter.get('executable', '')
                if executable and (skill_dir / executable).exists():
                    # 可执行技能
                    skill_data['executable'] = True
                    tools_list = frontmatter.get('tools', [])
                    if tools_list:
                        # 多工具: 从 tools 数组提取 toolNames
                        tool_names = [t.get('toolName') for t in tools_list if t.get('toolName')]
                        skill_data['toolNames'] = tool_names
                        skill_data['toolName'] = tool_names[0] if tool_names else skill_name_to_tool_name(skill_data['name'])
                        # 合并所有工具的 keywords
                        all_keywords = list(skill_data.get('keywords', []))
                        all_inputs = {}
                        for t in tools_list:
                            all_keywords.extend(t.get('keywords', []))
                            all_inputs.update(t.get('inputs', {}))
                        skill_data['keywords'] = list(set(all_keywords))
                        if all_inputs:
                            skill_data['inputs'] = all_inputs
                    else:
                        # 单工具
                        tool_name = skill_name_to_tool_name(skill_data['name'])
                        skill_data['toolName'] = tool_name
                        skill_data['toolNames'] = [tool_name]
                        if frontmatter.get('inputs'):
                            skill_data['inputs'] = frontmatter['inputs']
                else:
                    # 指令型技能
                    skill_data['toolType'] = 'instruction'
                    tool_name = skill_name_to_tool_name(skill_data['name'])
                    skill_data['toolName'] = tool_name
                    skill_data['toolNames'] = [tool_name]
                    if frontmatter.get('inputs'):
                        skill_data['inputs'] = frontmatter['inputs']

                skills.append(skill_data)

            # ── Deprecated fallback: manifest.json (兼容无 SKILL.md 的第三方技能) ──
            for manifest_path in skills_dir.rglob('manifest.json'):
                skill_dir = manifest_path.parent
                dir_key = str(skill_dir.resolve())
                skill_id = skill_dir.name
                
                if dir_key in seen or skill_id in seen_ids:
                    continue
                seen.add(dir_key)
                seen_ids.add(skill_id)

                try:
                    manifest = json.loads(manifest_path.read_text(encoding='utf-8'))
                except Exception:
                    continue

                skill_data = {
                    'id': skill_id,
                    'name': manifest.get('name', skill_dir.name),
                    'description': manifest.get('description', ''),
                    'location': source,
                    'path': str(skill_dir),
                    'status': 'active',
                    'enabled': True,
                    'keywords': manifest.get('keywords', []),
                }

                # 从 manifest 提取工具信息
                tools_list = manifest.get('tools', [])
                if not tools_list:
                    tools_list = [manifest]
                tool_names = [t.get('toolName') for t in tools_list if t.get('toolName')]
                if tool_names:
                    skill_data['toolNames'] = tool_names
                    skill_data['toolName'] = tool_names[0]
                    skill_data['executable'] = True
                skill_data['dangerLevel'] = manifest.get('dangerLevel', 'safe')
                skill_data['version'] = manifest.get('version', '1.0.0')

                skills.append(skill_data)

        self.send_json(skills)

    # ============================================
    # 🌌 Nexus 管理
    # ============================================

    def handle_nexuses(self):
        """GET /nexuses - 扫描 nexuses/ 目录，返回所有 Nexus 列表"""
        nexuses = []
        nexuses_dir = self.clawd_path / 'nexuses'

        if not nexuses_dir.exists():
            nexuses_dir.mkdir(parents=True, exist_ok=True)
            self.send_json([])
            return

        seen = set()

        for nexus_md in nexuses_dir.rglob('NEXUS.md'):
            nexus_dir = nexus_md.parent
            dir_key = str(nexus_dir.resolve())
            if dir_key in seen:
                continue
            seen.add(dir_key)

            frontmatter = parse_nexus_frontmatter(nexus_md)
            if not frontmatter or not frontmatter.get('name'):
                continue

            sop_content = extract_nexus_body(nexus_md)
            exp_dir = nexus_dir / 'experience'
            xp = count_experience_entries(exp_dir) if exp_dir.exists() else 0

            visual_dna = frontmatter.get('visual_dna', {})

            nexus_data = {
                'id': frontmatter.get('name', nexus_dir.name),
                'name': frontmatter.get('name', nexus_dir.name),
                'description': frontmatter.get('description', ''),
                'archetype': frontmatter.get('archetype', 'REACTOR'),
                'version': frontmatter.get('version', '1.0.0'),
                'skillDependencies': frontmatter.get('skill_dependencies', []),
                'tags': frontmatter.get('tags', []),
                'triggers': frontmatter.get('triggers', []),
                'visualDNA': visual_dna,
                'sopContent': sop_content,
                'xp': xp,
                'location': 'local',
                'path': str(nexus_dir),
                'status': 'active',
                # 目标函数驱动字段 (Objective-Driven Execution)
                'objective': frontmatter.get('objective', ''),
                'metrics': frontmatter.get('metrics', []),
                'strategy': frontmatter.get('strategy', ''),
            }
            nexuses.append(nexus_data)

        self.send_json(nexuses)

    def handle_nexuses_health(self):
        """GET /nexuses/health - 检查 nexuses 目录的配置健康状况"""
        nexuses_dir = self.clawd_path / 'nexuses'
        issues = []
        suggestions = []
        stats = {
            'valid_nexuses': 0,
            'orphan_files': 0,
            'missing_nexus_md': 0,
            'invalid_frontmatter': 0,
        }

        if not nexuses_dir.exists():
            self.send_json({
                'healthy': True,
                'issues': [],
                'suggestions': ['nexuses 目录为空，可以开始创建 Nexus'],
                'stats': stats
            })
            return

        # 收集所有有效的 Nexus 目录
        valid_dirs = set()
        for nexus_md in nexuses_dir.rglob('NEXUS.md'):
            nexus_dir = nexus_md.parent
            frontmatter = parse_nexus_frontmatter(nexus_md)
            if frontmatter and frontmatter.get('name'):
                valid_dirs.add(str(nexus_dir.resolve()))
                stats['valid_nexuses'] += 1
            else:
                stats['invalid_frontmatter'] += 1
                issues.append({
                    'type': 'invalid_frontmatter',
                    'path': str(nexus_md),
                    'message': f"NEXUS.md 缺少必要的 'name' 字段",
                })

        # 检查孤立文件（有 .json 但没有 NEXUS.md）
        for item in nexuses_dir.iterdir():
            if item.is_file() and item.suffix == '.json':
                # 检查是否有对应的 NEXUS.md 目录
                stem = item.stem.replace('.json', '')
                potential_dir = nexuses_dir / stem
                if not (potential_dir / 'NEXUS.md').exists():
                    stats['orphan_files'] += 1
                    issues.append({
                        'type': 'orphan_json',
                        'path': str(item),
                        'message': f"发现孤立的 JSON 文件，没有对应的 NEXUS.md",
                        'suggestion': f"创建 {stem}/NEXUS.md 或删除此文件",
                    })
                    suggestions.append(
                        f"文件 '{item.name}' 可能是 AI 生成的配置，需要转换为 NEXUS.md 格式才能被系统识别"
                    )

            # 检查目录但没有 NEXUS.md
            if item.is_dir() and not (item / 'NEXUS.md').exists():
                # 检查目录内是否有其他文件
                files = list(item.iterdir())
                if files:
                    stats['missing_nexus_md'] += 1
                    issues.append({
                        'type': 'missing_nexus_md',
                        'path': str(item),
                        'message': f"目录 '{item.name}' 缺少 NEXUS.md 文件",
                        'files': [f.name for f in files[:5]],
                    })

        healthy = len(issues) == 0
        self.send_json({
            'healthy': healthy,
            'issues': issues,
            'suggestions': suggestions,
            'stats': stats,
            'tip': '运行 /nexuses 查看所有有效的 Nexus' if healthy else '请修复上述问题后重新检查',
        })

    def handle_nexus_detail(self, nexus_name: str):
        """GET /nexuses/{name} - 获取单个 Nexus 完整信息"""
        nexus_dir = self._resolve_nexus_dir(nexus_name)
        if not nexus_dir:
            self.send_error_json(f"Nexus '{nexus_name}' not found", 404)
            return
        nexus_md = nexus_dir / 'NEXUS.md'
        if not nexus_md.exists():
            self.send_error_json(f"Nexus '{nexus_name}' not found", 404)
            return

        frontmatter = parse_nexus_frontmatter(nexus_md)
        sop_content = extract_nexus_body(nexus_md)
        exp_dir = nexus_dir / 'experience'
        xp = count_experience_entries(exp_dir) if exp_dir.exists() else 0

        # 加载最近经验条目
        recent_experiences = []
        for exp_file in ['successes.md', 'failures.md']:
            exp_path = exp_dir / exp_file
            if not exp_path.exists():
                continue
            try:
                content = exp_path.read_text(encoding='utf-8')
                outcome = 'success' if 'success' in exp_file else 'failure'
                entries = content.split('\n### ')
                for entry in entries[1:]:  # skip header
                    entry = entry.strip()
                    if not entry:
                        continue
                    lines = entry.split('\n')
                    title = lines[0].strip() if lines else ''
                    recent_experiences.append({
                        'title': title,
                        'outcome': outcome,
                        'content': '\n'.join(lines[1:]).strip(),
                    })
            except Exception:
                pass

        # 按时间倒序（标题通常包含日期）
        recent_experiences = recent_experiences[-10:][::-1]

        visual_dna = frontmatter.get('visual_dna', {})

        response = {
            'id': frontmatter.get('name', nexus_name),
            'name': frontmatter.get('name', nexus_name),
            'description': frontmatter.get('description', ''),
            'archetype': frontmatter.get('archetype', 'REACTOR'),
            'version': frontmatter.get('version', '1.0.0'),
            'skillDependencies': frontmatter.get('skill_dependencies', []),
            'tags': frontmatter.get('tags', []),
            'triggers': frontmatter.get('triggers', []),
            'visualDNA': visual_dna,
            'sopContent': sop_content,
            'xp': xp,
            'recentExperiences': recent_experiences,
            'location': 'local',
            'path': str(nexus_dir),
            'status': 'active',
            # 目标函数驱动字段 (Objective-Driven Execution)
            'objective': frontmatter.get('objective', ''),
            'metrics': frontmatter.get('metrics', []),
            'strategy': frontmatter.get('strategy', ''),
        }
        self.send_json(response)

    def _resolve_nexus_dir(self, nexus_name: str, auto_create: bool = False) -> Path | None:
        """根据 nexus id/name 找到实际目录（优先精确匹配目录名，其次匹配 frontmatter name）
        
        auto_create: 如果为 True 且找不到已有目录，则自动创建最小 Nexus 目录结构
                     （支持 Observer 自动创建的 Nexus，它们没有 NEXUS.md 文件）
        """
        nexuses_dir = self.clawd_path / 'nexuses'
        if not nexuses_dir.exists():
            if auto_create:
                nexuses_dir.mkdir(parents=True, exist_ok=True)
            else:
                return None
        # 1) 精确匹配目录名
        direct = nexuses_dir / nexus_name
        if (direct / 'NEXUS.md').exists():
            return direct
        # 2) 扫描所有 NEXUS.md，匹配 frontmatter 中的 name
        for nexus_md in nexuses_dir.rglob('NEXUS.md'):
            fm = parse_nexus_frontmatter(nexus_md)
            if fm.get('name') == nexus_name:
                return nexus_md.parent
        # 3) 自动创建最小目录结构 (Observer 创建的 Nexus)
        if auto_create:
            # 将 nexus id 中不安全的路径字符替换
            safe_name = re.sub(r'[<>:"/\\|?*]', '_', nexus_name)
            target_dir = nexuses_dir / safe_name
            target_dir.mkdir(parents=True, exist_ok=True)
            minimal_md = f"---\nname: {nexus_name}\ndescription: Auto-created Nexus\nversion: 1.0.0\nskill_dependencies: []\n---\n"
            (target_dir / 'NEXUS.md').write_text(minimal_md, encoding='utf-8')
            print(f'[Nexus] Auto-created directory for Observer Nexus: {nexus_name}', file=sys.stderr)
            return target_dir
        return None

    def handle_nexus_update_skills(self, nexus_name: str, data: dict):
        """POST /nexuses/{name}/skills - 更新 Nexus 技能依赖"""
        action = data.get('action', '')  # 'add' or 'remove'
        skill_id = data.get('skillId', '')

        if action not in ('add', 'remove') or not skill_id:
            self.send_error_json('Invalid: need action (add/remove) and skillId', 400)
            return

        nexus_dir = self._resolve_nexus_dir(nexus_name, auto_create=True)
        if not nexus_dir:
            self.send_error_json(f"Nexus '{nexus_name}' not found", 404)
            return

        nexus_md = nexus_dir / 'NEXUS.md'

        frontmatter = parse_nexus_frontmatter(nexus_md)
        deps = list(frontmatter.get('skill_dependencies', []))

        if action == 'add':
            if skill_id not in deps:
                deps.append(skill_id)
        elif action == 'remove':
            if len(deps) <= 1:
                self.send_error_json('Cannot remove last skill dependency', 400)
                return
            if skill_id in deps:
                deps.remove(skill_id)

        update_nexus_frontmatter(nexus_md, {'skill_dependencies': deps})

        self.send_json({
            'status': 'ok',
            'nexusId': nexus_name,
            'skillDependencies': deps,
        })

    def handle_nexus_update_meta(self, nexus_name: str, data: dict):
        """POST /nexuses/{name}/meta - 更新 Nexus 元数据(名称等)"""
        nexus_dir = self._resolve_nexus_dir(nexus_name, auto_create=True)
        if not nexus_dir:
            self.send_error_json(f"Nexus '{nexus_name}' not found", 404)
            return
        nexus_md = nexus_dir / 'NEXUS.md'

        new_name = data.get('name', '').strip()
        if not new_name:
            self.send_error_json('Invalid: name is required', 400)
            return
            
        update_nexus_frontmatter(nexus_md, {'name': new_name})

        self.send_json({
            'status': 'ok',
            'nexusId': nexus_name,
            'name': new_name
        })

    def handle_add_experience(self, nexus_name: str, data: dict):
        """POST /nexuses/{name}/experience - 为 Nexus 添加经验记录 (优化4: 结构化索引)"""
        nexus_dir = self._resolve_nexus_dir(nexus_name, auto_create=True)
        if not nexus_dir:
            self.send_error_json(f"Nexus '{nexus_name}' not found", 404)
            return

        task = data.get('task', '')
        tools_used = data.get('tools_used', [])
        outcome = data.get('outcome', 'success')
        key_insight = data.get('key_insight', '')

        if not task:
            self.send_error_json('Missing required field: task', 400)
            return

        # 确保 experience 目录存在
        exp_dir = nexus_dir / 'experience'
        exp_dir.mkdir(parents=True, exist_ok=True)

        # 构建 Markdown 条目
        from datetime import datetime
        timestamp = datetime.now().strftime('%Y-%m-%d')
        tool_seq = ' → '.join(tools_used) if tools_used else 'N/A'

        entry = f"\n### [{timestamp}] {task[:80]}\n"
        entry += f"- **Tools**: {tool_seq}\n"
        if key_insight:
            entry += f"- **Insight**: {key_insight}\n"
        entry += "---\n"

        # 追加到对应文件
        target_file = exp_dir / ('successes.md' if outcome == 'success' else 'failures.md')
        try:
            with target_file.open('a', encoding='utf-8') as f:
                f.write(entry)
        except Exception as e:
            self.send_error_json(f'Failed to write experience: {str(e)}', 500)
            return

        # 优化4: 同步更新结构化索引 (index.json)
        try:
            self._update_experience_index(exp_dir, {
                'type': 'success' if outcome == 'success' else 'failure',
                'task': task[:200],
                'tools': tools_used,
                'insight': key_insight,
                'timestamp': timestamp,
                'category': self._classify_experience(task, key_insight, tools_used),
            })
        except Exception as e:
            # 索引更新失败不影响主流程
            print(f"[WARN] Failed to update experience index: {e}")

        self.send_json({'status': 'ok', 'outcome': outcome})

    def _classify_experience(self, task: str, insight: str, tools: list) -> str:
        """优化4: 经验分类器 — 将经验归类为可检索的类别"""
        text = f"{task} {insight}".lower()
        if any(kw in text for kw in ['timeout', 'network', '超时', '网络', 'connection']):
            return 'network_error'
        if any(kw in text for kw in ['not found', 'enoent', '找不到', '不存在', 'path', '路径']):
            return 'path_error'
        if any(kw in text for kw in ['permission', '权限', 'access denied', 'forbidden']):
            return 'permission_error'
        if any(kw in text for kw in ['parameter', '参数', 'invalid', 'type error', '格式']):
            return 'param_error'
        if any(kw in text for kw in ['file', '文件', 'write', 'read', '写入', '读取']):
            return 'file_operation'
        if any(kw in text for kw in ['search', '搜索', 'web', '查询']):
            return 'search_operation'
        return 'general'

    def _update_experience_index(self, exp_dir: Path, entry: dict):
        """优化4: 维护结构化经验索引"""
        import json
        index_file = exp_dir / 'index.json'
        index = []

        # 读取现有索引
        if index_file.exists():
            try:
                with index_file.open('r', encoding='utf-8') as f:
                    index = json.load(f)
                    if not isinstance(index, list):
                        index = []
            except (json.JSONDecodeError, Exception):
                index = []

        # 追加新条目
        index.append(entry)

        # 限制索引大小: 保留最近 200 条
        if len(index) > 200:
            index = index[-200:]

        # 写入
        with index_file.open('w', encoding='utf-8') as f:
            json.dump(index, f, ensure_ascii=False, indent=2)

    # ============================================
    # 🧬 SOP Fitness API (SOP 演进系统)
    # ============================================

    def handle_nexus_fitness_get(self, nexus_name: str):
        """GET /nexuses/{name}/fitness - 读取 SOP fitness 数据"""
        nexus_dir = self._resolve_nexus_dir(nexus_name)
        if not nexus_dir:
            self.send_error_json(f"Nexus '{nexus_name}' not found", 404)
            return

        fitness_file = nexus_dir / 'sop-fitness.json'
        if not fitness_file.exists():
            # 返回空表示尚无数据，前端会使用默认值
            self.send_json(None)
            return

        try:
            with fitness_file.open('r', encoding='utf-8') as f:
                data = json.load(f)
            self.send_json(data)
        except (json.JSONDecodeError, Exception) as e:
            self.send_error_json(f'Failed to read fitness data: {str(e)}', 500)

    def handle_nexus_fitness_save(self, nexus_name: str, data: dict):
        """POST /nexuses/{name}/fitness - 保存 SOP fitness 数据"""
        nexus_dir = self._resolve_nexus_dir(nexus_name, auto_create=True)
        if not nexus_dir:
            self.send_error_json(f"Nexus '{nexus_name}' not found", 404)
            return

        fitness_file = nexus_dir / 'sop-fitness.json'
        try:
            with fitness_file.open('w', encoding='utf-8') as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
            self.send_json({'status': 'ok'})
        except Exception as e:
            self.send_error_json(f'Failed to save fitness data: {str(e)}', 500)

    def handle_nexus_sop_content_get(self, nexus_name: str):
        """GET /nexuses/{name}/sop-content - 读取 NEXUS.md 完整内容"""
        nexus_dir = self._resolve_nexus_dir(nexus_name)
        if not nexus_dir:
            self.send_error_json(f"Nexus '{nexus_name}' not found", 404)
            return

        nexus_md = nexus_dir / 'NEXUS.md'
        if not nexus_md.exists():
            self.send_json({'content': None})
            return

        try:
            content = nexus_md.read_text(encoding='utf-8')
            self.send_json({'content': content})
        except Exception as e:
            self.send_error_json(f'Failed to read NEXUS.md: {str(e)}', 500)

    def handle_nexus_sop_content_save(self, nexus_name: str, data: dict):
        """POST /nexuses/{name}/sop-content - 写入 NEXUS.md 完整内容"""
        nexus_dir = self._resolve_nexus_dir(nexus_name, auto_create=True)
        if not nexus_dir:
            self.send_error_json(f"Nexus '{nexus_name}' not found", 404)
            return

        content = data.get('content', '')
        if not content:
            self.send_error_json('Missing required field: content', 400)
            return

        nexus_md = nexus_dir / 'NEXUS.md'
        try:
            nexus_md.write_text(content, encoding='utf-8')
            self.send_json({'status': 'ok'})
        except Exception as e:
            self.send_error_json(f'Failed to write NEXUS.md: {str(e)}', 500)

    def handle_nexus_sop_history_get(self, nexus_name: str):
        """GET /nexuses/{name}/sop-history - 读取指定版本的 SOP 历史"""
        nexus_dir = self._resolve_nexus_dir(nexus_name)
        if not nexus_dir:
            self.send_error_json(f"Nexus '{nexus_name}' not found", 404)
            return

        # 解析 query 参数获取版本号
        parsed = urlparse(self.path)
        query = parse_qs(parsed.query)
        version = query.get('version', [None])[0]

        history_dir = nexus_dir / 'sop-history'
        if not history_dir.exists():
            self.send_json({'content': None, 'versions': []})
            return

        # 列出所有版本
        versions = sorted([
            f.stem  # e.g. 'v1', 'v2'
            for f in history_dir.glob('v*.md')
        ])

        if version:
            history_file = history_dir / f'{version}.md'
            if history_file.exists():
                content = history_file.read_text(encoding='utf-8')
                self.send_json({'content': content, 'versions': versions})
            else:
                self.send_json({'content': None, 'versions': versions})
        else:
            self.send_json({'content': None, 'versions': versions})

    def handle_nexus_sop_history_save(self, nexus_name: str, data: dict):
        """POST /nexuses/{name}/sop-history - 保存一个 SOP 版本到历史"""
        nexus_dir = self._resolve_nexus_dir(nexus_name, auto_create=True)
        if not nexus_dir:
            self.send_error_json(f"Nexus '{nexus_name}' not found", 404)
            return

        version = data.get('version', '')
        content = data.get('content', '')
        if not version or not content:
            self.send_error_json('Missing required fields: version, content', 400)
            return

        history_dir = nexus_dir / 'sop-history'
        history_dir.mkdir(parents=True, exist_ok=True)

        history_file = history_dir / f'v{version}.md'
        try:
            history_file.write_text(content, encoding='utf-8')
            self.send_json({'status': 'ok', 'version': version})
        except Exception as e:
            self.send_error_json(f'Failed to save SOP history: {str(e)}', 500)

    def handle_tools_list(self):
        """GET /tools - 列出所有已注册的工具"""
        self.send_json(self.registry.list_all())

    def handle_tools_reload(self, data=None):
        """POST /tools/reload - 热重载插件工具"""
        self.registry.plugin_tools.clear()
        self.registry.instruction_tools.clear()
        self.registry.scan_plugins()
        tools = self.registry.list_all()
        self.send_json({
            'status': 'ok',
            'message': f'Reloaded. {len(tools)} tools registered.',
            'tools': tools,
        })

    # ============================================
    # 🔌 MCP 服务器管理
    # ============================================

    def handle_mcp_servers_list(self):
        """GET /mcp/servers - 列出 MCP 服务器状态"""
        if not self.registry.mcp_manager:
            self.send_json({
                'status': 'ok',
                'enabled': False,
                'servers': {},
                'message': 'MCP support not initialized'
            })
            return

        status = self.registry.mcp_manager.get_server_status()
        mcp_tools = [t for t in self.registry.list_all() if t.get('type') == 'mcp']

        self.send_json({
            'status': 'ok',
            'enabled': True,
            'servers': status,
            'toolCount': len(mcp_tools),
            'tools': mcp_tools,
        })

    def handle_mcp_reload(self, data=None):
        """POST /mcp/reload - 重新加载 MCP 服务器"""
        if not HAS_MCP:
            self.send_json({
                'status': 'error',
                'message': 'MCP support not available'
            }, 400)
            return

        # 清理现有 MCP 工具
        self.registry.mcp_tools.clear()
        if self.registry.mcp_manager:
            self.registry.mcp_manager.shutdown_all()

        # 重新扫描
        self.registry.scan_mcp_servers()

        mcp_tools = [t for t in self.registry.list_all() if t.get('type') == 'mcp']
        self.send_json({
            'status': 'ok',
            'message': f'MCP reloaded. {len(mcp_tools)} tool(s) registered.',
            'tools': mcp_tools,
        })

    def handle_mcp_reconnect(self, server_name: str):
        """POST /mcp/servers/{name}/reconnect - 重连 MCP 服务器"""
        if not self.registry.mcp_manager:
            self.send_json({
                'status': 'error',
                'message': 'MCP support not initialized'
            }, 400)
            return

        success = self.registry.mcp_manager.reconnect_server(server_name)

        if success:
            # 更新工具注册
            self.registry.mcp_tools.clear()
            for tool_info in self.registry.mcp_manager.get_all_tools():
                tool_name = tool_info['name']
                if tool_name not in self.registry.builtin_tools and tool_name not in self.registry.plugin_tools:
                    self.registry.mcp_tools[tool_name] = {
                        'name': tool_name,
                        'server': tool_info.get('server', ''),
                        'description': tool_info.get('description', ''),
                        'inputs': tool_info.get('inputs', {}),
                        'dangerLevel': 'safe',
                        'version': '1.0.0',
                    }

            self.send_json({
                'status': 'ok',
                'message': f'Server {server_name} reconnected',
                'server': server_name,
            })
        else:
            self.send_json({
                'status': 'error',
                'message': f'Failed to reconnect server: {server_name}'
            }, 500)

    # ============================================
    # 🔍 Registry 在线搜索 (TF-IDF, 无 LLM)
    # ============================================

    def handle_registry_skills_search(self, query: dict):
        """GET /api/registry/skills?q={query} - 搜索可安装的技能"""
        q = query.get('q', [''])[0].strip().lower()
        
        # 读取 registry 文件
        registry_path = self.clawd_path / 'registry' / 'skills.json'
        if not registry_path.exists():
            self.send_json({'status': 'ok', 'results': [], 'message': 'Registry not found'})
            return
        
        try:
            registry = json.loads(registry_path.read_text(encoding='utf-8'))
            skills = registry.get('skills', [])
        except Exception as e:
            self.send_json({'status': 'error', 'message': f'Failed to read registry: {e}'}, 500)
            return
        
        # 如果没有查询词，返回所有
        if not q:
            self.send_json({
                'status': 'ok',
                'results': skills[:20],
                'total': len(skills)
            })
            return
        
        # TF-IDF 风格的关键词匹配
        tokens = self._tokenize(q)
        scored_results = []
        
        for skill in skills:
            score = self._compute_skill_score(skill, tokens)
            if score > 0:
                scored_results.append({**skill, 'score': score})
        
        # 按分数排序
        scored_results.sort(key=lambda x: x['score'], reverse=True)
        
        self.send_json({
            'status': 'ok',
            'results': scored_results[:10],
            'total': len(scored_results),
            'query': q
        })

    def handle_registry_mcp_search(self, query: dict):
        """GET /api/registry/mcp?q={query} - 搜索可安装的 MCP 服务器"""
        q = query.get('q', [''])[0].strip().lower()
        
        # 读取 registry 文件
        registry_path = self.clawd_path / 'registry' / 'mcp-servers.json'
        if not registry_path.exists():
            self.send_json({'status': 'ok', 'results': [], 'message': 'Registry not found'})
            return
        
        try:
            registry = json.loads(registry_path.read_text(encoding='utf-8'))
            servers = registry.get('servers', [])
        except Exception as e:
            self.send_json({'status': 'error', 'message': f'Failed to read registry: {e}'}, 500)
            return
        
        # 如果没有查询词，返回所有
        if not q:
            self.send_json({
                'status': 'ok',
                'results': servers[:20],
                'total': len(servers)
            })
            return
        
        # TF-IDF 风格的关键词匹配
        tokens = self._tokenize(q)
        scored_results = []
        
        for server in servers:
            score = self._compute_mcp_score(server, tokens)
            if score > 0:
                scored_results.append({**server, 'score': score})
        
        # 按分数排序
        scored_results.sort(key=lambda x: x['score'], reverse=True)
        
        self.send_json({
            'status': 'ok',
            'results': scored_results[:10],
            'total': len(scored_results),
            'query': q
        })

    def _tokenize(self, text: str) -> list:
        """分词：按空格和标点拆分"""
        import re
        tokens = re.split(r'[\s,，.。!！?？、;；:：\-—]+', text.lower())
        return [t for t in tokens if t and len(t) >= 1]

    def _compute_skill_score(self, skill: dict, tokens: list) -> float:
        """计算技能的匹配分数 (TF-IDF 简化版)"""
        score = 0.0
        name = skill.get('name', '').lower()
        desc = skill.get('description', '').lower()
        keywords = [k.lower() for k in skill.get('keywords', [])]
        full_text = f"{name} {desc} {' '.join(keywords)}"
        
        for token in tokens:
            # 词频 (TF)
            tf = full_text.count(token)
            # 长词权重更高 (简化 IDF)
            idf = 1.5 if len(token) > 3 else 1.0
            score += tf * idf
            
            # 精确匹配加权
            if token in name:
                score += 10
            if token in keywords:
                score += 5
        
        return min(score, 100)

    def _compute_mcp_score(self, server: dict, tokens: list) -> float:
        """计算 MCP 服务器的匹配分数"""
        score = 0.0
        name = server.get('name', '').lower()
        desc = server.get('description', '').lower()
        keywords = [k.lower() for k in server.get('keywords', [])]
        full_text = f"{name} {desc} {' '.join(keywords)}"
        
        for token in tokens:
            tf = full_text.count(token)
            idf = 1.5 if len(token) > 3 else 1.0
            score += tf * idf
            
            if token in name:
                score += 10
            if token in keywords:
                score += 5
        
        return min(score, 100)

    def handle_mcp_install(self, data: dict):
        """POST /mcp/install - 安装 MCP 服务器配置"""
        server_id = data.get('id', '')
        server_name = data.get('name', server_id)
        command = data.get('command', '')
        args = data.get('args', [])
        env = data.get('env', {})
        
        if not server_name or not command:
            self.send_error_json('Missing required fields: name, command', 400)
            return
        
        # 安全检查
        if '..' in server_name or '/' in server_name or '\\' in server_name:
            self.send_error_json('Invalid server name', 400)
            return
        
        # 读取现有配置
        config_path = self.clawd_path / 'mcp-servers.json'
        try:
            if config_path.exists():
                config = json.loads(config_path.read_text(encoding='utf-8'))
            else:
                config = {'servers': {}}
        except Exception as e:
            self.send_error_json(f'Failed to read config: {e}', 500)
            return
        
        # 检查是否已存在
        if server_name in config.get('servers', {}):
            self.send_error_json(f'Server already exists: {server_name}', 409)
            return
        
        # 添加新服务器配置
        config['servers'][server_name] = {
            'command': command,
            'args': args,
            'env': env,
            'enabled': False  # 默认禁用，需要用户手动启用
        }
        
        # 写回配置文件
        try:
            config_path.write_text(json.dumps(config, indent=2, ensure_ascii=False), encoding='utf-8')
        except Exception as e:
            self.send_error_json(f'Failed to write config: {e}', 500)
            return
        
        self.send_json({
            'status': 'ok',
            'serverName': server_name,
            'message': f'MCP server "{server_name}" added (disabled by default). Enable it in mcp-servers.json to use.',
            'configPath': str(config_path)
        })

    # ============================================
    # 📦 远程技能安装/卸载
    # ============================================

    def handle_skill_install(self, data):
        """POST /skills/install - 从 Git URL 安装技能"""
        source = data.get('source', '')
        name = data.get('name', '')

        if not source:
            self.send_error_json('Missing source parameter', 400)
            return

        if not source.startswith(('http://', 'https://', 'git@')):
            self.send_error_json('Unsupported source format. Use a Git URL (https://... or git@...)', 400)
            return

        try:
            # 从 URL 提取仓库名
            match = re.search(r'/([^/]+?)(?:\.git)?$', source)
            repo_name = name or (match.group(1) if match else 'downloaded-skill')
            # 安全化名称
            repo_name = re.sub(r'[^\w\-.]', '_', repo_name)

            target = self.clawd_path / 'skills' / repo_name
            if target.exists():
                self.send_error_json(f'Skill already exists: {repo_name}. Use /skills/uninstall first.', 409)
                return

            # 确保 skills/ 目录存在
            (self.clawd_path / 'skills').mkdir(parents=True, exist_ok=True)

            # Git clone (shallow, 限制深度)
            process = subprocess.run(
                ['git', 'clone', '--depth', '1', source, str(target)],
                capture_output=True,
                text=True,
                timeout=120,
            )

            if process.returncode != 0:
                # 清理失败的 clone
                if target.exists():
                    shutil.rmtree(target, ignore_errors=True)
                stderr = process.stderr[:500] if process.stderr else 'Unknown error'
                raise RuntimeError(f"git clone failed: {stderr}")

            # 验证: 必须有 SKILL.md (manifest.json 已 deprecated)
            has_skill_md = (target / 'SKILL.md').exists() or any(target.rglob('SKILL.md'))
            has_manifest = (target / 'manifest.json').exists()

            if not has_skill_md:
                if has_manifest:
                    print(f"[SkillInstall] ⚠️ DEPRECATED: {target.name} only has manifest.json, please add SKILL.md")
                else:
                    shutil.rmtree(target, ignore_errors=True)
                    self.send_error_json('Invalid skill: no SKILL.md found', 400)
                    return

            # 重新扫描注册
            self.registry.plugin_tools.clear()
            self.registry.instruction_tools.clear()
            self.registry.scan_plugins()

            self.send_json({
                'status': 'ok',
                'name': repo_name,
                'path': str(target),
                'message': f'Skill installed: {repo_name}',
                'toolCount': len(self.registry.list_all()),
            })

        except subprocess.TimeoutExpired:
            if target.exists():
                shutil.rmtree(target, ignore_errors=True)
            self.send_error_json('Git clone timed out (120s limit)', 500)
        except RuntimeError as e:
            self.send_error_json(str(e), 500)
        except Exception as e:
            self.send_error_json(f'Installation failed: {str(e)}', 500)

    def handle_skill_uninstall(self, data):
        """POST /skills/uninstall - 卸载技能"""
        name = data.get('name', '')

        if not name:
            self.send_error_json('Missing name parameter', 400)
            return

        # 安全检查: 不允许路径遍历
        if '..' in name or '/' in name or '\\' in name:
            self.send_error_json('Invalid skill name', 400)
            return

        target = self.clawd_path / 'skills' / name

        if not target.exists():
            self.send_error_json(f'Skill not found: {name}', 404)
            return

        if not target.is_dir():
            self.send_error_json(f'Not a directory: {name}', 400)
            return

        try:
            shutil.rmtree(target)

            # 重新扫描
            self.registry.plugin_tools.clear()
            self.registry.instruction_tools.clear()
            self.registry.scan_plugins()

            self.send_json({
                'status': 'ok',
                'name': name,
                'message': f'Skill uninstalled: {name}',
                'toolCount': len(self.registry.list_all()),
            })

        except Exception as e:
            self.send_error_json(f'Uninstall failed: {str(e)}', 500)

    def handle_clawhub_install(self, data):
        """POST /clawhub/install - 从 ClawHub 下载并安装技能"""
        slug = data.get('slug', '')
        archive_url = data.get('archive_url', '')
        skill_name = data.get('name', '')

        if not archive_url:
            self.send_error_json('Missing archive_url parameter', 400)
            return

        # 从 slug 提取技能名
        if not skill_name:
            skill_name = slug.split('/')[-1] if '/' in slug else slug
        skill_name = re.sub(r'[^\w\-.]', '_', skill_name)

        if not skill_name:
            self.send_error_json('Cannot determine skill name', 400)
            return

        target = self.clawd_path / 'skills' / skill_name

        if target.exists():
            self.send_error_json(f'Skill already exists: {skill_name}. Use /skills/uninstall first.', 409)
            return

        try:
            import tempfile
            import tarfile
            import io

            # 确保 skills/ 目录存在
            (self.clawd_path / 'skills').mkdir(parents=True, exist_ok=True)

            # 下载归档
            req = _urllib_req.Request(archive_url, headers={'User-Agent': 'DD-OS/3.0'})
            with _urllib_req.urlopen(req, timeout=60) as resp:
                archive_data = resp.read()

            # 解压 tar.gz
            with tarfile.open(fileobj=io.BytesIO(archive_data), mode='r:gz') as tar:
                # 安全检查: 确保没有路径遍历
                for member in tar.getmembers():
                    if member.name.startswith('/') or '..' in member.name:
                        raise RuntimeError(f'Unsafe path in archive: {member.name}')

                # 创建目标目录
                target.mkdir(parents=True, exist_ok=True)

                # 解压时去掉顶层目录 (如果存在)
                members = tar.getmembers()
                top_dirs = set()
                for m in members:
                    parts = m.name.split('/')
                    if parts[0]:
                        top_dirs.add(parts[0])

                strip_prefix = ''
                if len(top_dirs) == 1:
                    strip_prefix = top_dirs.pop() + '/'

                for member in members:
                    if strip_prefix and member.name.startswith(strip_prefix):
                        member.name = member.name[len(strip_prefix):]
                    if not member.name or member.name == '.':
                        continue
                    tar.extract(member, target)

            # 验证
            has_skill_md = (target / 'SKILL.md').exists()
            has_manifest = (target / 'manifest.json').exists()
            if not has_skill_md:
                if has_manifest:
                    print(f"[SkillInstall] ⚠️ DEPRECATED: {target.name} only has manifest.json, please add SKILL.md")
                else:
                    shutil.rmtree(target, ignore_errors=True)
                    self.send_error_json('Invalid skill: no SKILL.md found', 400)
                    return

            # 重新扫描注册
            self.registry.plugin_tools.clear()
            self.registry.instruction_tools.clear()
            self.registry.scan_plugins()

            self.send_json({
                'status': 'ok',
                'name': skill_name,
                'slug': slug,
                'path': str(target),
                'message': f'Skill installed from ClawHub: {skill_name}',
                'toolCount': len(self.registry.list_all()),
            })

        except Exception as e:
            if target.exists():
                shutil.rmtree(target, ignore_errors=True)
            self.send_error_json(f'ClawHub install failed: {str(e)}', 500)

    def handle_clawhub_publish(self, data):
        """POST /clawhub/publish - 打包本地技能供发布到 ClawHub"""
        skill_name = data.get('skill_name', '')

        if not skill_name:
            self.send_error_json('Missing skill_name parameter', 400)
            return

        if '..' in skill_name or '/' in skill_name or '\\' in skill_name:
            self.send_error_json('Invalid skill name', 400)
            return

        # 搜索技能目录
        skill_dir = None
        for skills_parent in [self.clawd_path / 'skills', (self.project_path or APP_DIR) / 'skills']:
            candidate = skills_parent / skill_name
            if candidate.exists() and candidate.is_dir():
                skill_dir = candidate
                break

        if not skill_dir:
            self.send_error_json(f'Skill not found: {skill_name}', 404)
            return

        try:
            import tarfile
            import io
            import base64

            # 读取 frontmatter 元数据
            skill_md = skill_dir / 'SKILL.md'
            frontmatter = parse_skill_frontmatter(skill_md) if skill_md.exists() else {}

            # 打包为 tar.gz (排除 __pycache__, node_modules, .git)
            buffer = io.BytesIO()
            exclude_patterns = {'__pycache__', 'node_modules', '.git', '.env', '.DS_Store'}

            with tarfile.open(fileobj=buffer, mode='w:gz') as tar:
                for item in skill_dir.rglob('*'):
                    # 排除不需要的文件/目录
                    parts = item.relative_to(skill_dir).parts
                    if any(p in exclude_patterns for p in parts):
                        continue
                    if item.is_file():
                        arcname = str(item.relative_to(skill_dir))
                        tar.add(str(item), arcname=arcname)

            archive_b64 = base64.b64encode(buffer.getvalue()).decode('ascii')

            # 收集文件列表
            file_list = []
            for item in skill_dir.rglob('*'):
                parts = item.relative_to(skill_dir).parts
                if any(p in exclude_patterns for p in parts):
                    continue
                if item.is_file():
                    file_list.append(str(item.relative_to(skill_dir)))

            self.send_json({
                'status': 'ok',
                'name': frontmatter.get('name', skill_name),
                'description': frontmatter.get('description', ''),
                'version': frontmatter.get('version', '1.0.0'),
                'tags': frontmatter.get('tags', []),
                'archive_base64': archive_b64,
                'file_list': file_list,
                'archive_size': len(buffer.getvalue()),
            })

        except Exception as e:
            self.send_error_json(f'Publish packaging failed: {str(e)}', 500)

    def handle_skill_raw(self, skill_name: str):
        """GET /skills/<name>/raw - 获取技能原始文件列表和内容"""
        if '..' in skill_name or '/' in skill_name or '\\' in skill_name:
            self.send_error_json('Invalid skill name', 400)
            return

        skill_dir = None
        for skills_parent in [self.clawd_path / 'skills', (self.project_path or APP_DIR) / 'skills']:
            candidate = skills_parent / skill_name
            if candidate.exists() and candidate.is_dir():
                skill_dir = candidate
                break

        if not skill_dir:
            self.send_error_json(f'Skill not found: {skill_name}', 404)
            return

        try:
            exclude_patterns = {'__pycache__', 'node_modules', '.git', '.env'}
            files = []
            for item in skill_dir.rglob('*'):
                parts = item.relative_to(skill_dir).parts
                if any(p in exclude_patterns for p in parts):
                    continue
                if item.is_file():
                    rel_path = str(item.relative_to(skill_dir))
                    try:
                        content = item.read_text(encoding='utf-8')
                    except Exception:
                        content = f'[binary file, {item.stat().st_size} bytes]'
                    files.append({
                        'path': rel_path,
                        'size': item.stat().st_size,
                        'content': content[:10000],  # 限制内容大小
                    })

            self.send_json({
                'name': skill_name,
                'path': str(skill_dir),
                'files': files,
            })

        except Exception as e:
            self.send_error_json(f'Failed to read skill: {str(e)}', 500)

    def handle_trace_save(self, data):
        """POST /api/traces/save - 保存执行追踪 (P2: 执行流记忆)"""
        if not data:
            self.send_error_json('Missing trace data', 400)
            return

        traces_dir = self.clawd_path / 'memory' / 'exec_traces'
        traces_dir.mkdir(parents=True, exist_ok=True)

        # 按月分片存储
        month = datetime.now().strftime('%Y-%m')
        trace_file = traces_dir / f'{month}.jsonl'

        # 敏感数据脱敏
        trace_json = json.dumps(data, ensure_ascii=False)
        import re
        trace_json = re.sub(
            r'(password|token|secret|api_key|apikey|auth)["\s:]*["\']([^"\']{3,})["\']',
            r'\1": "***"',
            trace_json,
            flags=re.IGNORECASE
        )

        try:
            with open(trace_file, 'a', encoding='utf-8') as f:
                f.write(trace_json + '\n')

            self.send_json({
                'status': 'ok',
                'message': f'Trace saved to {month}.jsonl',
            })
        except Exception as e:
            self.send_error_json(f'Failed to save trace: {e}', 500)

    # ============================================
    # 🧬 Gene Pool API
    # ============================================

    def handle_gene_save(self, data):
        """POST /api/genes/save - 保存/更新基因到基因库"""
        if not data:
            self.send_error_json('Missing gene data', 400)
            return

        gene_file = self.clawd_path / 'memory' / 'gene_pool.jsonl'
        gene_file.parent.mkdir(parents=True, exist_ok=True)

        gene_id = data.get('id', '')

        try:
            with ClawdDataHandler._gene_file_lock:
                # 如果基因已存在 (同 ID)，先读取并替换
                existing_lines = []
                replaced = False
                if gene_file.exists():
                    with open(gene_file, 'r', encoding='utf-8', errors='replace') as f:
                        for line in f:
                            line = line.strip()
                            if not line:
                                continue
                            try:
                                existing = json.loads(line)
                                if existing.get('id') == gene_id:
                                    existing_lines.append(json.dumps(data, ensure_ascii=False))
                                    replaced = True
                                else:
                                    existing_lines.append(line)
                            except (json.JSONDecodeError, UnicodeDecodeError):
                                # 跳过损坏的行，不再保留
                                continue

                if replaced:
                    # 覆写整个文件 (替换已有基因)
                    with open(gene_file, 'w', encoding='utf-8') as f:
                        for line in existing_lines:
                            f.write(line + '\n')
                else:
                    # 追写新基因
                    with open(gene_file, 'a', encoding='utf-8') as f:
                        f.write(json.dumps(data, ensure_ascii=False) + '\n')

            self.send_json({
                'status': 'ok',
                'message': f'Gene {"updated" if replaced else "saved"}: {gene_id}',
            })
        except Exception as e:
            self.send_error_json(f'Failed to save gene: {e}', 500)

    def handle_gene_load(self):
        """GET /api/genes/load - 加载全部基因"""
        gene_file = self.clawd_path / 'memory' / 'gene_pool.jsonl'

        genes = []
        if not gene_file.exists():
            self.send_json(genes)
            return

        try:
            with ClawdDataHandler._gene_file_lock:
                with open(gene_file, 'r', encoding='utf-8', errors='replace') as f:
                    for line in f:
                        line = line.strip()
                        if not line:
                            continue
                        try:
                            genes.append(json.loads(line))
                        except (json.JSONDecodeError, UnicodeDecodeError):
                            continue
        except Exception as e:
            self.send_error_json(f'Failed to load genes: {e}', 500)
            return

        self.send_json(genes)

    def handle_capsule_load(self):
        """GET /api/capsules/load - 加载全部胶囊"""
        capsule_file = self.clawd_path / 'memory' / 'capsules.json'

        if not capsule_file.exists():
            self.send_json([])
            return

        try:
            content = capsule_file.read_text(encoding='utf-8')
            capsules = json.loads(content) if content.strip() else []
            self.send_json(capsules)
        except Exception as e:
            self.send_error_json(f'Failed to load capsules: {e}', 500)

    def handle_capsule_save(self, data):
        """POST /api/capsules/save - 批量保存胶囊 (全量覆写)"""
        if not isinstance(data, list):
            self.send_error_json('Expected array of capsules', 400)
            return

        capsule_file = self.clawd_path / 'memory' / 'capsules.json'
        capsule_file.parent.mkdir(parents=True, exist_ok=True)

        try:
            # 只保留最近 100 条
            trimmed = data[-100:] if len(data) > 100 else data
            capsule_file.write_text(json.dumps(trimmed, ensure_ascii=False, indent=2), encoding='utf-8')
            self.send_json({
                'status': 'ok',
                'message': f'Saved {len(trimmed)} capsules',
            })
        except Exception as e:
            self.send_error_json(f'Failed to save capsules: {e}', 500)

    def handle_trace_search(self, query_params):
        """GET /api/traces/search?query=xxx&limit=5 - 检索执行追踪 (P2)"""
        query = query_params.get('query', [''])[0]
        limit = min(int(query_params.get('limit', ['5'])[0]), 20)

        if not query:
            self.send_json([])
            return

        traces_dir = self.clawd_path / 'memory' / 'exec_traces'
        if not traces_dir.exists():
            self.send_json([])
            return

        query_lower = query.lower()
        query_words = [w for w in query_lower.split() if len(w) > 1]
        results = []

        # 从最近的月份文件开始搜索
        for trace_file in sorted(traces_dir.glob('*.jsonl'), reverse=True)[:6]:
            try:
                for line in reversed(trace_file.read_text(encoding='utf-8').strip().split('\n')):
                    if not line.strip():
                        continue
                    try:
                        trace = json.loads(line)
                        task = trace.get('task', '').lower()
                        tags = [t.lower() for t in trace.get('tags', [])]
                        # 关键词匹配: task 描述或 tags
                        matched = any(w in task for w in query_words) or \
                                  any(w in ' '.join(tags) for w in query_words)
                        if matched:
                            results.append(trace)
                            if len(results) >= limit:
                                break
                    except json.JSONDecodeError:
                        continue
            except Exception:
                continue
            if len(results) >= limit:
                break

        self.send_json(results)
    
    def handle_trace_recent(self, query_params):
        """GET /api/traces/recent?days=3&limit=100 - 获取最近N天的执行日志 (供 Observer 分析)"""
        days = min(int(query_params.get('days', ['3'])[0]), 30)
        limit = min(int(query_params.get('limit', ['100'])[0]), 500)
        
        traces_dir = self.clawd_path / 'memory' / 'exec_traces'
        if not traces_dir.exists():
            self.send_json({'traces': [], 'stats': {}})
            return
        
        cutoff_time = datetime.now() - timedelta(days=days)
        cutoff_ts = cutoff_time.timestamp() * 1000  # 毫秒时间戳
        
        traces = []
        tool_freq = {}  # 工具使用频率
        nexus_freq = {}  # Nexus 使用频率
        total_turns = 0
        total_errors = 0
        
        # 从最近的月份文件开始读取
        for trace_file in sorted(traces_dir.glob('*.jsonl'), reverse=True)[:3]:
            try:
                for line in reversed(trace_file.read_text(encoding='utf-8').strip().split('\n')):
                    if not line.strip():
                        continue
                    try:
                        trace = json.loads(line)
                        ts = trace.get('timestamp', 0)
                        if ts < cutoff_ts:
                            continue  # 超出时间范围
                        
                        traces.append(trace)
                        
                        # 统计工具频率
                        for tool in trace.get('tools', []):
                            tool_name = tool.get('name', 'unknown')
                            tool_freq[tool_name] = tool_freq.get(tool_name, 0) + 1
                        
                        # 统计 Nexus 频率
                        nexus_id = trace.get('activeNexusId')
                        if nexus_id:
                            nexus_freq[nexus_id] = nexus_freq.get(nexus_id, 0) + 1
                        
                        # 统计轮次和错误
                        total_turns += trace.get('turnCount', 0)
                        total_errors += trace.get('errorCount', 0)
                        
                        if len(traces) >= limit:
                            break
                    except json.JSONDecodeError:
                        continue
            except Exception:
                continue
            if len(traces) >= limit:
                break
        
        # 按时间倒序排列
        traces.sort(key=lambda t: t.get('timestamp', 0), reverse=True)
        
        self.send_json({
            'traces': traces,
            'stats': {
                'totalExecutions': len(traces),
                'toolFrequency': tool_freq,
                'nexusFrequency': nexus_freq,
                'avgTurnsPerExecution': total_turns / len(traces) if traces else 0,
                'totalErrors': total_errors,
                'timeRangeDays': days,
            }
        })
    
    def handle_memories(self):
        memories = []
        
        memory_md = self.clawd_path / 'MEMORY.md'
        if memory_md.exists():
            try:
                content = memory_md.read_text(encoding='utf-8')
                memories.extend(parse_memory_md(content))
            except:
                pass
        
        memory_dir = self.clawd_path / 'memory'
        if memory_dir.exists() and memory_dir.is_dir():
            for item in memory_dir.iterdir():
                if item.is_file() and item.suffix == '.md':
                    try:
                        content = item.read_text(encoding='utf-8')
                        memories.append({
                            'id': f'file-{item.stem}',
                            'title': item.stem.replace('-', ' ').replace('_', ' ').title(),
                            'content': content[:500],
                            'type': 'long-term',
                            'timestamp': item.stat().st_mtime,
                            'tags': [],
                        })
                    except:
                        pass
        
        self.send_json(memories)
    
    def handle_all(self):
        data = {
            'soul': None,
            'identity': None,
            'skills': [],
            'memories': [],
            'files': list_files(self.clawd_path),
        }
        
        soul_path = self.clawd_path / 'SOUL.md'
        if soul_path.exists():
            try:
                data['soul'] = soul_path.read_text(encoding='utf-8')
            except:
                pass
        
        identity_path = self.clawd_path / 'IDENTITY.md'
        if identity_path.exists():
            try:
                data['identity'] = identity_path.read_text(encoding='utf-8')
            except:
                pass
        
        skills_dir = self.clawd_path / 'skills'
        if skills_dir.exists():
            for item in skills_dir.iterdir():
                if item.is_dir():
                    data['skills'].append({
                        'name': item.name,
                        'location': 'local',
                        'status': 'active',
                        'enabled': True,
                    })
        
        memory_md = self.clawd_path / 'MEMORY.md'
        if memory_md.exists():
            try:
                content = memory_md.read_text(encoding='utf-8')
                data['memories'] = parse_memory_md(content)
            except:
                pass
        
        self.send_json(data)
    
    def handle_task_execute(self, data):
        """兼容旧的任务执行接口"""
        prompt = data.get('prompt', '').strip()
        if not prompt:
            self.send_error_json('Missing prompt', 400)
            return
        
        task_id = str(uuid.uuid4())[:8]
        
        thread = threading.Thread(
            target=run_task_in_background,
            args=(task_id, prompt, self.clawd_path),
            daemon=True,
        )
        thread.start()
        
        self.send_json({
            'taskId': task_id,
            'status': 'running',
        })
    
    # ============================================
    # 🤖 子代理 API 处理器 (Quest 模式支持)
    # ============================================
    
    def handle_subagent_spawn(self, data):
        """启动子代理"""
        if not self.subagent_manager:
            self.send_error_json('SubagentManager not initialized', 500)
            return
        
        agent_type = data.get('type', 'explore')
        task = data.get('task', '')
        tools = data.get('tools', [])
        context = data.get('context', '')
        
        if not task:
            self.send_error_json('Missing task', 400)
            return
        
        try:
            agent_id = self.subagent_manager.spawn(agent_type, task, tools, context)
            self.send_json({
                'status': 'success',
                'agentId': agent_id,
                'message': f'Spawned {agent_type} agent'
            })
        except Exception as e:
            self.send_error_json(f'Failed to spawn agent: {e}', 500)
    
    def handle_subagent_status(self, agent_id):
        """获取子代理状态"""
        if not self.subagent_manager:
            self.send_error_json('SubagentManager not initialized', 500)
            return
        
        status = self.subagent_manager.get_status(agent_id)
        if status:
            self.send_json({'status': 'success', 'agent': status})
        else:
            self.send_error_json(f'Agent not found: {agent_id}', 404)
    
    def handle_subagent_collect(self, data):
        """收集多个子代理的结果"""
        if not self.subagent_manager:
            self.send_error_json('SubagentManager not initialized', 500)
            return
        
        agent_ids = data.get('agentIds', [])
        timeout = data.get('timeout', 60.0)
        
        if not agent_ids:
            # 返回所有代理状态
            all_status = self.subagent_manager.get_all_status()
            self.send_json({'status': 'success', 'agents': all_status})
            return
        
        try:
            results = self.subagent_manager.collect_results(agent_ids, timeout)
            self.send_json({'status': 'success', 'results': results})
        except Exception as e:
            self.send_error_json(f'Failed to collect results: {e}', 500)
    
    def handle_llm_proxy(self, data: dict):
        """代理转发 LLM API 请求（解决 CORS 问题）
        
        前端请求: POST /api/llm/proxy
        Body: { "url": "https://api.moonshot.cn/v1/chat/completions", "apiKey": "sk-...", "body": {...}, "stream": true }
        
        对于 stream=true，使用分块传输将 SSE 事件流式转发给前端。
        """
        target_url = data.get('url')
        api_key = data.get('apiKey')
        request_body = data.get('body')
        is_stream = data.get('stream', False)
        custom_headers = data.get('headers')  # 前端可传入完整 headers（Anthropic 等非 Bearer 认证）
        
        if not target_url or not request_body:
            self.send_error_json('Missing url or body', 400)
            return
        
        # 若前端未提供 apiKey 也未提供 headers，报错
        if not api_key and not custom_headers:
            self.send_error_json('Missing apiKey or headers', 400)
            return
        
        try:
            import requests as req_lib
            import urllib3
            urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)
        except ImportError as e:
            self.send_error_json(f'LLM proxy requires "requests" package: pip install requests', 500)
            return
        
        print(f'[LLM Proxy] -> {target_url} (stream={is_stream})', file=sys.stderr)
        
        # 创建独立 Session，禁止读取系统代理 (Windows 注册表可能配有本地代理如 Clash/V2Ray)
        session = req_lib.Session()
        session.trust_env = False
        
        try:
            # 构建请求头: 优先使用前端传入的 custom_headers，否则回退到 Bearer token
            if custom_headers and isinstance(custom_headers, dict):
                req_headers = {'Content-Type': 'application/json'}
                req_headers.update(custom_headers)
            else:
                req_headers = {
                    'Content-Type': 'application/json',
                    'Authorization': f'Bearer {api_key}',
                }
            
            resp = session.post(
                target_url,
                json=request_body,
                headers=req_headers,
                stream=is_stream,
                timeout=(10, 300),  # connect 10s, read 300s (长任务可能很久)
                verify=False,
            )
            
            if is_stream:
                # 流式转发: 使用 chunked transfer encoding
                self.send_response(resp.status_code)
                self.send_header('Content-Type', 'text/event-stream; charset=utf-8')
                self.send_header('Cache-Control', 'no-cache')
                self.send_header('Access-Control-Allow-Origin', '*')
                self.send_header('Access-Control-Allow-Methods', 'GET, POST, PUT, DELETE, OPTIONS')
                self.send_header('Access-Control-Allow-Headers', 'Content-Type, Authorization')
                self.send_header('Transfer-Encoding', 'chunked')
                self.end_headers()
                
                try:
                    for chunk in resp.iter_content(chunk_size=None):
                        if chunk:
                            # HTTP chunked encoding: size\r\ndata\r\n
                            chunk_data = chunk
                            self.wfile.write(f'{len(chunk_data):x}\r\n'.encode())
                            self.wfile.write(chunk_data)
                            self.wfile.write(b'\r\n')
                            self.wfile.flush()
                    # 终止 chunk
                    self.wfile.write(b'0\r\n\r\n')
                    self.wfile.flush()
                except (BrokenPipeError, ConnectionResetError):
                    pass  # 客户端断开
                finally:
                    resp.close()
            else:
                # 非流式: 直接转发响应
                if resp.ok:
                    try:
                        self.send_json(resp.json())
                    except Exception:
                        self.send_response(resp.status_code)
                        self.send_header('Content-Type', 'application/json')
                        self.send_header('Access-Control-Allow-Origin', '*')
                        self.end_headers()
                        self.wfile.write(resp.content)
                else:
                    error_text = resp.text[:500]
                    print(f'[LLM Proxy] HTTP error: {resp.status_code} - {error_text}', file=sys.stderr)
                    self.send_error_json(f'LLM API error ({resp.status_code}): {error_text}', resp.status_code)
        
        except req_lib.exceptions.ConnectTimeout:
            self.send_error_json('LLM API connect timeout', 504)
        except req_lib.exceptions.ReadTimeout:
            self.send_error_json('LLM API read timeout', 504)
        except req_lib.exceptions.ConnectionError as e:
            self.send_error_json(f'Failed to connect to LLM API: {str(e)[:200]}', 502)
        except Exception as e:
            print(f'[LLM Proxy] Error: {type(e).__name__}: {e}', file=sys.stderr)
            self.send_error_json(f'LLM proxy error: {type(e).__name__}: {str(e)[:200]}', 500)
        finally:
            session.close()

    def handle_evomap_proxy(self, path: str, data: dict):
        """代理转发 EvoMap API 请求（解决 CORS 问题）
        
        前端请求: /api/evomap/a2a/hello
        转发到:   https://evomap.ai/a2a/hello
        """
        try:
            import requests
            import urllib3
            urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)
        except ImportError as e:
            print(f'[EvoMap Proxy] Missing dependency: {e}', file=sys.stderr)
            self.send_error_json(f'EvoMap proxy requires "requests" package: pip install requests', 500)
            return
        
        # 提取目标路径: /api/evomap/a2a/hello -> /a2a/hello
        target_path = path[11:]  # strip '/api/evomap' (keep leading /)
        target_url = f'https://evomap.ai{target_path}'
        
        print(f'[EvoMap Proxy] {path} -> {target_url}', file=sys.stderr)
        
        try:
            # evomap.ai 响应较慢 (实测需要 30+ 秒)，使用较长超时
            session = requests.Session()
            session.trust_env = False  # 绕过系统代理
            response = session.post(
                target_url,
                json=data if data else {},
                headers={
                    'Content-Type': 'application/json',
                    'User-Agent': 'DD-OS/1.0 EvoMap-Proxy',
                },
                timeout=(10, 90),  # connect 10s, read 90s (evomap.ai 响应慢)
                verify=False,  # 跳过 SSL 验证 (企业安全工具可能替换证书)
            )
            
            if response.ok:
                result = response.json()
                print(f'[EvoMap Proxy] Success: {response.status_code}', file=sys.stderr)
                self.send_json(result)
            else:
                error_text = response.text[:500]
                print(f'[EvoMap Proxy] HTTP error: {response.status_code} - {error_text}', file=sys.stderr)
                # 尝试返回原始 JSON 错误（可能包含有用信息）
                try:
                    self.send_json(response.json(), response.status_code)
                except Exception:
                    self.send_error_json(f'EvoMap API error: {response.status_code}', response.status_code)
                
        except requests.exceptions.ConnectTimeout as e:
            print(f'[EvoMap Proxy] Connect timeout: {e}', file=sys.stderr)
            self.send_error_json('Failed to connect to EvoMap (connect timeout)', 504)
        except requests.exceptions.ReadTimeout as e:
            print(f'[EvoMap Proxy] Read timeout (evomap.ai unresponsive): {e}', file=sys.stderr)
            self.send_error_json('EvoMap server unresponsive (read timeout)', 504)
        except requests.exceptions.ConnectionError as e:
            print(f'[EvoMap Proxy] Connection error: {e}', file=sys.stderr)
            self.send_error_json(f'Failed to connect to EvoMap: {str(e)[:200]}', 502)
        except Exception as e:
            print(f'[EvoMap Proxy] Error: {type(e).__name__}: {e}', file=sys.stderr)
            self.send_error_json(f'EvoMap proxy error: {type(e).__name__}: {str(e)[:200]}', 500)
        finally:
            session.close()
    
    def handle_task_status(self, task_id, offset=0):
        with self.tasks_lock:
            task = self.tasks.get(task_id)
        
        if not task:
            self.send_error_json(f'Task not found: {task_id}', 404)
            return
        
        log_path = task.get('logPath')
        content = ''
        new_offset = offset
        has_more = False
        file_size = task.get('fileSize', 0)
        
        if log_path:
            content, new_offset, has_more = read_log_chunk(log_path, offset)
            try:
                file_size = Path(log_path).stat().st_size
            except:
                pass
        
        self.send_json({
            'taskId': task_id,
            'status': task['status'],
            'content': content,
            'offset': new_offset,
            'hasMore': has_more,
            'fileSize': file_size,
        })


# ============================================
# 辅助函数
# ============================================

def list_files(clawd_path):
    files = []
    try:
        for item in clawd_path.iterdir():
            if item.is_file():
                files.append(item.name)
    except:
        pass
    return sorted(files)


def parse_memory_md(content):
    memories = []
    sections = content.split('## ')
    
    for i, section in enumerate(sections[1:], 1):
        lines = section.strip().split('\n')
        if not lines:
            continue
        
        title = lines[0].strip()
        body = '\n'.join(lines[1:]).strip()
        
        if title:
            memories.append({
                'id': f'memory-{i}',
                'title': title,
                'content': body[:500] if body else title,
                'type': 'long-term',
                'timestamp': None,
                'tags': [],
            })
    
    return memories


def read_log_chunk(log_path, offset=0, max_bytes=51200):
    path = Path(log_path)
    if not path.exists():
        return ('', offset, False)
    
    try:
        file_size = path.stat().st_size
    except:
        return ('', offset, False)
    
    if offset >= file_size:
        return ('', offset, False)
    
    try:
        with open(path, 'rb') as f:
            f.seek(offset)
            raw = f.read(max_bytes)
        
        content = raw.decode('utf-8', errors='replace')
        new_offset = offset + len(raw)
        has_more = new_offset < file_size
        return (content, new_offset, has_more)
    except Exception as e:
        return (f'[日志读取错误: {e}]', offset, False)


def run_task_in_background(task_id, prompt, clawd_path):
    logs_dir = clawd_path / 'logs'
    logs_dir.mkdir(exist_ok=True)
    log_file = logs_dir / f"{task_id}.log"
    
    with ClawdDataHandler.tasks_lock:
        ClawdDataHandler.tasks[task_id] = {
            'taskId': task_id,
            'status': 'running',
            'logPath': str(log_file),
            'fileSize': 0,
        }
    
    try:
        with open(log_file, 'w', encoding='utf-8') as f:
            f.write(f"Task: {prompt}\n")
            f.write(f"Started: {datetime.now().isoformat()}\n")
            f.write("-" * 50 + "\n\n")
        
        # 尝试运行 clawdbot，如果不存在则模拟
        try:
            with open(log_file, 'ab') as f:
                process = subprocess.Popen(
                    ['clawdbot', 'agent', '--agent', 'main', '--message', prompt],
                    cwd=str(clawd_path),
                    stdout=f,
                    stderr=subprocess.STDOUT,
                )
                
                start_time = time.time()
                timeout = 300
                
                while process.poll() is None:
                    time.sleep(0.5)
                    try:
                        with ClawdDataHandler.tasks_lock:
                            ClawdDataHandler.tasks[task_id]['fileSize'] = log_file.stat().st_size
                    except:
                        pass
                    
                    if time.time() - start_time > timeout:
                        process.kill()
                        process.wait()
                        with ClawdDataHandler.tasks_lock:
                            ClawdDataHandler.tasks[task_id]['status'] = 'error'
                            ClawdDataHandler.tasks[task_id]['fileSize'] = log_file.stat().st_size
                        with open(log_file, 'a', encoding='utf-8') as ef:
                            ef.write(f'\n\n[错误] 任务执行超时 ({timeout}s)\n')
                        return
                
                process.wait()
            
            with ClawdDataHandler.tasks_lock:
                ClawdDataHandler.tasks[task_id]['status'] = 'done' if process.returncode == 0 else 'error'
                ClawdDataHandler.tasks[task_id]['fileSize'] = log_file.stat().st_size
        
        except FileNotFoundError:
            # clawdbot 不存在，使用 Native 模式提示
            with open(log_file, 'a', encoding='utf-8') as f:
                f.write("\n[DD-OS Native] clawdbot 未安装。\n")
                f.write("在 Native 模式下，请使用 /api/tools/execute 接口直接执行工具。\n")
                f.write("\n任务已记录，等待 AI 引擎处理。\n")
            
            with ClawdDataHandler.tasks_lock:
                ClawdDataHandler.tasks[task_id]['status'] = 'done'
                ClawdDataHandler.tasks[task_id]['fileSize'] = log_file.stat().st_size
    
    except Exception as e:
        with open(log_file, 'a', encoding='utf-8') as ef:
            ef.write(f'\n\n[错误] {str(e)}\n')
        with ClawdDataHandler.tasks_lock:
            ClawdDataHandler.tasks[task_id]['status'] = 'error'
            ClawdDataHandler.tasks[task_id]['fileSize'] = log_file.stat().st_size


def cleanup_old_logs(clawd_path, max_age_hours=24):
    logs_dir = clawd_path / 'logs'
    if not logs_dir.exists():
        return
    
    now = time.time()
    count = 0
    for f in logs_dir.glob('*.log'):
        try:
            age = now - f.stat().st_mtime
            if age > max_age_hours * 3600:
                f.unlink()
                count += 1
        except:
            pass
    
    if count > 0:
        print(f"[Cleanup] Removed {count} old log files")


def cleanup_old_traces(clawd_path, max_months=6):
    """清理过期的执行追踪文件 (P2: 保留最近N个月)"""
    traces_dir = clawd_path / 'memory' / 'exec_traces'
    if not traces_dir.exists():
        return

    files = sorted(traces_dir.glob('*.jsonl'))
    if len(files) <= max_months:
        return

    old_files = files[:-max_months]
    for f in old_files:
        try:
            f.unlink()
            print(f"[Cleanup] Removed old trace: {f.name}")
        except:
            pass


def cleanup_temp_uploads(clawd_path, max_age_hours=1):
    """清理超过指定时间的临时上传文件"""
    upload_dir = clawd_path / 'temp' / 'uploads'
    if not upload_dir.exists():
        return
    
    now = time.time()
    count = 0
    for f in upload_dir.iterdir():
        try:
            if f.is_file() and (now - f.stat().st_mtime) > max_age_hours * 3600:
                f.unlink()
                count += 1
        except:
            pass
    
    if count > 0:
        print(f"[Cleanup] Removed {count} old temp upload files")


def main():
    parser = argparse.ArgumentParser(description='DD-OS Native Server')
    parser.add_argument('--port', type=int, default=3001, help='Server port (default: 3001)')
    # 支持环境变量覆盖默认路径
    default_path = os.getenv('DDOS_DATA_PATH', '~/.ddos')
    parser.add_argument('--path', type=str, default=default_path, help='Data directory path (default: ~/.ddos)')
    parser.add_argument('--host', type=str, default='0.0.0.0', help='Server host (default: 0.0.0.0)')
    args = parser.parse_args()
    
    clawd_path = Path(args.path).expanduser().resolve()
    
    if not clawd_path.exists():
        print(f"Creating data directory: {clawd_path}")
        clawd_path.mkdir(parents=True, exist_ok=True)
        
        # 创建默认 SOUL.md
        soul_file = clawd_path / 'SOUL.md'
        soul_file.write_text("""# DD-OS Native Soul

You are DD-OS, a local AI operating system running directly on the user's computer.

## Core Principles
- Be helpful and efficient
- Protect user data and privacy
- Execute tasks safely
- Learn from interactions

## Available Tools
- readFile: Read file contents
- writeFile: Write file contents
- listDir: List directory contents
- runCmd: Execute shell commands

## Safety Rules
- Never delete system files
- Ask before destructive operations
- Keep execution logs
""", encoding='utf-8')
        print(f"Created default SOUL.md")
    
    logs_dir = clawd_path / 'logs'
    logs_dir.mkdir(exist_ok=True)
    
    memory_dir = clawd_path / 'memory'
    memory_dir.mkdir(exist_ok=True)
    
    skills_dir = clawd_path / 'skills'
    skills_dir.mkdir(exist_ok=True)
    
    cleanup_old_logs(clawd_path)
    
    # 🔌 初始化工具注册表
    registry = ToolRegistry(clawd_path)
    # 注册内置工具
    builtin_names = [
        'readFile', 'writeFile', 'appendFile', 'listDir', 'runCmd',
        'weather', 'webSearch', 'webFetch', 'saveMemory', 'searchMemory',
        'nexusBindSkill', 'nexusUnbindSkill', 'openInExplorer', 'parseFile',
        'generateSkill',
    ]
    for name in builtin_names:
        registry.register_builtin(name, name)  # handler resolved at dispatch time
    # 扫描插件工具
    registry.scan_plugins()
    # 扫描 MCP 服务器
    registry.scan_mcp_servers()

    # 清理过期执行追踪 (P2: 保留最近6个月)
    cleanup_old_traces(clawd_path)
    cleanup_temp_uploads(clawd_path)

    # V2: 初始化 SQLite 数据库
    global _db_conn
    db_path = clawd_path / 'ddos_v2.db'
    _db_conn = init_sqlite_db(db_path)

    # 项目目录 (脚本所在目录 / exe 所在目录)
    project_path = APP_DIR
    
    ClawdDataHandler.clawd_path = clawd_path
    ClawdDataHandler.project_path = project_path
    ClawdDataHandler.registry = registry
    ClawdDataHandler.subagent_manager = SubagentManager(registry)
    
    server = ThreadingHTTPServer((args.host, args.port), ClawdDataHandler)
    
    tool_names = [t['name'] for t in registry.list_all()]
    plugin_count = len(registry.plugin_tools)
    mcp_count = len(registry.mcp_tools)
    print(f"""
+==================================================================+
|              DD-OS Native Server v{VERSION}                         |
+==================================================================+
|  Mode:    NATIVE (standalone, no OpenClaw needed)                |
|  Server:  http://{args.host}:{args.port}                                    |
|  Data:    {str(clawd_path)[:50]:<50} |
+------------------------------------------------------------------+
|  Tools:   {len(tool_names)} registered ({len(builtin_names)} builtin + {plugin_count} plugins + {mcp_count} mcp)    |
|  API:     /api/tools/execute (POST)  |  /tools (GET)            |
+==================================================================+
    """)
    
    print(f"Press Ctrl+C to stop\n")
    
    try:
        server.serve_forever()
    except KeyboardInterrupt:
        print("\nShutting down...")
        _browser_manager.shutdown()
        server.shutdown()


if __name__ == '__main__':
    main()
